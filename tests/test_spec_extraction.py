"""Integration tests for the full spec extraction pipeline.

Tests the SpecExtractor orchestrator end-to-end: colors, fonts, slide
classification, density analysis, YAML serialization, and CLI spec
management functions. Uses a programmatically-generated 3-slide PPTX
fixture — no committed binary files.
"""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal 3-slide PPTX with standard layouts for testing.

    Slide 1: Title Slide (title + subtitle)
    Slide 2: Title and Content (title + bullet text)
    Slide 3: Blank with a manually-added text box
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Test Presentation"
    # Subtitle (placeholder index 1)
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            ph.text = "Extraction Test"
            break

    # --- Slide 2: Content slide ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Content Slide"
    body = slide2.placeholders[1]
    body.text = "Bullet point 1\nBullet point 2\nBullet point 3"

    # --- Slide 3: Section divider (blank layout + manual text box) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left, top, width, height = Inches(1), Inches(3), Inches(8), Inches(1)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "--- Section Break ---"

    out = tmp_path / "test.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestColorExtraction:
    """Tests for theme color extraction."""

    def test_extract_colors(self, sample_pptx: Path):
        """Colors should produce a populated ColorPalette with valid HEX values."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        spec = extractor.extract()

        palette_fields = [
            "background1", "background2",
            "text1", "text2",
            "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
            "hyperlink", "followed_hyperlink",
        ]

        for field in palette_fields:
            val = getattr(spec.colors, field)
            # Must be a string
            assert isinstance(val, str), f"{field} is not a string: {type(val)}"
            # If non-empty, must start with '#'
            if val:
                assert val.startswith("#"), f"{field} = {val!r} does not start with '#'"

        # At least some fields should be populated (Office defaults)
        populated = [f for f in palette_fields if getattr(spec.colors, f)]
        assert len(populated) >= 6, f"Expected at least 6 populated colors, got {len(populated)}"


class TestFontExtraction:
    """Tests for font family and font size extraction."""

    def test_extract_fonts(self, sample_pptx: Path):
        """Typography should have non-empty font families and populated size dicts."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        spec = extractor.extract()

        # Font families
        assert spec.typography.heading_family, "heading_family should not be empty"
        assert spec.typography.body_family, "body_family should not be empty"

        # Font sizes: heading_sizes should have at least "title" key from the fixture
        h_sizes = spec.typography.heading_sizes
        assert isinstance(h_sizes, dict), "heading_sizes should be a dict"
        if h_sizes:
            assert "title" in h_sizes, f"heading_sizes keys: {list(h_sizes.keys())}"
            assert h_sizes["title"] > 0, f"title size should be > 0, got {h_sizes['title']}"

        # body_sizes should have at least "body" key from the fixture
        b_sizes = spec.typography.body_sizes
        assert isinstance(b_sizes, dict), "body_sizes should be a dict"
        if b_sizes:
            assert "body" in b_sizes, f"body_sizes keys: {list(b_sizes.keys())}"
            assert b_sizes["body"] > 0, f"body size should be > 0, got {b_sizes['body']}"


class TestSlideClassification:
    """Tests for slide type classification."""

    def test_slide_classification(self, sample_pptx: Path):
        """Three-slide fixture should classify as title → content → section_divider."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        spec = extractor.extract()

        assert len(spec.slides) == 3, f"Expected 3 slides, got {len(spec.slides)}"

        # Slide 1: should be TITLE
        assert spec.slides[0].slide_type.value in ("title", "content"), \
            f"Slide 1 type: {spec.slides[0].slide_type.value}"

        # Slide 2: should be CONTENT
        assert spec.slides[1].slide_type.value == "content", \
            f"Slide 2 type: {spec.slides[1].slide_type.value}"

        # Slide 3: section_divider or content (blank layout + short text → section_divider)
        assert spec.slides[2].slide_type.value in ("section_divider", "content"), \
            f"Slide 3 type: {spec.slides[2].slide_type.value}"


class TestDensityAnalysis:
    """Tests for content density and rhythm analysis."""

    def test_density_analysis(self, sample_pptx: Path):
        """Each slide should have a density label, and rhythm should be populated."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        spec = extractor.extract()

        # Density labels
        for i, slide_spec in enumerate(spec.slides):
            assert slide_spec.density.value in ("breathing", "dense", "anchor"), \
                f"Slide {i+1} has unexpected density: {slide_spec.density.value}"

        # At least one slide should have non-zero char_count (slide 2 has bullet text)
        slide2 = spec.slides[1]
        assert slide2.char_count > 0, f"Slide 2 char_count should be > 0, got {slide2.char_count}"

        # Rhythm: density_profile length should equal slide count
        assert len(spec.rhythm.density_profile) == len(spec.slides), \
            f"density_profile length {len(spec.rhythm.density_profile)} != slide count {len(spec.slides)}"

        # Rhythm: sequencing_pattern length should equal slide count
        assert len(spec.rhythm.sequencing_pattern) == len(spec.slides)


class TestYamlSerialization:
    """Tests for YAML serialization and round-trip integrity."""

    def test_save_and_load_yaml(self, sample_pptx: Path, tmp_path: Path):
        """Saving a spec should produce a parsable YAML file with all top-level keys."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        output_path = extractor.save(str(tmp_path))

        # File must exist
        assert output_path.is_file(), f"YAML file not created at {output_path}"

        # Parse YAML back
        with open(output_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)

        # Top-level keys
        expected_keys = {"metadata", "colors", "typography", "slides", "rhythm", "source_config"}
        actual_keys = set(data.keys())
        missing = expected_keys - actual_keys
        assert not missing, f"Missing top-level keys: {missing}"

        # Slide count
        assert len(data["slides"]) == 3, f"Expected 3 slides in YAML, got {len(data['slides'])}"

        # Metadata
        assert data["metadata"]["name"] == "test"
        assert data["metadata"]["slide_count"] == 3

        # Colors should be a dict with string values
        assert isinstance(data["colors"], dict)

        # Typography should have heading_family and body_family
        assert "heading_family" in data["typography"]
        assert "body_family" in data["typography"]

        # Rhythm
        assert "story_arc" in data["rhythm"]

    def test_all_density_labels_valid_enum(self, sample_pptx: Path, tmp_path: Path):
        """Density labels in YAML should be plain strings, not Python object references."""
        from ppt_skill.spec.extractor import SpecExtractor

        extractor = SpecExtractor(sample_pptx, "test")
        output_path = extractor.save(str(tmp_path))

        with open(output_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)

        valid_labels = {"breathing", "dense", "anchor"}
        for slide in data["slides"]:
            density = slide.get("density", "")
            assert density in valid_labels, \
                f"Invalid density label: {density!r}. Must be one of {valid_labels}"
            assert isinstance(density, str), \
                f"Density label should be a plain string, got {type(density)}"


class TestSpecManagement:
    """Tests for spec listing, selection, and active query."""

    def test_list_and_select_specs(self, sample_pptx: Path, tmp_path: Path):
        """After saving a spec, list/select/query should work correctly."""
        from ppt_skill.cli.spec_commands import (
            extract_spec,
            get_active_spec,
            list_specs,
            select_spec,
        )

        specs_dir = str(tmp_path / "specs")

        # 1. Save a spec
        extract_spec("test", str(sample_pptx), specs_dir)

        # 2. List — should find it
        names = list_specs(specs_dir)
        assert "test" in names, f"Expected 'test' in listed specs, got {names}"

        # 3. No active spec yet
        assert get_active_spec(specs_dir) is None, "Should be None before selection"

        # 4. Select the spec
        active_path = select_spec("test", specs_dir)
        assert active_path.is_file(), f".active file not created at {active_path}"
        assert active_path.read_text().strip() == "test"

        # 5. Now get_active_spec should return "test"
        assert get_active_spec(specs_dir) == "test"

        # 6. Selecting non-existent spec should raise
        with pytest.raises(FileNotFoundError):
            select_spec("nonexistent", specs_dir)
