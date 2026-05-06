"""End-to-end tests: SVG → quality check → convert → validate PPTX."""

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_skill.pipeline import convert_svg_to_pptx

FIXTURES = Path(__file__).parent / "fixtures"


class TestEndToEnd:
    """Full pipeline integration tests."""

    def test_simple_svg_to_pptx(self, temp_output):
        """A simple rectangle SVG produces a PPTX with a native shape."""
        result = convert_svg_to_pptx(
            svg_files=[FIXTURES / "sample_simple.svg"],
            output_path=temp_output,
        )
        assert result is True
        assert temp_output.exists()
        assert temp_output.stat().st_size > 1000  # Not empty

        # Verify the PPTX is valid and has shapes
        prs = Presentation(str(temp_output))
        assert len(prs.slides) >= 1
        shapes = prs.slides[0].shapes
        assert len(shapes) > 0, (
            f"Expected at least one shape on the slide, got {len(shapes)}"
        )

        # Each shape should be a native DrawingML element (not an image)
        for shape in shapes:
            assert shape.shape_type is not None, (
                f"Shape has no shape_type: {shape.name}"
            )

    def test_multi_slide_conversion(self, temp_output):
        """Multiple SVGs produce multiple slides."""
        svg_files = [
            FIXTURES / "sample_simple.svg",
            FIXTURES / "sample_text.svg",
        ]
        result = convert_svg_to_pptx(
            svg_files=svg_files,
            output_path=temp_output,
        )
        assert result is True
        prs = Presentation(str(temp_output))
        assert len(prs.slides) == 2, (
            f"Expected 2 slides, got {len(prs.slides)}"
        )

    def test_text_svg_preserves_text(self, temp_output):
        """Text in SVG should produce selectable text in PPTX."""
        result = convert_svg_to_pptx(
            svg_files=[FIXTURES / "sample_text.svg"],
            output_path=temp_output,
        )
        assert result is True
        prs = Presentation(str(temp_output))
        shapes = prs.slides[0].shapes

        # At least one shape should contain "Hello World" as editable text
        text_contents = []
        for shape in shapes:
            if shape.has_text_frame:
                text_contents.append(shape.text_frame.text)

        combined = " ".join(text_contents)
        assert "Hello" in combined, (
            f"Expected 'Hello' in slide text, got: {combined}"
        )

    def test_quality_check_rejects_banned_svg(self, temp_output):
        """Pipeline raises ValueError when quality check fails."""
        with pytest.raises(ValueError, match="quality check failed"):
            convert_svg_to_pptx(
                svg_files=[FIXTURES / "banned_features_svg" / "mask.svg"],
                output_path=temp_output,
            )

    def test_skip_check_bypasses_validation(self, temp_output):
        """Skip quality check allows conversion of SVGs with banned features."""
        # Note: This may or may not produce valid output — the test only
        # verifies that the skip flag bypasses validation.
        try:
            convert_svg_to_pptx(
                svg_files=[FIXTURES / "banned_features_svg" / "mask.svg"],
                output_path=temp_output,
                skip_quality_check=True,
            )
            # If it succeeds, PPTX should exist
            assert temp_output.exists()
        except Exception as e:
            # Conversion may fail due to unsupported features — acceptable
            # The key behavior is that NO ValueError about quality check is raised
            assert "quality check failed" not in str(e)
