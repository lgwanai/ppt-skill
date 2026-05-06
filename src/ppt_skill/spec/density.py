"""Content density and presentation rhythm analysis.

Computes per-slide content metrics (character count, image count, shape count)
and classifies slides into breathing/dense/anchor categories using
**percentile-based thresholds** (not absolute character counts). This adapts
to any deck length and writing style — a 10-slide pitch deck and a 60-slide
corporate deck both get meaningful density classifications.

Also builds the ``PresentationRhythm`` dataclass — capturing sequencing
patterns, density flow, and a heuristic story arc estimation.

Exports:
    analyze_slide_density(slide, slide_index) -> dict   — per-slide metrics
    classify_density(densities) -> list[dict]            — percentile-based labels
    programmatic_percentile(values, pct) -> float         — linear interpolation percentile
    build_presentation_rhythm(densities, types) -> PresentationRhythm
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_skill.spec.spec_model import DensityLabel, PresentationRhythm

if TYPE_CHECKING:
    from pptx.slide import Slide


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def programmatic_percentile(values: list[int], pct: float) -> float:
    """Compute a percentile using linear interpolation (numpy default method).

    Args:
        values: List of integer values (e.g., char counts).
        pct: Percentile to compute, between 0 and 100 (e.g., 20 for 20th).

    Returns:
        The interpolated percentile value as a float.

    Raises:
        ValueError: If values is empty.
    """
    if not values:
        raise ValueError("Cannot compute percentile of empty list")

    sorted_values = sorted(values)
    n = len(sorted_values)

    # Linear interpolation index
    # i = (pct / 100) * (n - 1)
    idx = (pct / 100.0) * (n - 1)

    lo = int(idx)
    hi = lo + 1

    if hi >= n:
        return float(sorted_values[-1])

    frac = idx - lo
    return sorted_values[lo] + frac * (sorted_values[hi] - sorted_values[lo])


def analyze_slide_density(slide: Slide, slide_index: int) -> dict:
    """Compute per-slide content metrics: chars, images, shapes.

    Args:
        slide: A python-pptx Slide object.
        slide_index: 1-based slide index for identification.

    Returns:
        Dict with keys: "slide_index" (int), "char_count" (int),
        "image_count" (int), "shape_count" (int).
    """
    char_count: int = 0
    image_count: int = 0
    shape_count: int = 0

    for shape in slide.shapes:
        shape_count += 1

        # --- Character counting ---
        try:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if shape.text_frame is not None:
                    text = shape.text_frame.text
                    if text:
                        char_count += len(text)
        except (AttributeError, TypeError):
            pass

        # --- Image detection ---
        try:
            # Check via shape_type
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
                image_count += 1
                continue
        except (AttributeError, TypeError):
            pass

        try:
            # Check via image attribute
            if hasattr(shape, "image") and shape.image is not None:
                image_count += 1
        except (AttributeError, TypeError):
            pass

    return {
        "slide_index": slide_index,
        "char_count": char_count,
        "image_count": image_count,
        "shape_count": shape_count,
    }


def classify_density(densities: list[dict]) -> list[dict]:
    """Assign breathing/dense/anchor labels using percentile-based thresholds.

    Classification rules:
        - char_count <= 20th percentile → "breathing" (light content)
        - char_count >= 80th percentile → "anchor" (information-dense)
        - Between 20th and 80th → "dense" (normal)
        - Edge case: all identical → all "breathing"

    Args:
        densities: List of per-slide density dicts from analyze_slide_density().

    Returns:
        The same list with a "density" (str) field added to each entry.
    """
    if not densities:
        return []

    char_counts = [d["char_count"] for d in densities]

    # Edge case: all slides have identical char_count
    if len(set(char_counts)) == 1:
        for d in densities:
            d["density"] = DensityLabel.BREATHING.value
        return densities

    p20 = programmatic_percentile(char_counts, 20.0)
    p80 = programmatic_percentile(char_counts, 80.0)

    for d in densities:
        cc = d["char_count"]
        if cc <= p20:
            d["density"] = DensityLabel.BREATHING.value
        elif cc >= p80:
            d["density"] = DensityLabel.ANCHOR.value
        else:
            d["density"] = DensityLabel.DENSE.value

    return densities


def build_presentation_rhythm(
    slide_densities: list[dict],
    slide_types: list[str],
) -> PresentationRhythm:
    """Build a PresentationRhythm from ordered slide type and density sequences.

    Constructs a structured rhythm analysis capturing the sequencing pattern,
    density flow, and a heuristic story arc estimation.

    Args:
        slide_densities: Per-slide density dicts (with "density" field).
        slide_types: Ordered list of slide type strings matching density order.

    Returns:
        A PresentationRhythm dataclass instance.

    Story arc heuristic:
        - "opening": first 1–2 slides if they are title/section_divider types
        - "development": content/image_text slides in the middle
        - "climax": data slides or anchor-density slides
        - "closing": last 1–2 slides if section_divider/data
    """
    sequencing_pattern: list[str] = list(slide_types)
    density_profile: list[str] = [d["density"] for d in slide_densities]

    n = len(slide_types)
    if n == 0:
        return PresentationRhythm(
            sequencing_pattern=sequencing_pattern,
            density_profile=density_profile,
            story_arc={"opening": 0, "development": 0, "climax": 0, "closing": 0},
        )

    # --- Story arc heuristics ---
    opening_count: int = 0
    development_count: int = 0
    climax_count: int = 0
    closing_count: int = 0

    for i, (stype, dlabel) in enumerate(zip(slide_types, density_profile)):
        is_first_slide = i == 0
        is_last_slide = i == n - 1
        is_early = i <= max(1, n // 5)   # first ~20% of deck
        is_late = i >= n - max(1, n // 5)  # last ~20% of deck

        if stype == "title":
            opening_count += 1
            continue

        if stype == "section_divider":
            if is_early:
                opening_count += 1
            elif is_late:
                closing_count += 1
            else:
                development_count += 1
            continue

        if stype == "data":
            climax_count += 1
            continue

        # Check density for remaining content/image_text slides
        if dlabel == "anchor":
            # Anchor density anywhere suggests climax (information peak)
            if is_early:
                opening_count += 1
            elif is_late:
                closing_count += 1
            else:
                climax_count += 1  # mid-deck anchor = climax
        elif dlabel == "breathing":
            if is_early:
                opening_count += 1
            elif is_late:
                closing_count += 1
            else:
                development_count += 1
        else:
            # Normal dense content → development
            if is_early and stype == "title":
                opening_count += 1
            elif is_late and (stype == "section_divider" or dlabel == "breathing"):
                closing_count += 1
            else:
                development_count += 1

    story_arc: dict[str, int] = {
        "opening": opening_count,
        "development": development_count,
        "climax": climax_count,
        "closing": closing_count,
    }

    return PresentationRhythm(
        sequencing_pattern=sequencing_pattern,
        density_profile=density_profile,
        story_arc=story_arc,
    )


__all__ = [
    "analyze_slide_density",
    "build_presentation_rhythm",
    "classify_density",
    "programmatic_percentile",
]
