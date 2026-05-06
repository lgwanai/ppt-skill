"""Spatial layout measurement from python-pptx shape positions.

Measures per-slide margins, title positioning, and content regions by analyzing
shape bounding boxes. All measurements are converted from EMU (English Metric
Units, 914400 EMU = 1 inch) to inches and rounded to 2 decimal places.

Exports:
    analyze_slide_layout(slide) -> dict     — analyze a single slide's spatial layout
    analyze_all_slides(prs) -> list[dict]   — analyze all slides with index
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.util import Emu

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


def analyze_slide_layout(slide: Slide) -> dict:
    """Analyze a single slide's spatial layout from shape positions.

    Returns a dict with:
        margins: {top, bottom, left, right} in inches
        title_position: {x, y, width, height} in inches or None
        content_regions: list of {x, y, width, height} dicts in inches

    All measurements are rounded to 2 decimal places.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        Dict with "margins", "title_position", and "content_regions" keys.
    """
    # Slide dimensions in EMU
    slide_width: int = slide.slide_layout.slide_width   # EMU
    slide_height: int = slide.slide_layout.slide_height  # EMU

    shapes = list(slide.shapes)
    if not shapes:
        # Empty slide — return slide dimensions as margins
        w_inches = round(Emu(slide_width).inches, 2)
        h_inches = round(Emu(slide_height).inches, 2)
        return {
            "margins": {"top": h_inches, "bottom": h_inches, "left": w_inches, "right": w_inches},
            "title_position": None,
            "content_regions": [],
        }

    # Collect bounding boxes for all non-background shapes
    boxes: list[dict] = []
    title_box: dict | None = None

    for shape in shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.type is None:
                continue  # background/unsupported placeholder
        except (AttributeError, TypeError):
            pass

        try:
            left = shape.left   # EMU
            top = shape.top
            width = shape.width
            height = shape.height
        except (AttributeError, TypeError):
            continue  # shape without position info

        box = {
            "left": left,
            "top": top,
            "right": left + width,
            "bottom": top + height,
            "width": width,
            "height": height,
        }

        # Check if this is the title shape
        try:
            slide_title = slide.shapes.title
        except (AttributeError, TypeError):
            slide_title = None

        if slide_title is not None and shape is slide_title:
            title_box = box

        boxes.append(box)

    if not boxes:
        # All shapes were background/unpositioned
        w_inches = round(Emu(slide_width).inches, 2)
        h_inches = round(Emu(slide_height).inches, 2)
        return {
            "margins": {"top": h_inches, "bottom": h_inches, "left": w_inches, "right": w_inches},
            "title_position": None,
            "content_regions": [],
        }

    # --- Margin calculation ---
    top_margin_emu: int = min(b["top"] for b in boxes)
    left_margin_emu: int = min(b["left"] for b in boxes)
    bottom_margin_emu: int = slide_height - max(b["bottom"] for b in boxes)
    right_margin_emu: int = slide_width - max(b["right"] for b in boxes)

    margins = {
        "top": round(Emu(top_margin_emu).inches, 2),
        "bottom": round(Emu(bottom_margin_emu).inches, 2),
        "left": round(Emu(left_margin_emu).inches, 2),
        "right": round(Emu(right_margin_emu).inches, 2),
    }

    # --- Title detection ---
    title_position: dict | None = None
    if title_box is not None:
        title_position = {
            "x": round(Emu(title_box["left"]).inches, 2),
            "y": round(Emu(title_box["top"]).inches, 2),
            "width": round(Emu(title_box["width"]).inches, 2),
            "height": round(Emu(title_box["height"]).inches, 2),
        }
    else:
        # Fallback: find the highest text shape as title candidate
        text_boxes = [
            b for b in boxes if b["left"] == min(b["left"] for b in boxes)
        ] or boxes
        if text_boxes:
            highest = min(text_boxes, key=lambda b: b["top"])
            # Only treat as title if it's in the top ~20% of the slide
            if highest["top"] < slide_height * 0.2:
                title_position = {
                    "x": round(Emu(highest["left"]).inches, 2),
                    "y": round(Emu(highest["top"]).inches, 2),
                    "width": round(Emu(highest["width"]).inches, 2),
                    "height": round(Emu(highest["height"]).inches, 2),
                }

    # --- Content regions ---
    # Non-title shapes grouped by vertical overlap
    title_box_ref = title_box  # for closure
    content_boxes = [b for b in boxes if title_box_ref is None or b is not title_box_ref]

    content_regions: list[dict] = []
    for box in content_boxes:
        content_regions.append({
            "x": round(Emu(box["left"]).inches, 2),
            "y": round(Emu(box["top"]).inches, 2),
            "width": round(Emu(box["width"]).inches, 2),
            "height": round(Emu(box["height"]).inches, 2),
        })

    return {
        "margins": margins,
        "title_position": title_position,
        "content_regions": content_regions,
    }


def analyze_all_slides(prs: Presentation) -> list[dict]:
    """Analyze spatial layout for every slide in a presentation.

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        List of dicts from analyze_slide_layout(), each with "slide_index" added.
    """
    results: list[dict] = []
    for i, slide in enumerate(prs.slides):
        layout_data = analyze_slide_layout(slide)
        layout_data["slide_index"] = i + 1
        results.append(layout_data)
    return results


__all__ = [
    "analyze_all_slides",
    "analyze_slide_layout",
]
