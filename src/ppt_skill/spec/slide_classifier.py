"""Slide type classification — dual-strategy: layout name mapping + content-based fallback.

Strategy 1 (layout-name-based): Maps standard PowerPoint layout names to slide types
via LAYOUT_NAME_MAP. This covers ~90% of slides in standard templates. Comparison is
case-insensitive and whitespace-normalized to handle custom template variations.

Strategy 2 (content-based fallback): For unrecognized or custom layouts, inspects shape
types (placeholders, charts, tables, images) and title presence to classify the slide.

Exports:
    classify_slide(slide) -> str        — classify a single slide
    classify_all_slides(prs) -> list[dict]  — classify all slides with metadata
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from ppt_skill.spec.spec_model import SlideType

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


# ---------------------------------------------------------------------------
# Layout name → slide type mapping (case-insensitive, whitespace-normalized)
# ---------------------------------------------------------------------------

LAYOUT_NAME_MAP: dict[str, str] = {
    "title slide": SlideType.TITLE.value,
    "1_title slide": SlideType.TITLE.value,  # common enterprise variant
    "title and content": SlideType.CONTENT.value,
    "section header": SlideType.SECTION_DIVIDER.value,
    "two content": SlideType.IMAGE_TEXT.value,
    "picture with caption": SlideType.IMAGE_TEXT.value,
    "content with caption": SlideType.IMAGE_TEXT.value,
    "comparison": SlideType.CONTENT.value,
    "title only": SlideType.CONTENT.value,
    "blank": SlideType.CONTENT.value,
    "agenda": SlideType.CONTENT.value,
    "section": SlideType.SECTION_DIVIDER.value,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def classify_slide(slide: Slide) -> str:
    """Classify a slide into one of 5 types using dual-strategy.

    Strategy 1: Look up the slide's layout name (normalized) in LAYOUT_NAME_MAP.
    Strategy 2 (fallback): Content-based inspection — charts, tables, images, titles.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        A SlideType string value: "title", "content", "section_divider",
        "image_text", or "data".
    """
    # --- Strategy 1: layout name mapping ---
    layout_name: str = slide.slide_layout.name
    normalized: str = layout_name.strip().lower()
    if normalized in LAYOUT_NAME_MAP:
        return LAYOUT_NAME_MAP[normalized]

    # --- Strategy 2: content-based fallback ---
    has_chart: bool = False
    has_table: bool = False
    has_image: bool = False
    placeholder_count: int = 0

    for shape in slide.shapes:
        # Check for charts
        if not has_chart:
            try:
                if shape.has_table:
                    has_table = True
            except (AttributeError, TypeError):
                pass
            try:
                if shape.has_chart:
                    has_chart = True
            except (AttributeError, TypeError):
                pass

        # Check for images via shape_type or image attribute
        if not has_image:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
                    has_image = True
            except (AttributeError, TypeError):
                pass
            if not has_image and hasattr(shape, "image"):
                try:
                    if shape.image is not None:
                        has_image = True
                except (AttributeError, TypeError):
                    pass

        # Count placeholders
        try:
            if shape.is_placeholder:
                placeholder_count += 1
        except (AttributeError, TypeError):
            pass

    has_title: bool = slide.shapes.title is not None

    # Data slides: contain charts or tables
    if has_chart or has_table:
        return SlideType.DATA.value

    # Image-text slides: images present, no title
    if has_image and not has_title:
        return SlideType.IMAGE_TEXT.value

    # Section divider candidates: has title, very few shapes, short title text
    if has_title and placeholder_count <= 3:
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title_text = title_shape.text_frame.text.strip()
                if len(title_text) < 30:
                    return SlideType.SECTION_DIVIDER.value
        except (AttributeError, TypeError):
            pass

    # Content slides: has title with more shapes, or default
    if has_title:
        return SlideType.CONTENT.value

    # Default fallback
    return SlideType.CONTENT.value


def classify_all_slides(prs: Presentation) -> list[dict]:
    """Classify every slide in a presentation.

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        List of dicts with keys: "slide_index" (1-based), "type" (str),
        "layout_name" (str — raw layout name from PPTX).
    """
    results: list[dict] = []
    for i, slide in enumerate(prs.slides):
        slide_type: str = classify_slide(slide)
        layout_name: str = slide.slide_layout.name
        results.append({
            "slide_index": i + 1,
            "type": slide_type,
            "layout_name": layout_name,
        })
    return results


__all__ = [
    "LAYOUT_NAME_MAP",
    "classify_all_slides",
    "classify_slide",
]
