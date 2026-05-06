"""Per-slide font size and weight extraction with inheritance chain resolution.

Walks the run→paragraph inheritance chain for every text element on a slide,
resolving actual font sizes (in points) and bold weights. python-pptx's
``run.font.size`` returns None for inherited values — this module resolves
through the paragraph-level default, with a universal fallback of 18pt.

Separates heading text (title placeholders, top-of-slide large text) from
body text (everything else) for side-by-side size/weight statistics.

Closes the SPC-01 gap: theme.py only extracts font *families* from theme1.xml;
this module extracts actual *sizes and weights* applied to text.

Exports:
    extract_slide_fonts(slide) -> dict           — analyze a single slide
    extract_all_slide_fonts(prs) -> list[dict]   — analyze all slides
    compute_spec_typography_sizes(all_fonts) -> dict  — aggregate Typography-ready dicts
    median(values) -> float                       — pure-Python median helper
"""

from __future__ import annotations

import statistics
from typing import TYPE_CHECKING

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def median(values: list[float]) -> float:
    """Compute the median of a list of floats.

    Pure Python implementation — no numpy dependency. Uses statistics.median
    for Python 3.4+.

    Args:
        values: List of float values.

    Returns:
        The median value. Returns 0.0 for an empty list.
    """
    if not values:
        return 0.0
    return float(statistics.median(values))


def extract_slide_fonts(slide: Slide) -> dict:
    """Extract font sizes and weights from a single slide.

    Walks all shapes with text frames, iterates each paragraph's runs, and
    resolves font size through the run→paragraph inheritance chain. Font
    weight (bold) is resolved similarly.

    Text is classified as heading or body based on placeholder type and
    position on the slide.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        Dict with keys:
            heading_sizes: {min, max, median} in points (float)
            body_sizes: {min, max, median} in points (float)
            heading_bold: bool — True if any heading runs are bold
            body_bold: bool — True if any body runs are bold
            all_heading_sizes: list[float] — raw sizes for percentile computation
            all_body_sizes: list[float] — raw sizes for percentile computation
    """
    all_heading_sizes: list[float] = []
    all_body_sizes: list[float] = []
    heading_bold: bool = False
    body_bold: bool = False

    # Get reference to the title shape for classification
    try:
        slide_title_shape = slide.shapes.title
    except (AttributeError, TypeError):
        slide_title_shape = None

    for shape in slide.shapes:
        # Skip shapes without text frames
        if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
            continue

        try:
            text_frame = shape.text_frame
        except (AttributeError, TypeError):
            continue

        if text_frame is None:
            continue

        # Determine if this shape is heading or body
        is_heading: bool = False

        # Check if it's a title placeholder
        try:
            if shape.is_placeholder and shape.placeholder_format.type is not None:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:  # 1
                    is_heading = True
        except (AttributeError, TypeError):
            pass

        # Check if it's the slide's title shape
        if not is_heading and slide_title_shape is not None:
            if shape is slide_title_shape:
                is_heading = True

        # Check if it's a top-of-slide shape with large text
        if not is_heading:
            try:
                shape_top = shape.top  # EMU
                # 1.5 inches in EMU = 1.5 * 914400 = 1371600
                if shape_top < 1371600:
                    # Check if any text component is large (>28pt)
                    for paragraph in text_frame.paragraphs:
                        resolved_size = _resolve_font_size(paragraph)
                        if resolved_size is not None and resolved_size > 28.0:
                            is_heading = True
                            break
            except (AttributeError, TypeError):
                pass

        # Extract font sizes and weights from all paragraphs
        for paragraph in text_frame.paragraphs:
            # Resolve paragraph-level defaults
            para_font_size = _resolve_paragraph_font_size(paragraph)
            para_font_bold = _resolve_paragraph_font_bold(paragraph)

            for run in paragraph.runs:
                # --- Font size ---
                size_pt = _resolve_run_font_size(run, para_font_size)
                if size_pt is not None and size_pt > 0:
                    if is_heading:
                        all_heading_sizes.append(size_pt)
                    else:
                        all_body_sizes.append(size_pt)

                # --- Font weight ---
                run_bold = _resolve_run_font_bold(run, para_font_bold)
                if run_bold:
                    if is_heading:
                        heading_bold = True
                    else:
                        body_bold = True

    # Compute statistics — defensive: return 0.0 for empty lists
    heading_stats = {
        "min": min(all_heading_sizes) if all_heading_sizes else 0.0,
        "max": max(all_heading_sizes) if all_heading_sizes else 0.0,
        "median": median(all_heading_sizes),
    }

    body_stats = {
        "min": min(all_body_sizes) if all_body_sizes else 0.0,
        "max": max(all_body_sizes) if all_body_sizes else 0.0,
        "median": median(all_body_sizes),
    }

    return {
        "heading_sizes": heading_stats,
        "body_sizes": body_stats,
        "heading_bold": heading_bold,
        "body_bold": body_bold,
        "all_heading_sizes": all_heading_sizes,
        "all_body_sizes": all_body_sizes,
    }


def extract_all_slide_fonts(prs: Presentation) -> list[dict]:
    """Extract font sizes and weights for every slide in a presentation.

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        List of dicts from extract_slide_fonts(), each with "slide_index" added.
    """
    results: list[dict] = []
    for i, slide in enumerate(prs.slides):
        font_data = extract_slide_fonts(slide)
        font_data["slide_index"] = i + 1
        results.append(font_data)
    return results


def compute_spec_typography_sizes(all_fonts: list[dict]) -> dict:
    """Compute aggregate Typography heading_sizes/body_sizes from per-slide data.

    Aggregates raw font sizes across all slides and computes the typical
    size hierarchy. Designed to populate ``Typography.heading_sizes`` and
    ``Typography.body_sizes`` fields.

    Args:
        all_fonts: List of per-slide font dicts from extract_all_slide_fonts().

    Returns:
        Dict with:
            heading_sizes: {"title", "subtitle", "h1", "h2"} → float pt values
            body_sizes: {"body", "small", "caption"} → float pt values
    """
    all_heading: list[float] = []
    all_body: list[float] = []

    for font_data in all_fonts:
        all_heading.extend(font_data.get("all_heading_sizes", []))
        all_body.extend(font_data.get("all_body_sizes", []))

    max_heading: float = max(all_heading) if all_heading else 0.0
    median_heading: float = median(all_heading)
    min_body: float = min(all_body) if all_body else 0.0
    median_body: float = median(all_body)

    heading_sizes: dict[str, float] = {
        "title": max_heading,
        "subtitle": median_heading,
        "h1": max_heading,
        "h2": median_heading,
    }

    body_sizes: dict[str, float] = {
        "body": median_body,
        "small": min_body,
        "caption": min_body,
    }

    return {
        "heading_sizes": heading_sizes,
        "body_sizes": body_sizes,
    }


# ---------------------------------------------------------------------------
# Internal font resolution helpers
# ---------------------------------------------------------------------------


def _resolve_run_font_size(run, para_default: float | None) -> float | None:
    """Resolve font size for a run through the inheritance chain.

    Hierarchy: run.font.size → paragraph.font.size → Pt(18) fallback.

    Args:
        run: A python-pptx Run object.
        para_default: Font size resolved at the paragraph level (or None).

    Returns:
        Font size in points (float), or None if unresolved and no fallback.
    """
    try:
        if run.font.size is not None:
            return run.font.size.pt  # Pt object → float
    except (AttributeError, TypeError):
        pass

    if para_default is not None:
        return para_default

    # Universal default for body text
    return Pt(18).pt  # 18.0


def _resolve_run_font_bold(run, para_default: bool) -> bool:
    """Resolve font weight (bold) for a run through the inheritance chain.

    Hierarchy: run.font.bold → paragraph.font.bold → False.

    Args:
        run: A python-pptx Run object.
        para_default: Bold resolved at the paragraph level.

    Returns:
        True if the run is bold, False otherwise.
    """
    try:
        if run.font.bold is not None:
            return bool(run.font.bold)
    except (AttributeError, TypeError):
        pass

    return para_default


def _resolve_paragraph_font_size(paragraph) -> float | None:
    """Resolve the default font size at the paragraph level.

    Args:
        paragraph: A python-pptx Paragraph object.

    Returns:
        Font size in points (float) if explicitly set, None otherwise.
    """
    try:
        if paragraph.font.size is not None:
            return paragraph.font.size.pt
    except (AttributeError, TypeError):
        pass
    return None


def _resolve_paragraph_font_bold(paragraph) -> bool:
    """Resolve the default font weight at the paragraph level.

    Args:
        paragraph: A python-pptx Paragraph object.

    Returns:
        True if the paragraph default is bold, False otherwise.
    """
    try:
        if paragraph.font.bold is not None:
            return bool(paragraph.font.bold)
    except (AttributeError, TypeError):
        pass
    return False


def _resolve_font_size(paragraph) -> float | None:
    """Resolve the effective font size for a paragraph (best-effort).

    Checks paragraph-level default first, then first run's font size.

    Args:
        paragraph: A python-pptx Paragraph object.

    Returns:
        Font size in points (float) or None.
    """
    para_size = _resolve_paragraph_font_size(paragraph)
    if para_size is not None:
        return para_size

    # Check first run
    if paragraph.runs:
        try:
            first_run = paragraph.runs[0]
            if first_run.font.size is not None:
                return first_run.font.size.pt
        except (AttributeError, TypeError, IndexError):
            pass

    return None


__all__ = [
    "compute_spec_typography_sizes",
    "extract_all_slide_fonts",
    "extract_slide_fonts",
    "median",
]
