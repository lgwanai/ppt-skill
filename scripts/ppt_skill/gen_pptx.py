#!/usr/bin/env python3
"""从 JSON 文件生成 PPT，无需 LLM"""
import json, sys, os
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SHAPE_MAP = {
    "RECTANGLE": MSO_SHAPE.RECTANGLE, "RECT": MSO_SHAPE.RECTANGLE,
    "OVAL": MSO_SHAPE.OVAL, "ELLIPSE": MSO_SHAPE.OVAL,
    "ROUNDED_RECTANGLE": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ROUNDED_RECT": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ROUNDRECT": MSO_SHAPE.ROUNDED_RECTANGLE,
    "BLOCKARC": MSO_SHAPE.OVAL,
    "DONUT": MSO_SHAPE.DONUT,
    "TRIANGLE": MSO_SHAPE.ISOSCELES_TRIANGLE,
}

def parse_color(c):
    if not c or len(c) != 6:
        return None
    try:
        int(c, 16)
        return RGBColor.from_string(c)
    except:
        return None

def build_ppt(json_path, output='regenerated.pptx'):
    with open(json_path, 'r', encoding='utf-8') as f:
        elements = json.load(f)
    
    # Collect chart references for post-processing
    chart_items = [el for el in elements if el.get('type') == 'chart_frame' and el.get('chart_ref')]
    
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    for el in elements:
        el_type = el.get('type', '')
        if el_type == '_background':
            continue
        
        name = el.get('name', '')
        left = Emu(el.get('left', 0))
        top = Emu(el.get('top', 0))
        w = Emu(el.get('width', 914400))
        h = Emu(el.get('height', 914400))
        
        if el_type == 'picture':
            img = el.get('image_file', '')
            if img and os.path.exists(img):
                pic = slide.shapes.add_picture(img, left, top, w, h)
                pic.name = name
                # Apply image effects
                effects = el.get('img_effects')
                if effects:
                    blip = pic._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if blip is not None:
                        alpha = effects.get('alpha')
                        if alpha:
                            amf = etree.SubElement(blip, qn('a:alphaModFix'))
                            amf.set('amt', str(alpha))
                        clr = effects.get('clrChange')
                        if clr:
                            cc = etree.SubElement(blip, qn('a:clrChange'))
                            cf = etree.SubElement(cc, qn('a:clrFrom'))
                            fc = etree.SubElement(cf, qn('a:srgbClr'))
                            fc.set('val', clr.get('from', 'EFEFED'))
                            a1 = etree.SubElement(fc, qn('a:alpha'))
                            a1.set('val', '100000')
                            ct = etree.SubElement(cc, qn('a:clrTo'))
                            tc = etree.SubElement(ct, qn('a:srgbClr'))
                            tc.set('val', clr.get('to', 'EFEFED'))
                            a2 = etree.SubElement(tc, qn('a:alpha'))
                            a2.set('val', '100000')
        
        elif el_type == 'textbox':
            txBox = slide.shapes.add_textbox(left, top, w, h)
            txBox.name = name
            tf = txBox.text_frame
            tf.word_wrap = True
            
            autofit = el.get('autofit', '')
            if autofit == 'shrink':
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Vertical anchor (t=top, ctr=middle, b=bottom)
            v_anchor = el.get('v_anchor', '')
            if v_anchor:
                txBody = txBox._element.find(qn('p:txBody'))
                if txBody is not None:
                    bodyPr = txBody.find(qn('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('anchor', v_anchor)
            
            for i, run_info in enumerate(el.get('text_runs', [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                align = (run_info.get('alignment') or 'LEFT').upper()
                amap = {'CENTER': PP_ALIGN.CENTER, 'LEFT': PP_ALIGN.LEFT, 'RIGHT': PP_ALIGN.RIGHT,
                        'DISTRIBUTED': PP_ALIGN.DISTRIBUTE, 'JUSTIFY': PP_ALIGN.JUSTIFY}
                if align in amap:
                    p.alignment = amap[align]
                run = p.add_run()
                run.text = run_info.get('text', '')
                fs = run_info.get('font_size')
                if fs:
                    run.font.size = Emu(int(fs) * 12700)
                fn = run_info.get('font_name')
                if fn:
                    run.font.name = fn
                if run_info.get('bold'):
                    run.font.bold = True
                c = parse_color(run_info.get('color'))
                if c:
                    run.font.color.rgb = c
        
        elif el_type == 'auto_shape':
            subtype = el.get('shape_subtype', 'RECTANGLE')
            st = SHAPE_MAP.get(subtype, MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(st, left, top, w, h)
            shape.name = name
            
            # Fix blockArc: change prst from 'ellipse' to 'blockArc'
            if subtype == 'BLOCKARC':
                prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
                if prstGeom is not None:
                    prstGeom.set('prst', 'blockArc')
            
            # Rotation
            rot = el.get('rotation')
            if rot:
                shape.rotation = float(rot)
            # Flip
            flip = el.get('flip', '')
            if flip:
                spPr = shape._element.spPr
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    parts = flip.split(',')
                    if parts[0] == '1': xfrm.set('flipH', '1')
                    if len(parts) > 1 and parts[1] == '1': xfrm.set('flipV', '1')
            
            # Fill
            fill_type = el.get('fill_type', '无')
            fill_color = el.get('fill_color')
            gradient = el.get('gradient_stops')
            
            if fill_type == '纯色' and fill_color:
                shape.fill.solid()
                c = parse_color(fill_color)
                if c:
                    shape.fill.fore_color.rgb = c
            elif fill_type == '渐变' and gradient:
                has_alpha = any('alpha' in s for s in gradient)
                _set_gradient(shape, gradient, has_alpha)
            else:
                shape.fill.background()
            
            # Border
            border = el.get('border')
            if border and border.get('color'):
                shape.line.fill.solid()
                c = parse_color(border['color'])
                if c:
                    shape.line.color.rgb = c
                if border.get('width'):
                    shape.line.width = Emu(border['width'])
            else:
                try: shape.line.fill.background()
                except: pass
            
            # Shadow
            shadow = el.get('shadow')
            if shadow and '/' in shadow:
                blur, dist = shadow.split('/')
                _add_shadow(shape, blur, dist)
            
            # adj values (corner radius for roundRect, arc params for blockArc)
            adj = el.get('adj')
            if adj:
                if isinstance(adj, dict):
                    # Multiple adj values (e.g. blockArc: adj1/adj2/adj3)
                    _set_adj_values(shape, adj)
                elif 'ROUND' in subtype:
                    _set_adj_values(shape, {'adj': adj})
            
            # Handle text inside auto_shape (e.g. numbered circles "01")
            text_runs = el.get('text_runs', [])
            if text_runs:
                tf = shape.text_frame
                tf.word_wrap = True
                # Clear default margins added by add_shape
                txBody = shape._element.find(qn('p:txBody'))
                if txBody is not None:
                    bodyPr = txBody.find(qn('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('lIns', '0')
                        bodyPr.set('tIns', '0')
                        bodyPr.set('rIns', '0')
                        bodyPr.set('bIns', '0')
                v_anchor = el.get('v_anchor', '')
                if v_anchor:
                    txBody = shape._element.find(qn('p:txBody'))
                    if txBody is not None:
                        bodyPr = txBody.find(qn('a:bodyPr'))
                        if bodyPr is not None:
                            bodyPr.set('anchor', v_anchor)
                for i, run_info in enumerate(text_runs):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    align = (run_info.get('alignment') or 'CENTER').upper()
                    amap = {'CENTER': PP_ALIGN.CENTER, 'LEFT': PP_ALIGN.LEFT, 'RIGHT': PP_ALIGN.RIGHT}
                    if align in amap: p.alignment = amap[align]
                    run = p.add_run()
                    run.text = run_info.get('text', '')
                    fs = run_info.get('font_size')
                    if fs: run.font.size = Emu(int(fs) * 12700)
                    fn = run_info.get('font_name')
                    if fn: run.font.name = fn
                    if run_info.get('bold'): run.font.bold = True
                    c = parse_color(run_info.get('color'))
                    if c: run.font.color.rgb = c
        
        elif el_type == 'raw_sp':
            try:
                sp_xml = el.get('xml', '')
                if sp_xml:
                    sp = etree.fromstring(sp_xml.encode('utf-8'))
                    # Update position to absolute coordinates
                    xfrm = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                    if xfrm is not None:
                        off = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                        ext = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                        if off is not None:
                            off.set('x', str(el.get('left', 0)))
                            off.set('y', str(el.get('top', 0)))
                        if ext is not None:
                            ext.set('cx', str(el.get('width', 1)))
                            ext.set('cy', str(el.get('height', 1)))
                    slide.shapes._spTree.append(sp)
            except Exception as e:
                 print(f"  skip {el.get('name')}: {e}")
        
        elif el_type == 'chart_frame':
            # Save chart frame for post-processing (copy from original)
            pass
        
        elif el_type == 'connector':
            try:
                from pptx.oxml import parse_xml
                cxn = parse_xml(
                    f'<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                    f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    f'<p:nvCxnSpPr><p:cNvPr id="0" name="{name}"/>'
                    f'<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
                    f'<p:spPr><a:xfrm><a:off x="{left}" y="{top}"/>'
                    f'<a:ext cx="{w}" cy="{h}"/></a:xfrm>'
                    f'<a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>'
                    f'<a:ln w="12700"><a:solidFill><a:srgbClr val="0445FE"/>'
                    f'</a:solidFill></a:ln></p:spPr></p:cxnSp>'
                )
                slide.shapes._spTree.append(cxn)
            except:
                pass
    
    prs.save(output)
    
    # Post-process: copy chart files from source pptx
    if chart_items:
        _copy_charts(output, chart_items, json_path)
    
    print(f'PPT 已保存: {output}')

def _set_gradient(shape, stops, has_alpha):
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split('}')[-1]
        if tag in ('solidFill', 'noFill', 'gradFill', 'blipFill'):
            spPr.remove(child)
    gf = etree.SubElement(spPr, qn('a:gradFill'))
    gf.set('flip', 'none')
    gf.set('rotWithShape', '1')
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    for s in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', s.get('pos', '0'))
        sc = etree.SubElement(gs, qn('a:srgbClr'))
        sc.set('val', s.get('color', '0445FE'))
        if has_alpha and 'alpha' in s:
            al = etree.SubElement(sc, qn('a:alpha'))
            al.set('val', s['alpha'])
    lin = etree.SubElement(gf, qn('a:lin'))
    lin.set('ang', '0')
    lin.set('scaled', '1')

def _set_adj_values(shape, adj_dict):
    """Set adj values (corner radius for roundRect, arc params for blockArc)"""
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    else:
        for gd in list(avLst.findall(qn('a:gd'))):
            avLst.remove(gd)
    for name, val in adj_dict.items():
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', name)
        gd.set('fmla', f'val {val}')

def _add_shadow(shape, blur, dist):
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}effectLst'):
            spPr.remove(child)
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', blur)
    sh.set('dist', dist)
    sh.set('dir', '5400000')
    sh.set('algn', 't')
    sh.set('rotWithShape', '0')
    sc = etree.SubElement(sh, qn('a:schemeClr'))
    sc.set('val', 'accent1')
    lm = etree.SubElement(sc, qn('a:lumMod'))
    lm.set('val', '50000')
    al = etree.SubElement(sc, qn('a:alpha'))
    al.set('val', '10000')

if __name__ == '__main__':
    input_json = sys.argv[1] if len(sys.argv) > 1 else 'ppt_page2.json'
    output_ppt = sys.argv[2] if len(sys.argv) > 2 else 'regenerated.pptx'
    build_ppt(input_json, output_ppt)
