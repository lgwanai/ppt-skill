#!/usr/bin/env python3
"""读取 course_content_pages.json，逐页合并模板+内容，生成多页 PPT"""
import json, sys, os, re, requests, glob, copy
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from ppt_skill.config import load_config, get_llm_config
_cfg = load_config()
_llm = get_llm_config(_cfg)
LLM_MODEL = _llm["model"]
LLM_API_KEY = _llm["api_key"] or os.environ.get("LLM_API_KEY", "")
LLM_API_BASE = _llm["api_base"] or os.environ.get("LLM_API_BASE", "")

SHAPE_MAP = {
    "RECTANGLE": MSO_SHAPE.RECTANGLE, "RECT": MSO_SHAPE.RECTANGLE,
    "OVAL": MSO_SHAPE.OVAL, "ELLIPSE": MSO_SHAPE.OVAL,
    "ROUNDED_RECTANGLE": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ROUNDED_RECT": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ROUNDRECT": MSO_SHAPE.ROUNDED_RECTANGLE,
    "BLOCKARC": MSO_SHAPE.OVAL,
    "DONUT": MSO_SHAPE.DONUT,
    "TRIANGLE": MSO_SHAPE.ISOSCELES_TRIANGLE,
}

def call_llm(system, user):
    resp = requests.post(
        f"{LLM_API_BASE}/chat/completions",
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {LLM_API_KEY}"},
        json={"model": LLM_MODEL, "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user}],
            "temperature": 0.1, "max_tokens": 64000},
        timeout=300)
    resp.raise_for_status()
    raw = resp.json()["choices"][0]["message"]["content"]
    raw = re.sub(r"^```(?:json)?\s*", "", raw.strip())
    raw = re.sub(r"\s*```$", "", raw)
    return raw

def build_page_map():
    mapping = {}
    for f in glob.glob("ppt_demo1*.json"):
        if f == "ppt_demo1.json": mapping[1] = f
        else:
            m = re.search(r'page(\d+)', f)
            if m: mapping[int(m.group(1))] = f
    return mapping

def merge_page(template_elements, content_page):
    clean_tpl = [e for e in template_elements if e.get('type') not in ('_background', 'raw_sp', 'chart_frame')]
    item_count = content_page.get('item_count', 1)
    
    system = """你是 PPT 文本填充专家。模板有几组你用几组，不扩不缩。只替换文字，模板原文一字不留。不改任何几何属性。输出完整 JSON 数组。"""

    user = f"""模板 JSON:
{json.dumps(clean_tpl, ensure_ascii=False)[:30000]}

新内容 ({item_count}项):
{json.dumps(content_page, ensure_ascii=False)}"""

    raw = call_llm(system, user)
    try: merged = json.loads(raw)
    except:
        m = re.search(r'\[[\s\S]*\]', raw)
        merged = json.loads(m.group(0)) if m else None
    if merged:
        merged = _fit_to_page(merged)
    return merged

def _count_rows(elements):
    min_y = min((e.get('top', 0) for e in elements if e.get('top', 0) > 0), default=0)
    circles = sum(1 for e in elements if e.get('shape_subtype') in ('OVAL','ELLIPSE') and e.get('top', 0) > min_y + 500000)
    return max(circles, 1)

def _fit_to_page(elements):
    """Python 几何修正：缩放+平移让所有元素在页面内"""
    PAGE_W, PAGE_H = 12192000, 6858000
    xs = [e.get('left', 0) + e.get('width', 0) for e in elements if e.get('width')]
    ys = [e.get('top', 0) + e.get('height', 0) for e in elements if e.get('height')]
    if not ys: return elements
    
    max_x = max(xs) if xs else PAGE_W
    max_y = max(ys) if ys else PAGE_H
    
    if max_y <= PAGE_H and max_x <= PAGE_W:
        return elements  # already fits
    
    # Scale down to fit
    scale = min(PAGE_W / max(max_x, 1), PAGE_H / max(max_y, 1), 1.0)
    for el in elements:
        if el.get('left', 0) > 0: el['left'] = int(el['left'] * scale)
        if el.get('top', 0) > 0: el['top'] = int(el['top'] * scale)
        if el.get('width'): el['width'] = int(el['width'] * scale)
        if el.get('height'): el['height'] = int(el['height'] * scale)
        if el.get('font_size'): el['font_size'] = int(el['font_size'] * scale)
        # Keep circles circular
        if el.get('shape_subtype') in ('OVAL', 'ELLIPSE'):
            mn = min(el['width'], el['height'])
            el['width'] = mn
            el['height'] = mn
    return elements

def _count_rows(elements):
    """Count content rows by counting circles (each content group has a circle)"""
    min_y = min((e.get('top', 0) for e in elements if e.get('top', 0) > 0), default=0)
    circles = sum(1 for e in elements if e.get('shape_subtype') in ('OVAL','ELLIPSE') and e.get('top', 0) > min_y + 500000)
    return max(circles, 1)

PAGE_W, PAGE_H = 12192000, 6858000

def _clip_to_bounds(elements):
    for el in elements:
        l, t, w, h = el.get('left', 0), el.get('top', 0), el.get('width', 0), el.get('height', 0)
        if w and l + w > PAGE_W: el['width'] = PAGE_W - l - 100000
        if h and t + h > PAGE_H: el['height'] = PAGE_H - t - 100000
        if el.get('left', 0) < 0: el['left'] = 0
        if el.get('top', 0) < 0: el['top'] = 0
    return elements

def parse_color(c):
    if not c or len(c) != 6: return None
    try: int(c, 16); return RGBColor.from_string(c)
    except: return None

def add_element_to_slide(slide, el):
    t = el.get("type", "")
    name = el.get("name", "")
    left, top = Emu(el.get("left", 0)), Emu(el.get("top", 0))
    w, h = Emu(el.get("width", 914400)), Emu(el.get("height", 914400))
    
    if t == '_background' or t == 'chart_frame':
        return
    if t == 'picture':
        img = el.get("image_file", "")
        if img and os.path.exists(img):
            pic = slide.shapes.add_picture(img, left, top, w, h)
            pic.name = name
        return
    if t == 'raw_sp':
        try:
            sp = etree.fromstring(el.get('xml','').encode('utf-8'))
            xfrm = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                off = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                ext = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                if off is not None: off.set('x', str(el.get('left',0))); off.set('y', str(el.get('top',0)))
                if ext is not None: ext.set('cx', str(el.get('width',1))); ext.set('cy', str(el.get('height',1)))
            slide.shapes._spTree.append(sp)
        except: pass
        return
    
    if t == 'textbox':
        txBox = slide.shapes.add_textbox(left, top, w, h)
        txBox.name = name
        tf = txBox.text_frame; tf.word_wrap = True
        va = el.get('v_anchor', '')
        if va:
            txBody = txBox._element.find(qn('p:txBody'))
            if txBody is not None:
                bp = txBody.find(qn('a:bodyPr'))
                if bp is not None: bp.set('anchor', va)
        if el.get('autofit') == 'shrink': tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for i, ri in enumerate(el.get('text_runs', [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            al = (ri.get('alignment') or 'LEFT').upper()
            am = {'CENTER': PP_ALIGN.CENTER, 'LEFT': PP_ALIGN.LEFT, 'RIGHT': PP_ALIGN.RIGHT}
            if al in am: p.alignment = am[al]
            run = p.add_run()
            run.text = ri.get('text', '')
            if ri.get('font_size'): run.font.size = Emu(int(ri['font_size']) * 12700)
            if ri.get('font_name'): run.font.name = ri['font_name']
            if ri.get('bold'): run.font.bold = True
            c = parse_color(ri.get('color'))
            if c: run.font.color.rgb = c
        return
    
    if t == 'auto_shape':
        st = SHAPE_MAP.get(el.get('shape_subtype', 'RECTANGLE'), MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(st, left, top, w, h)
        shape.name = name
        if el.get('rotation'): shape.rotation = float(el['rotation'])
        fc = parse_color(el.get('fill_color'))
        if fc:
            if el.get('fill_gradient'):
                _set_gradient(shape, el.get('gradient_stops', []), False)
            else:
                shape.fill.solid(); shape.fill.fore_color.rgb = fc
        else: shape.fill.background()
        border = el.get('border')
        if border and border.get('color'):
            shape.line.fill.solid()
            bc = parse_color(border['color'])
            if bc: shape.line.color.rgb = bc
            if border.get('width'): shape.line.width = Emu(border['width'])
        else:
            try: shape.line.fill.background()
            except: pass
        if el.get('shadow'):
            parts = el['shadow'].split('/')
            _add_shadow(shape, parts[0], parts[1] if len(parts) > 1 else '38100')
        adj = el.get('adj')
        if adj:
            prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
            if prstGeom is not None:
                avLst = prstGeom.find(qn('a:avLst'))
                if avLst is None: avLst = etree.SubElement(prstGeom, qn('a:avLst'))
                if isinstance(adj, dict):
                    for n, v in adj.items():
                        gd = etree.SubElement(avLst, qn('a:gd')); gd.set('name', n); gd.set('fmla', f'val {v}')
                else:
                    gd = etree.SubElement(avLst, qn('a:gd')); gd.set('name', 'adj'); gd.set('fmla', f'val {adj}')
        if el.get('shape_subtype') == 'BLOCKARC':
            prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
            if prstGeom is not None: prstGeom.set('prst', 'blockArc')
        flip = el.get('flip', '')
        if flip:
            xfrm = shape._element.spPr.find(qn('a:xfrm'))
            parts = flip.split(',')
            if parts[0] == '1': xfrm.set('flipH', '1')
            if len(parts) > 1 and parts[1] == '1': xfrm.set('flipV', '1')
        for ri in el.get('text_runs', []):
            tf = shape.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            al = (ri.get('alignment') or 'CENTER').upper()
            am = {'CENTER': PP_ALIGN.CENTER, 'LEFT': PP_ALIGN.LEFT, 'RIGHT': PP_ALIGN.RIGHT}
            if al in am: p.alignment = am[al]
            run = p.add_run(); run.text = ri.get('text', '')
            if ri.get('font_size'): run.font.size = Emu(int(ri['font_size']) * 12700)
            if ri.get('font_name'): run.font.name = ri['font_name']
            if ri.get('bold'): run.font.bold = True
            c = parse_color(ri.get('color'))
            if c: run.font.color.rgb = c

def _set_gradient(shape, stops, has_alpha):
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split('}')[-1]
        if tag in ('solidFill', 'noFill', 'gradFill', 'blipFill'): spPr.remove(child)
    gf = etree.SubElement(spPr, qn('a:gradFill')); gf.set('flip', 'none'); gf.set('rotWithShape', '1')
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    for s in stops:
        gs = etree.SubElement(gsLst, qn('a:gs')); gs.set('pos', s.get('pos', '0'))
        sc = etree.SubElement(gs, qn('a:srgbClr')); sc.set('val', s.get('color', '0445FE'))
        if has_alpha and 'alpha' in s:
            al = etree.SubElement(sc, qn('a:alpha')); al.set('val', s['alpha'])
    lin = etree.SubElement(gf, qn('a:lin')); lin.set('ang', '0'); lin.set('scaled', '1')

def _add_shadow(shape, blur, dist):
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}effectLst'): spPr.remove(child)
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', blur); sh.set('dist', dist)
    sh.set('dir', '5400000'); sh.set('algn', 't'); sh.set('rotWithShape', '0')
    sc = etree.SubElement(sh, qn('a:schemeClr')); sc.set('val', 'accent1')
    lm = etree.SubElement(sc, qn('a:lumMod')); lm.set('val', '50000')
    al = etree.SubElement(sc, qn('a:alpha')); al.set('val', '10000')

def main(pages_path, output_ppt="course.pptx", start=0, end=None):
    with open(pages_path, 'r', encoding='utf-8') as f:
        all_pages = json.load(f)
    
    if end is not None:
        pages = all_pages[start:end]
    else:
        pages = all_pages[start:]
    
    page_map = build_page_map()
    print(f"模板: {sorted(page_map.keys())}")
    print(f"页面: {start}-{start+len(pages)-1} (共 {len(pages)} 页)")
    print()
    
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    
    for i, pg in enumerate(pages):
        tp = pg.get("template_page")
        title = pg.get("title", "")[:40]
        if tp is None or tp not in page_map:
            print(f"  [{i:02d}] ⚠️ 跳过 [{tp}] {title}")
            continue
        
        with open(page_map[tp], 'r', encoding='utf-8') as f:
            tpl = json.load(f)
        
        # Calculate template slots
        tpl_slots = _count_rows([e for e in tpl if e.get('type') not in ('_background','raw_sp','chart_frame')])
        item_count = pg.get('item_count', 1)
        chunks = item_count // max(tpl_slots, 1) + (1 if item_count % max(tpl_slots, 1) else 0)
        
        if chunks > 1:
            print(f"  [{i:02d}] 📄 [{tp}] {title}  ({item_count}项/{tpl_slots}槽 → {chunks}页)")
            # Split content into chunks
            content_text = pg.get('content', '')
            items = _extract_items(content_text)
            for c in range(chunks):
                chunk_start = c * tpl_slots
                chunk_items = items[chunk_start:chunk_start + tpl_slots]
                chunk_pg = {**pg, 'title': f'{title} ({c+1}/{chunks})',
                           'content': '\n'.join(chunk_items),
                           'item_count': len(chunk_items)}
                merged = merge_page(tpl, chunk_pg)
                if merged:
                    for e in tpl:
                        if e.get('type') in ('_background', 'raw_sp', 'chart_frame'):
                            merged.append(e)
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    for el in merged:
                        add_element_to_slide(slide, el)
        else:
            print(f"  [{i:02d}] 🔄 [{tp}] {title}")
            merged = merge_page(tpl, pg)
            if merged is None:
                print(f"       ❌ 跳过")
                continue
            for e in tpl:
                if e.get('type') in ('_background', 'raw_sp', 'chart_frame'):
                    merged.append(e)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            for el in merged:
                add_element_to_slide(slide, el)

def _extract_items(text):
    """Split content into numbered list items only"""
    items = re.split(r'\n(?=\d+[\.\)] )', text)
    items = [i.strip() for i in items if i.strip() and i.strip()[0].isdigit()]
    return items
    
    prs.save(output_ppt)
    print(f"\n✅ {output_ppt} ({len(prs.slides)} 页)")

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法: python merge_and_generate.py <pages.json> [起始页] [结束页]")
        print("  例: python merge_and_generate.py course_content_pages.json        (全部)")
        print("  例: python merge_and_generate.py course_content_pages.json 0 5   (第0-4页)")
        sys.exit(1)
    start = int(sys.argv[2]) if len(sys.argv) > 2 else 0
    end = int(sys.argv[3]) if len(sys.argv) > 3 else None
    out = sys.argv[4] if len(sys.argv) > 4 else f"course_p{start}-{end if end else 'end'}.pptx"
    main(sys.argv[1], out, start, end)
