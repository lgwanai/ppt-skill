"""Style comparison on PPTX shapes — no SVG, no regex.

Compares generated PPTX slide elements against spec data:
  - Font family, size, weight, color
  - Element position (IoU-based)
  - Background color
  - Overall similarity score

All checks run on python-pptx Shape objects, not SVG text.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any


@dataclass
class StyleReport:
    overall: float = 0.0
    font_score: float = 0.0
    size_score: float = 0.0
    color_score: float = 0.0
    position_score: float = 0.0
    background_score: float = 0.0
    issues: list[str] = field(default_factory=list)
    passed: bool = False

    WEIGHTS = {"font": 0.25, "size": 0.25, "color": 0.25, "position": 0.15, "background": 0.10}


def _hex_to_rgb(h: str) -> tuple[int, int, int] | None:
    h = h.lstrip("#")
    if len(h) != 6: return None
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _color_distance(c1: str, c2: str) -> float:
    """Perceptual color distance 0-1. Lower = more similar."""
    r1 = _hex_to_rgb(c1); r2 = _hex_to_rgb(c2)
    if not r1 or not r2: return 1.0
    dr = (r1[0] - r2[0]) / 255.0; dg = (r1[1] - r2[1]) / 255.0; db = (r1[2] - r2[2]) / 255.0
    return (dr * dr * 0.299 + dg * dg * 0.587 + db * db * 0.114) ** 0.5


def evaluate_slide(slide, spec_page: dict, slide_w: int, slide_h: int) -> StyleReport:
    """Compare generated PPTX slide against spec page data.

    Args:
        slide: python-pptx Slide object (generated)
        spec_page: dict from spec YAML (expected style)
        slide_w, slide_h: EMU dimensions

    Returns StyleReport with scores and issues.
    """
    report = StyleReport()
    expected_elements = spec_page.get("elements", [])
    if not expected_elements:
        report.passed = True
        return report

    generated_texts: list[dict] = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame') or not shape.text_frame.text.strip():
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                try: c = str(r.font.color.rgb)
                except: c = ""
                generated_texts.append({
                    "text": r.text[:30],
                    "font": r.font.name or "",
                    "size": r.font.size / 12700 if r.font.size else 0,
                    "bold": r.font.bold or False,
                    "color": c,
                    "x": (shape.left or 0) / slide_w,
                    "y": (shape.top or 0) / slide_h,
                    "w": (shape.width or 0) / slide_w,
                    "h": (shape.height or 0) / slide_h,
                })
                break
            break

    if not generated_texts:
        report.issues.append("No text elements found in generated slide")
        return report

    # Match generated elements to expected by position proximity
    matched_pairs = []
    for exp in expected_elements:
        ts = exp.get("text_style", {}) or {}
        pos = exp.get("position", {})
        exp_font = ts.get("font_family", "")
        exp_size = ts.get("font_size_pt", 0) or 0
        exp_color = ts.get("font_color", "")
        exp_x = pos.get("x", 0); exp_y = pos.get("y", 0)
        role = exp.get("role", "body")

        # Find closest generated element by position
        best_gen = None; best_dist = float("inf")
        for g in generated_texts:
            dx = abs(g["x"] - exp_x); dy = abs(g["y"] - exp_y)
            dist = dx + dy
            if dist < best_dist:
                best_dist = dist; best_gen = g

        # Only match if within reasonable distance (0.3 normalized)
        if best_gen and best_dist < 0.3:
            matched_pairs.append((exp, best_gen))

    # Score matched pairs
    font_scores = []; size_scores = []; color_scores = []; pos_scores = []
    for exp, gen in matched_pairs:
        ts = exp.get("text_style", {}) or {}
        pos = exp.get("position", {})
        exp_font = ts.get("font_family", "")
        exp_size = ts.get("font_size_pt", 0) or 0
        exp_color = ts.get("font_color", "")
        exp_x = pos.get("x", 0); exp_y = pos.get("y", 0)
        if gen["font"] and exp_font:
            gf = gen["font"].lower(); ef = exp_font.lower()
            if gf == ef or ef in gf or gf in ef:
                font_scores.append(1.0)
            else:
                font_scores.append(0.3)
                report.issues.append(f"Font mismatch: expected '{exp_font}', got '{gen['font']}'")
        else:
            font_scores.append(0.5)

        # Font size match
        if gen["size"] > 0 and exp_size > 0:
            ratio = min(gen["size"], exp_size) / max(gen["size"], exp_size)
            size_scores.append(ratio)
            if ratio < 0.85:
                report.issues.append(f"Size mismatch: expected {exp_size:.0f}pt, got {gen['size']:.0f}pt")
        else:
            size_scores.append(0.5)

        # Color match
        if gen["color"] and exp_color:
            dist = _color_distance(gen["color"], exp_color)
            cs = max(0, 1 - dist * 2)
            color_scores.append(cs)
            if cs < 0.7:
                report.issues.append(f"Color mismatch: expected {exp_color}, got {gen['color']}")
        else:
            color_scores.append(0.5)

        # Position match
        if gen["x"] > 0:
            x_diff = abs(gen["x"] - exp_x)
            y_diff = abs(gen["y"] - exp_y)
            max_diff = max(x_diff, y_diff)
            ps = max(0, 1 - max_diff * 5)  # 5 = tolerance factor
            pos_scores.append(ps)
            if ps < 0.5:
                report.issues.append(
                    f"Position mismatch: expected ({exp_x:.3f},{exp_y:.3f}), got ({gen['x']:.3f},{gen['y']:.3f})")
        else:
            pos_scores.append(0.5)

    # Background check
    bg_color = spec_page.get("background_color", "")
    if bg_color:
        # Check first large rect or background of the slide
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'fill'):
                    from pptx.enum.dml import MSO_FILL_TYPE
                    if shape.fill.type == MSO_FILL_TYPE.SOLID and shape.fill.fore_color.rgb:
                        shape_color = str(shape.fill.fore_color.rgb)
                        dist = _color_distance(shape_color, bg_color)
                        report.background_score = max(0, 1 - dist * 2)
                        break
            except Exception:
                pass
    else:
        report.background_score = 1.0

    w = StyleReport.WEIGHTS
    report.font_score = sum(font_scores) / len(font_scores) if font_scores else 0.5
    report.size_score = sum(size_scores) / len(size_scores) if size_scores else 0.5
    report.color_score = sum(color_scores) / len(color_scores) if color_scores else 0.5
    report.position_score = sum(pos_scores) / len(pos_scores) if pos_scores else 0.5

    report.overall = (
        report.font_score * w["font"] + report.size_score * w["size"]
        + report.color_score * w["color"] + report.position_score * w["position"]
        + report.background_score * w["background"]
    )
    report.passed = report.overall >= 0.85
    return report


def evaluate_slide_vs_spec(prs, slide_index: int, spec_page: dict) -> StyleReport:
    """Evaluate one slide in a presentation against its spec."""
    slides = list(prs.slides)
    if slide_index >= len(slides):
        return StyleReport(issues=[f"Slide {slide_index} not found"])
    return evaluate_slide(slides[slide_index], spec_page, prs.slide_width, prs.slide_height)
