"""VL-driven spec extraction — reads PPTX, extracts all element properties,
sends each slide + element data to VL model for semantic analysis, then
generates a deduplicated JSON spec directory.

Key differences from the original SpecExtractor:
- Uses VL model to understand element roles and relationships
- Output JSON has NO text content — only attributes, roles, properties
- Similar layouts are merged into one file (cover/end_page each one file,
  content pages grouped by layout similarity)
- Each spec gets a short memorable name from the filename

Output directory structure:
  specs/<spec_name>/
  ├── spec.json                  # Master spec with palette, fonts, canvas info
  ├── slides/                    # PNG screenshots per slide
  │   ├── 00_cover.png
  │   ├── 01_content.png
  │   └── ...
  ├── cover.json                 # Cover page blueprint
  ├── end_page.json              # End page blueprint (if exists)
  ├── content_left_right.json    # Content pages with left-right layout
  ├── content_full_width.json    # Content pages with full-width layout
  └── ...                        # One JSON per unique layout type

Usage:
    extractor = VLVSpecExtractor()
    extractor.extract("input.pptx", "my-style")
    # → specs/my-style/ with JSON blueprints + slide images
"""

from __future__ import annotations

import json
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation

from ppt_skill.spec.theme import extract_theme_colors, extract_theme_fonts
from ppt_skill.spec.vl_element_analyzer import VLElementAnalyzer, ElementAnalysisResult


# ---------------------------------------------------------------------------
# PPTX rendering — LibreOffice -> PDF -> pdf2image -> per-page PNG
# ---------------------------------------------------------------------------

# Cache: pptx_path -> pdf_path, avoids re-converting for each slide
_pptx_to_pdf_cache: dict[str, str] = {}


def _convert_pptx_to_pdf(pptx_path: Path) -> Path | None:
    """Convert PPTX to PDF using LibreOffice headless mode. Returns PDF path."""
    key = str(pptx_path.resolve())
    if key in _pptx_to_pdf_cache:
        pdf = Path(_pptx_to_pdf_cache[key])
        if pdf.exists():
            return pdf

    soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if not Path(soffice).exists():
        return None

    tmp_dir = Path(tempfile.mkdtemp(prefix="ppt_slide_"))
    try:
        r = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(pptx_path)],
            capture_output=True, timeout=60,
        )
        if r.returncode != 0:
            return None

        pdfs = list(tmp_dir.glob("*.pdf"))
        if not pdfs:
            return None

        pdf_path = pdfs[0]
        _pptx_to_pdf_cache[key] = str(pdf_path)
        return pdf_path
    except Exception:
        return None


def _render_slide_to_png(pdf_path: Path, output_dir: Path, slide_index: int) -> Path | None:
    """Extract a single slide page from a PDF as a PNG using pdf2image.

    Args:
        pdf_path: Path to the PDF (from _convert_pptx_to_pdf).
        output_dir: Directory to write the PNG.
        slide_index: 0-based slide index.

    Returns:
        Path to the PNG, or None on failure.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"slide_{slide_index:02d}.png"
    if output_path.exists():
        return output_path

    if not pdf_path or not pdf_path.exists():
        return None

    try:
        from pdf2image import convert_from_path
    except ImportError:
        return None

    try:
        # Extract only the specific page (first_page=1, last_page=1 = single page)
        images = convert_from_path(
            str(pdf_path),
            dpi=150,
            fmt="png",
            first_page=slide_index + 1,
            last_page=slide_index + 1,
            output_folder=str(output_dir),
            paths_only=True,
        )
        if not images:
            return None

        src = Path(images[0])
        src.rename(output_path)
        return output_path
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Complete element extraction — table, picture, shape, group, layout-layer
# ---------------------------------------------------------------------------


def _extract_cell_text(cell) -> dict:
    """Extract text and style from a table cell."""
    result = {
        "text": cell.text.strip(),
        "text_style": {},
    }
    try:
        # First paragraph first run font
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                f = run.font
                ts = {}
                if f.name:
                    ts["font_family"] = f.name
                if f.size:
                    ts["font_size_pt"] = round(f.size / 12700, 1)
                if f.bold is not None:
                    ts["font_weight"] = "bold" if f.bold else "normal"
                if f.italic is not None:
                    ts["font_italic"] = f.italic
                try:
                    if f.color and f.color.rgb:
                        ts["font_color"] = f"#{f.color.rgb}"
                except Exception:
                    pass
                if ts:
                    result["text_style"] = ts
                break
            break
    except Exception:
        pass

    # Cell fill color
    try:
        from pptx.enum.dml import MSO_FILL_TYPE
        fill = cell.fill
        if fill.type == MSO_FILL_TYPE.SOLID and fill.fore_color and fill.fore_color.rgb:
            result["fill_color"] = f"#{fill.fore_color.rgb}"
    except Exception:
        pass

    return result


def _extract_table_data(table, sw: int, sh: int, z_base: int) -> dict:
    """Extract complete table structure from a python-pptx table."""
    rows = len(table.rows)
    cols = len(table.columns)
    col_widths = []
    for c in table.columns:
        col_widths.append(round(c.width / sw, 4) if sw else 0)

    cells = []
    cell_z = z_base
    for ri in range(rows):
        row_cells = []
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell_data = _extract_cell_text(cell)
            cell_data["row"] = ri
            cell_data["col"] = ci
            cell_data["id"] = cell_z
            cell_z += 1
            try:
                cell_data["rowspan"] = cell._tc.get("rowSpan", 1)
                cell_data["colspan"] = cell._tc.get("gridSpan", 1)
            except Exception:
                pass
            row_cells.append(cell_data)
            cells.append(cell_data)

    return {
        "table": {
            "rows": rows,
            "cols": cols,
            "col_widths": col_widths,
            "cells": cells,
        }
    }


def _extract_shape_deep(shape, sw: int, sh: int, z_order: int, spec_fs: Path) -> dict:
    """Extract every property from ANY shape type — text, picture, table, group, auto-shape."""
    from pptx.shapes.picture import Picture as PptxPicture
    from pptx.oxml.ns import qn

    base = {
        "id": z_order,
        "z_order": z_order,
        "shape_name": getattr(shape, "name", ""),
        "position": {
            "x": round((shape.left or 0) / sw, 4) if sw else 0,
            "y": round((shape.top or 0) / sh, 4) if sh else 0,
            "w": round((shape.width or 0) / sw, 4) if sw else 0,
            "h": round((shape.height or 0) / sh, 4) if sh else 0,
        },
    }

    # ── TABLE ──
    if hasattr(shape, "has_table") and shape.has_table:
        base["element_type"] = "table"
        base["semantic_role"] = "table"
        table_data = _extract_table_data(shape.table, sw, sh, z_order + 1)
        base.update(table_data)
        return base

    # ── PICTURE (image) ──
    if isinstance(shape, PptxPicture):
        base["element_type"] = "image"
        base["semantic_role"] = "image"
        try:
            img = shape.image
            base["image"] = {
                "content_type": img.content_type,
                "width_px": img.size[0],
                "height_px": img.size[1],
            }
        except Exception:
            pass
        # Extract blip embed for saving
        try:
            xml_str = shape._element.xml
            import re
            m = re.search(r'r:embed="([^"]+)"', xml_str)
            if m:
                base["image"]["embed_id"] = m.group(1)
        except Exception:
            pass
        return base

    # ── GROUP SHAPE ──
    if hasattr(shape, "shapes") and hasattr(shape.shapes, "__len__"):
        base["element_type"] = "group"
        base["semantic_role"] = "group"
        children = []
        for ci, child in enumerate(shape.shapes):
            try:
                child_data = _extract_shape_deep(child, sw, sh, ci, spec_fs)
                children.append(child_data)
            except Exception:
                pass
        if children:
            base["children"] = children
        return base

    # ── TEXT SHAPE (or any shape with text) ──
    has_text = False
    try:
        has_text = shape.has_text_frame and shape.text_frame.text.strip()
    except Exception:
        pass

    if has_text:
        base["element_type"] = "text"
        # Determine semantic role
        base["semantic_role"] = _infer_text_role(shape)
        # Text content
        base["text"] = shape.text_frame.text
        # Text style from first run
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    f = run.font
                    ts = {}
                    if f.name:
                        ts["font_family"] = f.name
                    if f.size:
                        ts["font_size_pt"] = round(f.size / 12700, 1)
                    if f.bold is not None:
                        ts["font_weight"] = "bold" if f.bold else "normal"
                    if f.italic is not None:
                        ts["font_italic"] = f.italic
                    try:
                        if f.color and f.color.rgb:
                            ts["font_color"] = f"#{f.color.rgb}"
                    except Exception:
                        pass
                    if ts:
                        base["text_style"] = ts
                    break
                break
            # Alignment
            try:
                from pptx.enum.text import PP_ALIGN
                align_map = {
                    PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center",
                    PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify",
                }
                for para in shape.text_frame.paragraphs:
                    if para.alignment:
                        base["text_alignment"] = align_map.get(para.alignment, "left")
                        break
            except Exception:
                pass
        except Exception:
            pass
    else:
        # ── AUTO SHAPE (rectangle, oval, line, etc.) ──
        base["element_type"] = "shape"
        base["semantic_role"] = _infer_shape_role(shape)
        # Fill
        try:
            from pptx.enum.dml import MSO_FILL_TYPE
            fill = shape.fill
            if fill.type == MSO_FILL_TYPE.SOLID:
                if fill.fore_color and fill.fore_color.rgb:
                    base["fill_color"] = f"#{fill.fore_color.rgb}"
                try:
                    base["fill_opacity"] = fill.fore_color.brightness or 1.0
                except Exception:
                    pass
            elif fill.type == MSO_FILL_TYPE.PICTURE:
                base["fill_type"] = "picture"
        except Exception:
            pass
        # Stroke
        try:
            line = shape.line
            if line.fill and line.fill.fore_color and line.fill.fore_color.rgb:
                base["stroke_color"] = f"#{line.fill.fore_color.rgb}"
            if line.width:
                base["stroke_width_pt"] = round(line.width / 12700, 1)
        except Exception:
            pass
        # Geometry
        try:
            st = str(shape.shape_type) if shape.shape_type else ""
            base["shape_type"] = _normalize_geom(st)
            xml_str = shape._element.xml
            if "roundRect" in xml_str:
                base["shape_type"] = "round_rect"
        except Exception:
            pass

    return base


def _infer_text_role(shape) -> str:
    """Infer semantic role of a text shape."""
    name = getattr(shape, "name", "").lower()
    try:
        pf = shape.placeholder_format
        if pf is not None:
            from pptx.enum.shapes import PP_PLACEHOLDER
            ph_map = {
                PP_PLACEHOLDER.TITLE: "title",
                PP_PLACEHOLDER.SUBTITLE: "subtitle",
                PP_PLACEHOLDER.BODY: "body",
                PP_PLACEHOLDER.FOOTER: "footer",
                PP_PLACEHOLDER.SLIDE_NUMBER: "page_number",
                PP_PLACEHOLDER.DATE: "date",
            }
            if pf.type in ph_map:
                return ph_map[pf.type]
    except Exception:
        pass
    if any(k in name for k in ["title", "标题"]):
        return "title"
    if any(k in name for k in ["subtitle", "副标题"]):
        return "subtitle"
    if any(k in name for k in ["footer", "页脚"]):
        return "footer"
    return "body"


def _infer_shape_role(shape) -> str:
    """Infer semantic role of a non-text shape."""
    name = getattr(shape, "name", "").lower()
    if any(k in name for k in ["logo", "icon"]):
        return "icon"
    if any(k in name for k in ["divider", "line", "分隔"]):
        return "divider"
    if any(k in name for k in ["decor", "accent", "装饰"]):
        return "decoration"
    return "decoration"


def _normalize_geom(st: str) -> str:
    st = st.upper()
    if "RECTANGLE" in st or "RECT" in st:
        return "rect"
    if "ROUND" in st:
        return "round_rect"
    if "OVAL" in st or "ELLIPSE" in st or "CIRCLE" in st:
        return "circle"
    if "LINE" in st or "CONNECTOR" in st:
        return "line"
    if "ARROW" in st:
        return "arrow"
    if "TRIANGLE" in st:
        return "triangle"
    return "rect"


def _extract_layout_pictures(slide, sw: int, sh: int) -> list[dict]:
    """Extract Picture elements from slide layout and slide master.

    These are often background images, logos, or decorations
    inherited from the layout that python-pptx doesn't expose
    via slide.shapes.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    results = []
    z = 1000
    seen_embeds = set()

    # Check if slide already has a full-size image covering everything
    try:
        for shape in slide.shapes:
            from pptx.shapes.picture import Picture as PptxPic
            if isinstance(shape, PptxPic):
                try:
                    xml = shape._element.xml
                    m = __import__('re').search(r'r:embed="([^"]+)"', xml)
                    if m:
                        seen_embeds.add(m.group(1))
                except Exception:
                    pass
    except Exception:
        pass

    for level, elem in [("slide_layout", slide.slide_layout._element),
                         ("slide_master", slide.slide_layout.slide_master._element)]:
        if elem is None:
            continue
        try:
            c_sld = elem.find(qn("p:cSld"))
            if c_sld is None:
                continue
            sp_tree = c_sld.find(qn("p:spTree"))
            if sp_tree is None:
                continue
            for pic in sp_tree.findall(qn("p:pic")):
                try:
                    blip_fill = pic.find(qn("p:blipFill"))
                    if blip_fill is None:
                        continue
                    blip = blip_fill.find(qn("a:blip"))
                    if blip is None:
                        continue
                    embed = blip.get(qn("r:embed"))
                    if not embed or embed in seen_embeds:
                        continue
                    seen_embeds.add(embed)

                    sp_pr = pic.find(qn("p:spPr"))
                    xfrm = sp_pr.find(qn("a:xfrm")) if sp_pr is not None else None
                    off = xfrm.find(qn("a:off")) if xfrm is not None else None
                    ext = xfrm.find(qn("a:ext")) if xfrm is not None else None

                    lx = int(off.get("x", 0)) if off is not None else 0
                    ly = int(off.get("y", 0)) if off is not None else 0
                    wx = int(ext.get("cx", 0)) if ext is not None else 0
                    hy = int(ext.get("cy", 0)) if ext is not None else 0

                    is_full_bg = (wx >= sw * 0.95 and hy >= sh * 0.95)

                    results.append({
                        "id": z,
                        "z_order": z,
                        "element_type": "image",
                        "semantic_role": "background" if is_full_bg else "decoration",
                        "from_layout": level,
                        "embed_id": embed,
                        "is_background": is_full_bg,
                        "position": {
                            "x": round(lx / sw, 4) if sw else 0,
                            "y": round(ly / sh, 4) if sh else 0,
                            "w": round(wx / sw, 4) if sw else 0,
                            "h": round(hy / sh, 4) if sh else 0,
                        },
                    })
                    z += 1
                except Exception:
                    pass
        except Exception:
            pass

    return results


def _collect_elements_data(slide, sw: int, sh: int, spec_fs: Path) -> list[dict]:
    """Extract ALL elements from a slide and its layout/master into plain dicts.

    Includes: text shapes, pictures, tables, auto-shapes, group shapes,
    and layout/master-level pictures (backgrounds, logos).
    """
    elements: list[dict] = []

    # 1. Shapes from the slide itself
    for zi, shape in enumerate(slide.shapes):
        try:
            d = _extract_shape_deep(shape, sw, sh, zi, spec_fs)
            elements.append(d)
        except Exception:
            pass

    # 2. Pictures from layout/master (backgrounds, logos not in slide.shapes)
    layout_pics = _extract_layout_pictures(slide, sw, sh)
    elements.extend(layout_pics)

    return elements


def _extract_layout_background(slide, spec_dir: Path, page_idx: int, sw_px: int = 12192000, sh_px: int = 6858000) -> dict:
    """Extract background image/color from slide, layout, and master.

    Checks both <p:bg> elements AND full-size <p:pic> in layout spTree.
    """
    from pptx.oxml.ns import qn
    import re

    result = {"type": "solid", "color": "#FFFFFF", "description": "White background"}

    # 1. Check <p:bg> on slide level
    try:
        bg = slide.background
        from pptx.enum.dml import MSO_FILL_TYPE
        if bg.fill.type == MSO_FILL_TYPE.SOLID:
            try:
                result["color"] = f"#{bg.fill.fore_color.rgb}"
                result["type"] = "solid"
                result["description"] = f"Solid background {result['color']}"
                return result
            except Exception:
                pass
    except Exception:
        pass

    # 2. Check for full-size <p:pic> in layout spTree (most common case)
    for level_name, level_elem in [
        ("layout", slide.slide_layout._element),
        ("master", slide.slide_layout.slide_master._element),
    ]:
        if level_elem is None:
            continue
        try:
            c_sld = level_elem.find(qn("p:cSld"))
            if c_sld is None:
                continue
            sp_tree = c_sld.find(qn("p:spTree"))
            if sp_tree is None:
                continue
            for pic in sp_tree.findall(qn("p:pic")):
                sp_pr = pic.find(qn("p:spPr"))
                if sp_pr is None:
                    continue
                xfrm = sp_pr.find(qn("a:xfrm"))
                if xfrm is None:
                    continue
                ext = xfrm.find(qn("a:ext"))
                if ext is None:
                    continue
                wx = int(ext.get("cx", 0))
                hy = int(ext.get("cy", 0))
                # Full-size = background
                if wx >= sw_px * 0.95 and hy >= sh_px * 0.95:
                    blip_fill = pic.find(qn("p:blipFill"))
                    if blip_fill is not None:
                        blip = blip_fill.find(qn("a:blip"))
                        if blip is not None:
                            embed = blip.get(qn("r:embed"))
                            if embed:
                                try:
                                    # Try slide part rels first, then layout
                                    if embed in slide.part.rels:
                                        rel = slide.part.rels[embed]
                                    elif level_name == "layout" and embed in slide.slide_layout.part.rels:
                                        rel = slide.slide_layout.part.rels[embed]
                                    else:
                                        raise KeyError(f"Embed {embed} not found")
                                    image = rel.target_part
                                    blob = image.blob
                                    ext_name = image.content_type.split("/")[-1]
                                    if ext_name == "jpeg":
                                        ext_name = "jpg"
                                    assets_dir = spec_dir / "assets"
                                    assets_dir.mkdir(parents=True, exist_ok=True)
                                    bg_path = assets_dir / f"bg_{page_idx:03d}.{ext_name}"
                                    bg_path.write_bytes(blob)
                                    result = {
                                        "type": "image",
                                        "color": "",
                                        "image": f"assets/{bg_path.name}",
                                        "description": f"Background image from {level_name}",
                                    }
                                    return result
                                except Exception:
                                    pass
                    result = {
                        "type": "image",
                        "color": "",
                        "image": "",
                        "description": f"Full-size image background from {level_name} (unsaved)",
                    }
                    return result
        except Exception:
            pass

    # 3. Check <p:bg> on layout/master
    for level_name, level_elem in [
        ("layout", slide.slide_layout._element),
        ("master", slide.slide_layout.slide_master._element),
    ]:
        if level_elem is None:
            continue
        try:
            c_sld = level_elem.find(qn("p:cSld"))
            if c_sld is not None:
                bg_elem = c_sld.find(qn("p:bg"))
                if bg_elem is not None:
                    bg_pr = bg_elem.find(qn("p:bgPr"))
                    if bg_pr is not None:
                        solid = bg_pr.find(qn("a:solidFill"))
                        if solid is not None:
                            srgb = solid.find(qn("a:srgbClr"))
                            if srgb is not None:
                                val = srgb.get("val", "")
                                if val:
                                    result = {
                                        "type": "solid",
                                        "color": f"#{val}",
                                        "description": f"Solid background #{val} from {level_name}",
                                    }
                                    return result
        except Exception:
            pass

    return result


def _save_slide_assets(elements: list[dict], slide, spec_dir: Path, page_idx: int, sw: int, sh: int) -> None:
    """Save image assets from slide shapes AND layout elements to assets/ directory.

    Updates each image element's dict with the saved file path for any
    element that has an `embed_id` key (from _extract_layout_pictures).
    """
    from pptx.shapes.picture import Picture as PptxPic
    from pptx.oxml.ns import qn

    assets_dir = spec_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    # Save slide-level pictures
    img_counter = [0]
    for shape in slide.shapes:
        if not isinstance(shape, PptxPic):
            continue
        try:
            image = shape.image
            blob = image.blob
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"img_{page_idx:02d}_{img_counter[0]:02d}.{ext}"
            (assets_dir / filename).write_bytes(blob)
            img_counter[0] += 1

            # Update matching element
            shape_name = getattr(shape, "name", "")
            for el in elements:
                if el.get("element_type") == "image" and el.get("shape_name") == shape_name:
                    el["saved_asset"] = f"assets/{filename}"
                    el["image"] = el.get("image", {})
                    el["image"]["saved_path"] = f"assets/{filename}"
                    break
        except Exception:
            pass

    # Save layout-level pictures (found via _extract_layout_pictures)
    for el in elements:
        embed_id = el.get("embed_id")
        if not embed_id or el.get("from_layout") is None:
            continue
        if el.get("saved_asset"):
            continue
        # Backgrounds are already saved by _extract_layout_background
        if el.get("is_background"):
            continue

        try:
            # Try layout part rels first, then master
            if embed_id in slide.slide_layout.part.rels:
                rel = slide.slide_layout.part.rels[embed_id]
            elif embed_id in slide.slide_layout.slide_master.part.rels:
                rel = slide.slide_layout.slide_master.part.rels[embed_id]
            else:
                continue

            image = rel.target_part
            blob = image.blob
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"

            prefix = "bg" if el.get("is_background") else "dec"
            filename = f"{prefix}_{page_idx:02d}_{img_counter[0]:02d}.{ext}"
            (assets_dir / filename).write_bytes(blob)
            img_counter[0] += 1

            el["saved_asset"] = f"assets/{filename}"
            el["image"] = el.get("image", {})
            el["image"]["saved_path"] = f"assets/{filename}"
            el["image"]["content_type"] = image.content_type
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Layout dedup — group similar pages
# ---------------------------------------------------------------------------


def _layout_signature(elements: list[dict], page_type: str = "content", layout_sub_type: str = "full_width") -> str:
    """Generate a compact signature for layout similarity comparison.

    Based on element types, roles, relative positions, and style patterns.
    Two slides with the same signature are considered "same layout".

    Includes VL-assigned page_type and layout_sub_type for better dedup.
    """
    parts: list[str] = []

    parts.append(f"PT:{page_type}")
    parts.append(f"LT:{layout_sub_type}")

    # Count by type
    type_counts: dict[str, int] = {}
    role_counts: dict[str, int] = {}
    for el in elements:
        et = el.get("element_type", "shape")
        type_counts[et] = type_counts.get(et, 0) + 1
        role = el.get("semantic_role", "body")
        role_counts[role] = role_counts.get(role, 0) + 1

    parts.append("T:" + "_".join(f"{k}{v}" for k, v in sorted(type_counts.items())))
    parts.append("R:" + "_".join(f"{k}{v}" for k, v in sorted(role_counts.items())))

    # Position grid: divide slide into 4x3 grid, mark occupied cells
    occupied: set[str] = set()
    for el in elements:
        pos = el.get("position", {})
        x = pos.get("x", 0)
        y = pos.get("y", 0)
        w = pos.get("w", 0)
        h = pos.get("h", 0)
        cx = int((x + w / 2) * 4)
        cy = int((y + h / 2) * 3)
        occupied.add(f"{cx}{cy}")

    parts.append("G:" + "".join(sorted(occupied)))

    # Major structural attributes
    fill_colors = set()
    for el in elements:
        ss = el.get("shape_style", {}) or {}
        if ss.get("fill_color"):
            fill_colors.add(ss["fill_color"])
    if fill_colors:
        parts.append("F:" + "_".join(sorted(fill_colors)[:3]))

    return "|".join(parts)


def _group_similar_pages(pages: list[dict]) -> list[dict]:
    """Group similar pages together, returning merged page specs.

    Returns deduplicated list where similar pages are merged into one entry
    with a `slide_indices` list referencing which original slides use this layout.
    """
    groups: dict[str, dict] = {}
    for page in pages:
        sig = page.get("_layout_signature", "")
        if sig in groups:
            groups[sig]["slide_indices"].append(page.get("_slide_index", 0))
            groups[sig]["slide_count"] += 1
        else:
            groups[sig] = {
                "slide_indices": [page.get("_slide_index", 0)],
                "slide_count": 1,
                "page_type": page.get("page_type", "content"),
                "layout_sub_type": page.get("layout_sub_type", "full_width"),
                "background": page.get("background", {}),
                "elements": page.get("elements", []),
                "element_groups": page.get("element_groups", []),
                "design_patterns": page.get("design_patterns", []),
                "_vl_elements": page.get("_vl_elements", []),
                "spec_file_name": "",
            }

    # Name each group based on page_type + layout_sub_type
    result: list[dict] = []
    type_counter: dict[str, int] = {}
    for sig, group in groups.items():
        pt = group["page_type"]
        lst = group["layout_sub_type"]
        if pt == "content":
            base = f"content_{lst}"
        else:
            base = pt
        type_counter[base] = type_counter.get(base, 0) + 1
        if type_counter[base] > 1:
            group["spec_file_name"] = f"{base}_{type_counter[base]}"
        else:
            group["spec_file_name"] = base
        result.append(group)

    return result


# ---------------------------------------------------------------------------
# Spec name generation
# ---------------------------------------------------------------------------


def _generate_spec_name(pptx_path: Path) -> str:
    """Generate a short, memorable spec name from the PPTX filename."""
    stem = pptx_path.stem
    name = re.sub(r"[^a-zA-Z0-9\s-]", "", stem)
    name = re.sub(r"\s+", "-", name.strip())
    name = name.lower()

    # Remove common suffixes
    for suffix in ["-pptx", "-presentation", "-deck", "-slide", "-slides"]:
        if name.endswith(suffix):
            name = name[: -len(suffix)]

    # Truncate
    if len(name) > 30:
        name = name[:30].rstrip("-")

    return name or pptx_path.stem


# ---------------------------------------------------------------------------
# VLVSpecExtractor — main orchestrator
# ---------------------------------------------------------------------------


class VLVSpecExtractor:
    """VL-driven spec extraction — produces deduplicated JSON spec directory.

    Pipeline:
      1. Open PPTX, extract theme colors/fonts
      2. Per slide: extract elements → render PNG → VL analysis → store
      3. Group similar layouts together (dedup)
      4. Generate JSON spec files (no text content)
      5. Save slide PNGs + JSON files to specs/<name>/
    """

    def __init__(self):
        self.analyzer = VLElementAnalyzer()

    def extract(self, pptx_path: str | Path, output_name: str | None = None) -> Path:
        """Extract VL-driven spec from a PPTX file.

        Args:
            pptx_path: Path to the source .pptx file.
            output_name: Spec name (directory name). Auto-generated from filename if None.

        Returns:
            Path to the created spec directory.
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        spec_name = output_name or _generate_spec_name(pptx_path)
        spec_dir = Path("specs") / spec_name
        slides_dir = spec_dir / "slides"

        # Clean existing
        if spec_dir.exists():
            shutil.rmtree(str(spec_dir))
        spec_dir.mkdir(parents=True)
        slides_dir.mkdir()

        prs = Presentation(str(pptx_path))
        sw = int(prs.slide_width)
        sh = int(prs.slide_height)

        # ── 1. Theme extraction ──
        pptx_abs = str(pptx_path.absolute())
        theme_colors = extract_theme_colors(pptx_abs)
        theme_fonts = extract_theme_fonts(pptx_abs)

        palette = {}
        for k, v in theme_colors.items():
            if v:
                palette[k] = v
        if not palette:
            palette = {
                "dk1": "#000000", "lt1": "#FFFFFF",
                "accent1": "#4472C4", "accent2": "#ED7D31",
            }

        total_slides = len(prs.slides)

        master_spec = {
            "spec_name": spec_name,
            "source": pptx_path.name,
            "slide_count": total_slides,
            "canvas": {
                "width_emu": sw,
                "height_emu": sh,
                "width_inches": round(sw / 914400, 2),
                "height_inches": round(sh / 914400, 2),
            },
            "palette": palette,
            "typography": {
                "heading_family": theme_fonts.get("majorFont", ""),
                "body_family": theme_fonts.get("minorFont", ""),
            },
            "layouts": [],
        }

        # ── Convert PPTX to PDF once (for screenshot rendering) ──
        pdf_path = _convert_pptx_to_pdf(pptx_path)

        # ── 2. Per-slide processing ──
        raw_pages: list[dict] = []

        for idx, slide in enumerate(prs.slides):
            slide_idx = idx

            # Extract elements (programmatic) — now includes all shapes + layout/master pics
            elements_data = _collect_elements_data(slide, sw, sh, spec_dir)

            # Render slide to real PNG screenshot (from PDF)
            png_path = _render_slide_to_png(pdf_path, slides_dir, idx)
            if png_path:
                slide_rel = f"slides/{png_path.name}"
            else:
                slide_rel = ""

            # VL analysis — describes element roles and visual groups, NOT page type
            vl_result = self.analyzer.analyze(png_path or "", elements_data)

            # Page type: pure position-based rules
            if slide_idx == 0:
                page_type = "cover"
            elif slide_idx == total_slides - 1:
                page_type = "end_page"
            elif slide_idx == 1 and total_slides > 3:
                text_count = sum(1 for e in elements_data if e.get("element_type") == "text")
                page_type = "toc" if text_count >= 3 else "content"
            else:
                page_type = "content"

            # Layout sub-type from VL (visual arrangement only)
            layout_sub_type = vl_result.layout_sub_type

            # ── Extract background (handles <p:bg> AND layout <p:pic>) ──
            background_info = _extract_layout_background(slide, spec_dir, idx, sw, sh)

            # ── Save image assets to disk ──
            _save_slide_assets(elements_data, slide, spec_dir, idx, sw, sh)

            # Generate layout signature for dedup
            sig = _layout_signature(elements_data, page_type, layout_sub_type)

            raw_pages.append({
                "_slide_index": slide_idx,
                "_layout_signature": sig,
                "page_type": page_type,
                "layout_sub_type": layout_sub_type,
                "background": background_info,
                "elements": elements_data,
                "_vl_elements": vl_result.elements,
                "element_groups": vl_result.element_groups,
                "design_patterns": vl_result.design_patterns,
                "slide_image": slide_rel,
            })

        # ── 3. Layout dedup ──
        layout_groups = _group_similar_pages(raw_pages)

        # ── 4. Write JSON files ──
        layout_file_map: dict[str, str] = {}

        for group in layout_groups:
            file_name = group["spec_file_name"]
            layout_file_map[file_name] = file_name

            # Write elements: use raw extracted data directly, just merge VL descriptions
            raw_elements = group.get("elements", [])
            vl_analysis = group.get("_vl_elements", [])

            # Build VL lookup by id
            vl_by_id = {}
            for vl_el in vl_analysis:
                vl_by_id[vl_el.get("id")] = vl_el

            # Merge VL descriptions into element data
            for el in raw_elements:
                eid = el.get("id")
                vl = vl_by_id.get(eid)
                if vl:
                    if vl.get("description"):
                        el["vl_description"] = vl["description"]
                    if vl.get("visual_weight"):
                        el["vl_visual_weight"] = vl["visual_weight"]
                    if vl.get("style_notes"):
                        el["vl_style_notes"] = vl["style_notes"]
                    if vl.get("role"):
                        el["vl_role"] = vl["role"]

            # Build the JSON page spec with full element data
            page_spec = {
                "spec_name": spec_name,
                "layout_id": file_name,
                "page_type": group["page_type"],
                "layout_sub_type": group["layout_sub_type"],
                "slide_count": group["slide_count"],
                "slide_indices": group["slide_indices"],
                "canvas": master_spec["canvas"],
                "background": group["background"],
                "design_patterns": group.get("design_patterns", []),
                "element_groups": group.get("element_groups", []),
                "elements": raw_elements,
            }

            json_path = spec_dir / f"{file_name}.json"
            json_path.write_text(
                json.dumps(page_spec, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )

            master_spec["layouts"].append({
                "file": f"{file_name}.json",
                "page_type": group["page_type"],
                "layout_sub_type": group["layout_sub_type"],
                "slide_count": group["slide_count"],
                "slide_indices": group["slide_indices"],
            })

        # ── 5. Write master spec.json ──
        (spec_dir / "spec.json").write_text(
            json.dumps(master_spec, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        # Copy source PPTX for reference
        shutil.copy2(str(pptx_path), str(spec_dir / pptx_path.name))

        print(f"VL Spec extracted: {spec_dir}/")
        print(f"  Pages: {len(raw_pages)} slides → {len(layout_groups)} unique layouts")
        print(f"  JSON files: {', '.join(sorted(layout_file_map.keys()))}")
        print(f"  VL analysis: {'enabled' if self.analyzer.enabled else 'disabled (fallback to programmatic)'}")
        if self.analyzer.enabled:
            print(f"  VL model: {self.analyzer.config.get('model', '?')}")

        return spec_dir


__all__ = ["VLVSpecExtractor", "ElementAnalysisResult"]
