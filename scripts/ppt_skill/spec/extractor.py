"""SpecExtractor orchestrator — ties all extraction modules into a single pipeline.

Orchestrates the full end-to-end extraction: theme colors/fonts/background,
slide classification, layout analysis, density rhythm, and font size analysis.
Produces a complete DesignSpec object and serializes it to YAML.

Usage:
    extractor = SpecExtractor("input.pptx", "corporate-blue")
    spec = extractor.extract()          # → DesignSpec dataclass
    path = extractor.save("specs")      # → specs/corporate-blue.yaml
"""

from __future__ import annotations

import logging
import warnings
from dataclasses import asdict, is_dataclass
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation

from ppt_skill.config import DESIGN_COLORS
from ppt_skill.spec.spec_model import (
    ColorPalette,
    DesignSpec,
    DensityLabel,
    PresentationRhythm,
    SlideSpec,
    SlideType,
    Typography,
)

# Extraction module imports
from ppt_skill.spec import theme as theme_mod
from ppt_skill.spec import slide_classifier as sc_mod
from ppt_skill.spec import layout_analysis as la_mod
from ppt_skill.spec import density as dens_mod
from ppt_skill.spec import font_analysis as fa_mod

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# OOXML scheme name → ColorPalette field mapping
# ---------------------------------------------------------------------------

_SCHEME_TO_PALETTE: dict[str, str] = {
    "dk1":      "background1",
    "dk2":      "background2",
    "lt1":      "text1",
    "lt2":      "text2",
    "accent1":  "accent1",
    "accent2":  "accent2",
    "accent3":  "accent3",
    "accent4":  "accent4",
    "accent5":  "accent5",
    "accent6":  "accent6",
    "hlink":    "hyperlink",
    "folHlink": "followed_hyperlink",
}


# ---------------------------------------------------------------------------
# Serialization helpers
# ---------------------------------------------------------------------------

def _enum_value(obj: Any) -> Any:
    """Extract .value from Enum types, or return the object unchanged."""
    if isinstance(obj, Enum):
        return obj.value
    return obj


def _dataclass_to_dict(obj: Any) -> Any:
    """Recursively convert a dataclass (or nested structure) to plain dicts.

    Uses dataclasses.asdict() for the initial conversion, then walks through
    the result to extract .value from Enum members.
    """
    if is_dataclass(obj) and not isinstance(obj, type):
        raw = asdict(obj)
        return _walk_enum_values(raw)
    elif isinstance(obj, dict):
        return {k: _dataclass_to_dict(v) for k, v in obj.items()}
    elif isinstance(obj, (list, tuple)):
        return [_dataclass_to_dict(item) for item in obj]
    else:
        return _enum_value(obj)


def _walk_enum_values(d: dict) -> dict:
    """Walk a nested dict structure and extract Enum .value at any level."""
    result: dict = {}
    for k, v in d.items():
        if isinstance(v, Enum):
            result[k] = v.value
        elif isinstance(v, dict):
            result[k] = _walk_enum_values(v)
        elif isinstance(v, list):
            result[k] = [
                _walk_enum_values(item) if isinstance(item, dict)
                else item.value if isinstance(item, Enum)
                else item
                for item in v
            ]
        else:
            result[k] = v
    return result


# ---------------------------------------------------------------------------
# SpecExtractor — main orchestrator
# ---------------------------------------------------------------------------


class SpecExtractor:
    """Orchestrates PPTX spec extraction end-to-end.

    Ties together all extraction modules (colors, fonts, layouts,
    classifications, density, font sizes) into a single pipeline that
    produces a fully populated DesignSpec dataclass.

    Usage::

        extractor = SpecExtractor("deck.pptx", "my-style")
        spec = extractor.extract()       # → DesignSpec
        path = extractor.save("specs")   # → specs/my-style.yaml
    """

    def __init__(self, pptx_path: str | Path, spec_name: str):
        self.pptx_path = Path(pptx_path)
        self.spec_name = spec_name
        self.prs = Presentation(str(self.pptx_path))
        self.spec = DesignSpec()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def extract(self) -> DesignSpec:
        """Run full extraction pipeline, return populated DesignSpec."""

        # --- 1. Metadata ---
        self.spec.metadata = {
            "name": self.spec_name,
            "source": self.pptx_path.name,
            "extracted_at": datetime.now().isoformat(),
            "slide_count": len(self.prs.slides),
        }

        # --- 2. Colors ---
        self._extract_colors()

        # --- 3. Theme fonts (families) ---
        self._extract_theme_fonts()

        # --- 4. Per-slide loop ---
        density_list: list[dict] = []
        slide_types: list[str] = []

        for i, slide in enumerate(self.prs.slides):
            slide_idx = i + 1

            # Classify
            try:
                slide_type_str = sc_mod.classify_slide(slide)
            except Exception:
                logger.warning("Slide %d: classification failed, defaulting to content", slide_idx)
                slide_type_str = SlideType.CONTENT.value
            slide_types.append(slide_type_str)

            # Layout analysis
            try:
                layout_data = la_mod.analyze_slide_layout(slide)
            except Exception:
                logger.warning("Slide %d: layout analysis failed", slide_idx)
                layout_data = {}

            # Density
            try:
                density_data = dens_mod.analyze_slide_density(slide, slide_idx)
            except Exception:
                logger.warning("Slide %d: density analysis failed", slide_idx)
                density_data = {
                    "slide_index": slide_idx,
                    "char_count": 0,
                    "image_count": 0,
                    "shape_count": 0,
                }
            density_list.append(density_data)

            # Background
            try:
                bg_data = theme_mod.extract_slide_background(slide, str(self.pptx_path))
            except Exception:
                logger.warning("Slide %d: background extraction failed", slide_idx)
                bg_data = None

            # Build SlideSpec
            slide_spec = SlideSpec(
                slide_index=slide_idx,
                slide_type=SlideType(slide_type_str),
                layout_name=slide.slide_layout.name,
                density=DensityLabel.DENSE,  # placeholder; overwritten in step 5
                char_count=density_data.get("char_count", 0),
                image_count=density_data.get("image_count", 0),
                shape_count=density_data.get("shape_count", 0),
                background=bg_data,
            )
            self.spec.slides.append(slide_spec)

        # --- 5. Density classification ---
        if density_list:
            classified = dens_mod.classify_density(density_list)
            for i, ddata in enumerate(classified):
                if i < len(self.spec.slides):
                    self.spec.slides[i].density = DensityLabel(ddata["density"])

        # --- 6. Font sizes ---
        try:
            all_fonts = fa_mod.extract_all_slide_fonts(self.prs)
            typo_sizes = fa_mod.compute_spec_typography_sizes(all_fonts)
            self.spec.typography.heading_sizes = typo_sizes.get("heading_sizes", {})
            self.spec.typography.body_sizes = typo_sizes.get("body_sizes", {})
        except Exception:
            logger.warning("Font size extraction failed — leaving empty heading_sizes/body_sizes")

        # --- 7. Rhythm ---
        try:
            self.spec.rhythm = dens_mod.build_presentation_rhythm(density_list, slide_types)
        except Exception:
            logger.warning("Rhythm construction failed")

        # --- 8. Source config ---
        self._derive_source_config()

        return self.spec

    def save(self, specs_dir: str | Path = "specs") -> Path:
        """Extract and save YAML spec file. Returns file path.

        Calls extract() if spec has no slides (i.e., hasn't been extracted yet).
        Creates the specs_dir directory if it doesn't exist.
        """
        out_dir = Path(specs_dir)

        # Run extraction if spec is empty (no slides yet)
        if not self.spec.slides and not self.spec.metadata:
            self.extract()

        out_dir.mkdir(parents=True, exist_ok=True)

        # Use custom to_dict() if available (handles backward-compat properties)
        if hasattr(self.spec, "to_dict") and callable(self.spec.to_dict):
            spec_data = self.spec.to_dict()
        else:
            spec_data = _dataclass_to_dict(self.spec)
        out_path = out_dir / f"{self.spec_name}.yaml"

        with open(out_path, "w", encoding="utf-8") as f:
            yaml.dump(
                spec_data, f,
                default_flow_style=False,
                sort_keys=False,
                allow_unicode=True,
            )

        return out_path

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _extract_colors(self) -> None:
        """Extract theme colors and populate ColorPalette, falling back to config defaults."""
        try:
            scheme_colors = theme_mod.extract_theme_colors(str(self.pptx_path))
        except Exception as exc:
            logger.warning("Theme color extraction failed: %s — using DESIGN_COLORS defaults", exc)
            scheme_colors = {}

        # Map OOXML scheme names → ColorPalette field names
        palette_kwargs: dict[str, str] = {}
        for scheme_name, palette_field in _SCHEME_TO_PALETTE.items():
            hex_val = scheme_colors.get(scheme_name, "")
            if not hex_val:
                # Fall back to config.py DESIGN_COLORS
                hex_val = DESIGN_COLORS.get(scheme_name, "")
            palette_kwargs[palette_field] = hex_val

        self.spec.colors = ColorPalette(**palette_kwargs)

    def _extract_theme_fonts(self) -> None:
        """Extract theme font families and populate Typography fields."""
        try:
            font_data = theme_mod.extract_theme_fonts(str(self.pptx_path))
        except Exception as exc:
            logger.warning("Theme font extraction failed: %s", exc)
            font_data = {}

        self.spec.typography.heading_family = font_data.get("majorFont", "")
        self.spec.typography.body_family = font_data.get("minorFont", "")

    def _derive_source_config(self) -> None:
        """Derive canvas format hints from slide dimensions.

        16:9 ratio ≈ 1.778, 4:3 ratio ≈ 1.333.
        """
        try:
            width = self.prs.slide_width   # EMU
            height = self.prs.slide_height  # EMU
        except Exception:
            self.spec.source_config = {}
            return

        if width and height:
            ratio = width / height
            if 1.70 < ratio < 1.85:
                canvas = "ppt169"
            elif 1.28 < ratio < 1.38:
                canvas = "ppt43"
            else:
                # Approximate — pick closest
                canvas = "ppt169" if ratio > 1.5 else "ppt43"
            self.spec.source_config = {"canvas": canvas}
        else:
            self.spec.source_config = {}


# ---------------------------------------------------------------------------
# Module-level helpers for YAML serialization
# ---------------------------------------------------------------------------

def _spec_to_dict(spec: DesignSpec) -> dict:
    """Convert a DesignSpec to plain dicts/lists for YAML serialization.

    Uses dataclasses.asdict() as the base, then extracts .value from Enum
    members to produce clean string output suitable for yaml.dump().
    """
    return _dataclass_to_dict(spec)


__all__ = [
    "SpecExtractor",
    "_dataclass_to_dict",
    "_spec_to_dict",
]
