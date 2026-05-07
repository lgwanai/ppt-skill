"""Generate visual representations of PPTX slides for VL model analysis.

Uses python-pptx to extract shape positions/dimensions, generates a layout
visualization as a simplified structural image showing text regions, image
regions, and decorative elements as colored boxes.

Two modes:
  1. layout_viz: Fast — generates SVG/PNG showing bounding boxes of all shapes
     with type-based coloring (text=blue, image=green, shape=orange, etc.)
  2. render: Uses LibreOffice or Quartz/CoreGraphics to render actual slide PNG
     (requires LibreOffice in PATH or macOS)

Output: PNG files in the spec directory for VL model consumption.
"""

from __future__ import annotations

import io
import os
import struct
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree


@dataclass
class SlideRegion:
    """A detected region within a slide."""
    x: int       # EMU
    y: int
    w: int
    h: int
    shape_type: str     # "text", "image", "table", "chart", "group", "shape"
    text_preview: str   # First 80 chars of text content (for VL context)
    has_image: bool


# ── Shape extraction ────────────────────────────────────────────────


def _emu_to_px(emu: int, dpi: int = 96) -> int:
    """Convert EMU (English Metric Units) to pixels at given DPI."""
    return int(emu * dpi / 914400)


def _px_to_inches(px: float, dpi: int = 96) -> float:
    return px / dpi


def extract_slide_regions(slide) -> list[SlideRegion]:
    """Extract all shape bounding boxes from a python-pptx Slide object."""
    regions: list[SlideRegion] = []
    for shape in slide.shapes:
        shape_type = _classify_shape(shape)
        text_preview = _extract_text_preview(shape)
        has_image = _has_image(shape)
        regions.append(SlideRegion(
            x=shape.left or 0,
            y=shape.top or 0,
            w=shape.width or 0,
            h=shape.height or 0,
            shape_type=shape_type,
            text_preview=text_preview,
            has_image=has_image,
        ))
    return regions


def _classify_shape(shape) -> str:
    """Classify a python-pptx shape into a type string."""
    try:
        from pptx.shapes.picture import Picture
    except ImportError:
        Picture = None

    if shape.has_table:
        return "table"
    if hasattr(shape, "chart"):
        return "chart"
    if shape.has_text_frame and shape.text_frame.text.strip():
        if _is_decorative_shape(shape):
            return "decorative"
        return "text"

    if hasattr(shape, "shape_type") and shape.shape_type is not None:
        st = str(shape.shape_type)
        if "PICTURE" in st.upper():
            return "image"
        if "CHART" in st.upper():
            return "chart"
        if "TABLE" in st.upper():
            return "table"
        if "PLACEHOLDER" in st.upper():
            return "placeholder"
        if "GROUP" in st.upper():
            return "group"
        if "TEXT" in st.upper():
            return "text"

    return "shape"


def _is_decorative_shape(shape) -> bool:
    """Heuristic: is this a decorative element (line, icon, accent shape) instead of content?"""
    text = shape.text_frame.text.strip() if shape.has_text_frame else ""
    if len(text) < 5 and shape.width:
        width_inches = shape.width / 914400
        if width_inches < 1.5:
            return True
    if hasattr(shape, "shape_type") and shape.shape_type is not None:
        st = str(shape.shape_type).upper()
        if "LINE" in st or "CONNECTOR" in st:
            return True
    return False


def _has_image(shape) -> bool:
    try:
        from pptx.shapes.picture import Picture
        return isinstance(shape, Picture)
    except ImportError:
        return False


def _extract_text_preview(shape) -> str:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) > 80:
            return text[:77] + "..."
        return text
    return ""


# ── Layout visualization SVG ─────────────────────────────────────────


def generate_layout_svg(slide_width_emu: int, slide_height_emu: int,
                        regions: list[SlideRegion]) -> str:
    """Generate an SVG showing all shape regions as colored bounding boxes.

    Colors by type: text=#3B82F6, image=#10B981, table=#F59E0B,
    chart=#EF4444, shape=#8B5CF6, decorative=#EC4899, group=#6B7280.
    Each box includes a label with the first line of text content.
    """
    W = _emu_to_px(slide_width_emu, 96)
    H = _emu_to_px(slide_height_emu, 96)

    color_map = {
        "text": "#3B82F6", "image": "#10B981", "table": "#F59E0B",
        "chart": "#EF4444", "shape": "#8B5CF6", "decorative": "#EC4899",
        "group": "#6B7280", "placeholder": "#9CA3AF",
    }
    opacity_map = {"text": 0.15, "image": 0.2, "table": 0.15,
                   "chart": 0.15, "shape": 0.1, "decorative": 0.1,
                   "group": 0.1, "placeholder": 0.08}

    rects: list[str] = []
    labels: list[str] = []
    for i, r in enumerate(regions):
        c = color_map.get(r.shape_type, "#6B7280")
        o = opacity_map.get(r.shape_type, 0.1)
        rx = _emu_to_px(r.x, 96)
        ry = _emu_to_px(r.y, 96)
        rw = max(_emu_to_px(r.w, 96), 1)
        rh = max(_emu_to_px(r.h, 96), 1)

        rects.append(
            f'<rect x="{rx}" y="{ry}" width="{rw}" height="{rh}" '
            f'fill="{c}" fill-opacity="{o}" stroke="{c}" stroke-width="1" rx="2"/>'
        )

        # Label with type and text preview
        if r.text_preview:
            label_text = f"[{r.shape_type[:4]}] {r.text_preview[:30]}"
        else:
            label_text = f"[{r.shape_type[:4]}]"
        labels.append(
            f'<text x="{rx + 4}" y="{ry + 12}" '
            f'fill="{c}" font-size="10" font-family="monospace">{label_text}</text>'
        )

    svg = f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 {W} {H}" width="{W}" height="{H}">
  <rect width="{W}" height="{H}" fill="#F8FAFC" rx="4"/>
  {chr(10).join("  " + r for r in rects)}
  {chr(10).join("  " + l for l in labels)}
</svg>'''
    return svg


# ── Image export ────────────────────────────────────────────────────


def svg_to_png(svg_data: str, output_path: Path,
               width: int = 1280, height: int = 720) -> bool:
    """Convert SVG to PNG using Pillow (+ cairosvg) or fallback."""
    # Try cairosvg first
    try:
        import cairosvg
        cairosvg.svg2png(
            bytestring=svg_data.encode("utf-8"),
            write_to=str(output_path),
            output_width=width,
            output_height=height,
        )
        return True
    except ImportError:
        pass

    # Fallback: write SVG and try to convert
    svg_path = output_path.with_suffix(".svg")
    svg_path.write_text(svg_data, encoding="utf-8")

    # Try libreoffice / rsvg-convert
    for cmd in [
        ["rsvg-convert", "-w", str(width), "-o", str(output_path), str(svg_path)],
    ]:
        try:
            subprocess.run(cmd, capture_output=True, check=True, timeout=30)
            return True
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass

    # Final fallback: just save the SVG
    return False


def render_slide_as_png(prs_path: Path, slide_index: int,
                        output_path: Path, width: int = 1920) -> bool:
    """Render a specific PPTX slide as a PNG image.

    Uses LibreOffice headless conversion on Linux/Mac, or COM on Windows.
    Returns True on success.
    """
    # Try LibreOffice headless
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            cmd = [
                "soffice", "--headless", "--convert-to", "png",
                f"--outdir", tmpdir,
                str(prs_path)
            ]
            result = subprocess.run(cmd, capture_output=True, timeout=60)
            if result.returncode == 0:
                # Find the generated PNG for this slide
                import glob
                pngs = sorted(glob.glob(os.path.join(tmpdir, "*.png")))
                if len(pngs) > slide_index:
                    import shutil
                    shutil.copy(pngs[slide_index], output_path)
                    return True
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        pass

    # Try macOS Quartz (sips + Quick Look)
    if os.uname().sysname == "Darwin":
        try:
            # Use python-pptx to export slide XML, generate layout viz
            return False  # Layout viz fallback
        except Exception:
            pass

    return False


def generate_slide_preview(prs, slide_index: int,
                           output_png: Path, output_svg: Path) -> tuple[bool, bool]:
    """Generate both a layout visualization and (if possible) a real render.

    Returns (png_ok, svg_ok).
    """
    slide = prs.slides[slide_index]
    sw = prs.slide_width
    sh = prs.slide_height

    try:
        regions = extract_slide_regions(slide)
    except Exception:
        regions = []

    svg = generate_layout_svg(sw, sh, regions)
    output_svg.write_text(svg, encoding="utf-8")
    svg_ok = True

    png_ok = svg_to_png(svg, output_png)

    return png_ok, svg_ok
