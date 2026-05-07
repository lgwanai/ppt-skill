"""Element-level extraction — capture every visual detail from python-pptx shapes.

For each shape on a slide, produces an Element with:
  - Text: font_family, font_size_pt, bold, italic, color, alignment, line_spacing
  - Shape: fill_color, stroke_color, stroke_width, corner_radius, shape_type
  - Image: src path, crop info
  - Position: normalized (0–1) coordinates
  - Semantic role: title, body, subtitle, footer, decoration, etc.
"""

from __future__ import annotations

import re
import io
import os
from pathlib import Path
from typing import Any

from ppt_skill.spec.spec_model import (
    Element,
    ElementType,
    ImageStyle,
    Pos,
    SemanticRole,
    ShapeStyle,
    TextStyle,
)


# ── EMU conversion ───────────────────────────────────────────────────

EMU_PER_INCH = 914400
EMU_PER_PT = 12700
PT_PER_EMU = 1.0 / EMU_PER_PT


def emu_to_inches(emu: int) -> float:
    return emu / EMU_PER_INCH


def emu_to_pt(emu: int) -> float:
    return emu * PT_PER_EMU


# ── Main extraction ──────────────────────────────────────────────────


def extract_element(shape, slide_width_emu: int,
                    slide_height_emu: int,
                    z_order: int,
                    spec_dir: Path) -> Element:
    """Extract complete element info from a python-pptx shape.

    Returns Element with detailed text/shape/image styles.
    """
    sw = slide_width_emu
    sh = slide_height_emu

    pos = Pos(
        x=(shape.left or 0) / sw if sw else 0,
        y=(shape.top or 0) / sh if sh else 0,
        w=(shape.width or 0) / sw if sw else 0,
        h=(shape.height or 0) / sh if sh else 0,
    )

    shape_name = getattr(shape, "name", "")

    # Determine semantic role
    role = _infer_role(shape)

    # Build element
    elem = Element(
        id=z_order,
        element_type=_get_element_type(shape),
        semantic_role=role,
        position=pos,
        z_order=z_order,
        shape_name=shape_name,
    )

    # Extract text details
    try:
        if shape.has_text_frame:
            elem.text = shape.text_frame.text
            elem.text_style = _extract_text_style(shape.text_frame)
            # Build text from all paragraphs, skipping empty leading ones
            parts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            if parts:
                elem.text = "\n".join(parts)
    except (AttributeError, ValueError):
        pass

    # Extract shape details
    elem.shape_style = _extract_shape_style(shape)

    # Extract image details
    try:
        from pptx.shapes.picture import Picture
        if isinstance(shape, Picture):
            elem.image_style = _extract_image_style(shape, spec_dir, z_order)
            if elem.image_style:
                elem.element_type = ElementType.IMAGE
                elem.semantic_role = SemanticRole.IMAGE
    except ImportError:
        pass

    # Handle group shapes
    if hasattr(shape, "shapes"):
        for ci, child in enumerate(shape.shapes):
            child_elem = extract_element(child, sw, sh, ci, spec_dir)
            if child_elem.element_type == ElementType.TEXT and not elem.text:
                elem.text = child_elem.text
                elem.text_style = child_elem.text_style
            elem.children.append(child_elem)

    # Bubble up text from children
    if not elem.text and elem.children:
        for child in elem.children:
            if child.text:
                elem.text = child.text
                elem.text_style = child.text_style
                break

    return elem


# ── Element type detection ────────────────────────────────────────────


def _get_element_type(shape) -> ElementType:
    try:
        if hasattr(shape, "has_table") and shape.has_table:
            return ElementType.TABLE
    except Exception: pass
    try:
        if hasattr(shape, "chart"):
            shape.chart
            return ElementType.CHART
    except Exception: pass
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return ElementType.TEXT
    except Exception: pass
    if hasattr(shape, "shapes"):
        return ElementType.GROUP
    return ElementType.SHAPE


# ── Semantic role inference ──────────────────────────────────────────


def _infer_role(shape) -> SemanticRole:
    name = getattr(shape, "name", "").lower()
    text = shape.text_frame.text if shape.has_text_frame else ""

    if shape.has_text_frame:
        text_lower = text.lower()
        text_len = len(text.strip())
        try:
            pf = shape.placeholder_format
            if pf is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER
                ph_type = pf.type
                if ph_type == PP_PLACEHOLDER.TITLE:
                    return SemanticRole.TITLE
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    return SemanticRole.SUBTITLE
                elif ph_type == PP_PLACEHOLDER.BODY:
                    return SemanticRole.BODY
                elif ph_type in (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE):
                    return SemanticRole.PAGE_NUMBER
        except (ValueError, AttributeError, Exception):
            pass

    # Name-based inference
    if any(kw in name for kw in ["title", "标题", "heading"]):
        return SemanticRole.TITLE
    if any(kw in name for kw in ["subtitle", "副标题", "sub"]):
        return SemanticRole.SUBTITLE
    if any(kw in name for kw in ["footer", "页脚", "page number", "slide number"]):
        return SemanticRole.FOOTER
    if any(kw in name for kw in ["logo", "icon", "picture"]):
        return SemanticRole.LOGO
    if any(kw in name for kw in ["divider", "line", "分隔", "直线", "connector"]):
        return SemanticRole.DIVIDER
    if any(kw in name for kw in ["decor", "accent", "shape", "rectangle", "矩形", "oval", "椭圆"]):
        return SemanticRole.DECORATION

    # Text-based inference for shapes with text
    if shape.has_text_frame and text_len > 0:
        # Short bold text → likely title
        if text_len < 30 and "\n" not in text.strip():
            return SemanticRole.TITLE
        # Bullet list → body
        if "\n" in text and text_len > 50:
            return SemanticRole.BODY
        # Short footer-like text
        if text_len < 15:
            return SemanticRole.FOOTER

    # No text shape → decoration
    if not shape.has_text_frame or not text.strip():
        return SemanticRole.DECORATION

    return SemanticRole.BODY


# ── Text style extraction ────────────────────────────────────────────


def _extract_text_style(text_frame) -> TextStyle:
    """Extract font details from the first text run across all paragraphs."""
    style = TextStyle()

    if not text_frame.paragraphs:
        return style

    # Find first paragraph with runs
    first_run = None
    first_para = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_para = para
            first_run = para.runs[0]
            break

    if first_para is None:
        return style

    # Alignment from the first paragraph with content
    try:
        from pptx.enum.text import PP_ALIGN
        align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }
        style.text_alignment = align_map.get(first_para.alignment, "left")
    except Exception:
        pass

    # Line spacing
    if first_para.line_spacing is not None and first_para.line_spacing > 0:
        style.line_spacing = first_para.line_spacing

    # Extract from first run
    font = first_run.font

    # Font family
    if font.name:
        style.font_family = font.name

    # Font size
    if font.size:
        style.font_size_pt = font.size / EMU_PER_PT

    # Bold / italic / underline
    style.font_weight = "bold" if font.bold else "normal"
    style.font_italic = bool(font.italic) if font.italic is not None else False
    style.font_underline = bool(font.underline) if font.underline is not None else False

    # Font color
    try:
        if font.color and font.color.rgb:
            style.font_color = f"#{font.color.rgb}"
    except AttributeError:
        # Theme color — try type_color
        try:
            if font.color and font.color.type_color:
                style.font_color = f"#{font.color.type_color}"
        except Exception:
            pass
    except Exception:
        pass

    return style


# ── Shape style extraction ────────────────────────────────────────────


def _extract_shape_style(shape) -> ShapeStyle:
    """Extract fill, stroke, and geometry details from a shape."""
    style = ShapeStyle()

    if not hasattr(shape, "fill"):
        return style

    # Fill
    try:
        fill = shape.fill
        if fill.type is not None:
            from pptx.enum.dml import MSO_FILL_TYPE
            if fill.type == MSO_FILL_TYPE.SOLID:
                if fill.fore_color and fill.fore_color.rgb:
                    style.fill_color = f"#{fill.fore_color.rgb}"
                try:
                    style.fill_opacity = fill.fore_color.brightness or 1.0
                except Exception:
                    pass
    except Exception:
        pass

    # Stroke / Line
    try:
        if hasattr(shape, "line") and shape.line:
            line = shape.line
            if line.fill.type is not None:
                from pptx.enum.dml import MSO_FILL_TYPE
                if line.fill.type == MSO_FILL_TYPE.SOLID:
                    if line.fill.fore_color and line.fill.fore_color.rgb:
                        style.stroke_color = f"#{line.fill.fore_color.rgb}"
            if line.width:
                style.stroke_width_pt = line.width / EMU_PER_PT
    except Exception:
        pass

    # Geometry type
    try:
        style.shape_type = _get_geometry_type(shape)
    except Exception:
        pass

    # Corner radius for rounded rects
    try:
        if hasattr(shape, "_element"):
            xml = shape._element.xml
            import re
            m = re.search(r'<a:prstGeom[^>]*prst="([^"]*)"', xml)
            if m:
                style.shape_type = m.group(1)
            # Look for rounded corner radius
            m2 = re.search(r'<a:roundRect[^>]*av="(\d+)"', xml)
            if m2:
                style.corner_radius_pt = int(m2.group(1)) / EMU_PER_PT
    except Exception:
        pass

    return style


def _get_geometry_type(shape) -> str:
    """Get human-readable geometry type name."""
    if not hasattr(shape, "shape_type") or shape.shape_type is None:
        return "rect"

    st = str(shape.shape_type).upper()

    if "RECTANGLE" in st:
        # Check if rounded
        try:
            xml = shape._element.xml
            if "roundRect" in xml:
                return "round_rect"
        except Exception:
            pass
        return "rect"
    if "ROUNDED" in st:
        return "round_rect"
    if "OVAL" in st or "ELLIPSE" in st:
        return "circle"
    if "LINE" in st or "CONNECTOR" in st:
        return "line"
    if "TRIANGLE" in st:
        return "triangle"
    if "ARROW" in st or "CHEVRON" in st:
        return "arrow"
    if "PARALLELOGRAM" in st:
        return "parallelogram"
    if "DIAMOND" in st or "RHOMBUS" in st:
        return "diamond"
    if "HEXAGON" in st:
        return "hexagon"
    if "PENTAGON" in st:
        return "pentagon"
    if "CIRCLE" in st:
        return "circle"
    if "TRAPEZOID" in st:
        return "trapezoid"
    if "PIE" in st:
        return "pie"
    if "CLOUD" in st or "CALLOUT" in st:
        return "callout"
    if "STAR" in st:
        return "star"
    if "WAVE" in st:
        return "wave"
    if "GROUP" in st:
        return "group"
    if "PLACEHOLDER" in st:
        return "placeholder"

    return "rect"


# ── Image extraction ─────────────────────────────────────────────────


def _extract_image_style(picture, spec_dir: Path, idx: int) -> ImageStyle | None:
    """Extract image details and save to assets/."""
    try:
        image = picture.image
        ext = image.content_type.split("/")[-1]
        if ext == "jpeg":
            ext = "jpg"

        assets_dir = spec_dir / "assets"
        assets_dir.mkdir(parents=True, exist_ok=True)

        filename = f"img_{idx:03d}.{ext}"
        filepath = assets_dir / filename
        filepath.write_bytes(image.blob)

        return ImageStyle(
            src=f"assets/{filename}",
            original_width_px=getattr(image, "width", 0) or 0,
            original_height_px=getattr(image, "height", 0) or 0,
            crop_left=getattr(picture, "crop_left", 0.0) or 0.0,
            crop_right=getattr(picture, "crop_right", 0.0) or 0.0,
            crop_top=getattr(picture, "crop_top", 0.0) or 0.0,
            crop_bottom=getattr(picture, "crop_bottom", 0.0) or 0.0,
        )
    except Exception:
        return None


# ── Visual hierarchy summary ──────────────────────────────────────────


def build_hierarchy_summary(elements: list[Element]) -> list[str]:
    """Build a human-readable visual hierarchy description."""
    summary: list[str] = []
    seen_roles: set[str] = set()

    for elem in sorted(elements, key=lambda e: (e.position.y, e.position.x)):
        role = elem.semantic_role.value
        if role in seen_roles:
            continue
        seen_roles.add(role)

        if elem.text_style and elem.text_style.font_family:
            ts = elem.text_style
            summary.append(
                f"{role}: {ts.font_family} {ts.font_size_pt:.0f}pt "
                f"{ts.font_weight} {ts.font_color}"
            )
        elif elem.shape_style and elem.shape_style.fill_color:
            summary.append(
                f"{role}: {elem.shape_style.shape_type} "
                f"fill={elem.shape_style.fill_color}"
            )
        elif elem.image_style:
            summary.append(f"{role}: image")

    return summary
