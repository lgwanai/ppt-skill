"""Enhanced SpecExtractor — capture every visual detail for faithful style reproduction.

Each page is extracted as a complete blueprint (PageSpec with Element list) —
given this spec alone, one should be able to redraw the original slide.

Output: specs/<name>/ ├── spec.yaml ├── pages/{cover,toc,content,end_page}/page_N.yaml
├── assets/ └── logic.yaml
"""

from __future__ import annotations

import datetime
import os
import yaml
from pathlib import Path
from typing import Any

from pptx import Presentation

from ppt_skill.spec.spec_model import (
    ColorPalette,
    DesignSpec,
    Element,
    LayoutSubType,
    PageSpec,
    PageType,
    Pos,
    PresentationLogic,
    SemanticRole,
    Typography,
    VLModelConfig,
)
from ppt_skill.spec.theme import extract_theme_colors, extract_theme_fonts, extract_slide_background
from ppt_skill.spec.element_extractor import (
    Element,
    extract_element,
    build_hierarchy_summary,
    _infer_role,
)
from ppt_skill.spec.vision import VLClient
from ppt_skill.spec.asset_extractor import (
    extract_background_image,
    extract_shape_assets,
    AssetInfo,
)


# ── Page type detection ──────────────────────────────────────────────

_PAGE_TYPE_KEYWORDS: dict[str, list[str]] = {
    "cover": ["title slide", "title", "cover", "封面", "首页", "标题幻灯片"],
    "toc": ["table of contents", "contents", "agenda", "目录", "outline", "目 录"],
    "transition": ["section", "divider", "section header", "过渡", "分隔", "section divider"],
    "end_page": ["thank", "thanks", "q&a", "contact", "end", "close", "谢谢", "感谢", "致谢", "尾页"],
}

_COVER_KEYWORDS = ["关于", "项目", "方案", "计划", "申请", "报告", "立项", "预算", "采购", "述职", "总结", "规划"]


def _detect_page_type(layout_name: str, slide_idx: int,
                      total: int, elements: list[Element]) -> str:
    name_lower = layout_name.lower()
    text = " ".join(e.text for e in elements if e.text)

    if slide_idx == 0:
        for kw in _PAGE_TYPE_KEYWORDS["cover"]:
            if kw in name_lower:
                return "cover"
        # Check shape names for title indicators (Chinese: 标题)
        title_elems = [e for e in elements if e.semantic_role == SemanticRole.TITLE
                       or (e.text and len(e.text) > 10 and "\n" not in e.text.strip()
                           and e.position.y < 0.5)]
        # Also check element with "标题" in shape_name
        has_title_shape = any("标题" in e.shape_name or "title" in e.shape_name.lower()
                              for e in elements)
        if has_title_shape or len(title_elems) >= 1:
            if any(kw in text for kw in _COVER_KEYWORDS):
                return "cover"
            if len([e for e in elements if e.element_type.value == "text"]) >= 2:
                return "cover"

    if slide_idx == total - 1:
        for kw in _PAGE_TYPE_KEYWORDS["end_page"]:
            if kw in name_lower or kw in text.lower():
                return "end_page"
        if len(text.strip()) < 100 and total > 3:
            return "end_page"

    for kw in _PAGE_TYPE_KEYWORDS["toc"]:
        if kw in name_lower:
            return "toc"
    if slide_idx <= 2:
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if len(lines) >= 3 and all(len(l) < 80 for l in lines):
            return "toc"

    for kw in _PAGE_TYPE_KEYWORDS["transition"]:
        if kw in name_lower:
            return "transition"
    if len(text.strip()) < 200 and any(kw in text.lower() for kw in ["section", "part", "章节"]):
        return "transition"

    return "content"


def _detect_layout_sub_type(elements: list[Element]) -> str:
    img_elems = [e for e in elements if e.element_type.value == "image"]
    text_elems = [e for e in elements if e.element_type.value == "text"]

    if img_elems and text_elems:
        img_x = min(e.position.x for e in img_elems)
        text_x = min(e.position.x for e in text_elems)
        if img_x < text_x:
            return "image_left"
        else:
            return "image_right"

    if len(text_elems) == 2:
        xs = sorted(e.position.x for e in text_elems)
        if xs[1] - xs[0] > 0.3:
            return "left_right"

    if len(text_elems) == 3:
        return "left_middle_right"

    if all(e.position.w > 0.8 for e in text_elems):
        return "full_width"

    if len(text_elems) >= 3 and not img_elems:
        return "grid"

    return "top_bottom"


# ── Config ────────────────────────────────────────────────────────────

def _load_config() -> VLModelConfig:
    config_path = Path("config.txt")
    if not config_path.exists():
        return VLModelConfig(enabled=False)
    config = {}
    for line in config_path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if "=" in line:
            k, _, v = line.partition("=")
            config[k.strip()] = v.strip().strip('"').strip("'")
    return VLModelConfig(
        provider=config.get("VL_PROVIDER", "openai"),
        model=config.get("VL_MODEL", "gpt-4o"),
        api_key=config.get("VL_API_KEY", ""),
        api_base=config.get("VL_API_BASE", ""),
        enabled=config.get("VL_ENABLED", "false").lower() == "true",
    )


# ── SpecExtractor ────────────────────────────────────────────────────


class SpecExtractor:
    """Extract a complete design spec from a reference PPTX."""

    def __init__(self):
        self.config = _load_config()
        self.vl_client = VLClient(self.config) if self.config.enabled else None

    def extract(self, pptx_path: Path | str) -> DesignSpec:
        pptx_path = Path(pptx_path)
        prs = Presentation(str(pptx_path))
        total = len(prs.slides)

        # Metadata
        metadata = {
            "name": pptx_path.stem,
            "source_pptx": str(pptx_path.absolute()),
            "extracted_at": datetime.datetime.now().isoformat(),
            "slide_count": total,
            "format": f"{int(prs.slide_width) // 914400}x{int(prs.slide_height) // 914400}in",
        }

        # Theme
        pptx_str = str(pptx_path.absolute())
        theme_colors = extract_theme_colors(pptx_str)
        palette = ColorPalette.from_theme_scheme(theme_colors)
        fonts_data = extract_theme_fonts(pptx_str)
        typography = Typography(
            heading_family=fonts_data.get("majorFont", ""),
            body_family=fonts_data.get("minorFont", ""),
        )

        # Per-page
        pages: list[PageSpec] = []
        page_types: set[str] = set()
        sub_types: set[str] = set()
        seq: list[str] = []
        density_seq: list[str] = []

        for idx, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name if slide.slide_layout else ""
            sw, sh = int(prs.slide_width), int(prs.slide_height)

            # Extract all elements
            elements: list[Element] = []
            for zi, shape in enumerate(slide.shapes):
                try:
                    elem = extract_element(shape, sw, sh, zi, Path("specs") / metadata["name"])
                    if elem is not None:
                        elements.append(elem)
                except Exception:
                    continue

            # Page type
            ptype = _detect_page_type(layout_name, idx, total, elements)
            page_types.add(ptype)

            # Layout sub-type
            stype = "full_width"
            if ptype == "content":
                stype = _detect_layout_sub_type(elements)
                sub_types.add(stype)

            # ── Background ──
            spec_dir_path = Path("specs") / metadata["name"]
            bg_result = extract_background_image(slide, spec_dir_path, idx)
            bg_color = bg_result.get("color", "#FFFFFF") if bg_result else "#FFFFFF"
            bg_image = bg_result.get("image", "") if bg_result else ""
            bg_type = bg_result.get("type", "solid") if bg_result else "solid"
            bg_desc = bg_result.get("description", "White background") if bg_result else "White background"
            gradient_stops = bg_result.get("gradient_stops", []) if bg_result else []

            # ── Shape assets (classify + extract, discard content images) ──
            shape_assets: list[AssetInfo] = []
            for shape in slide.shapes:
                try:
                    assets = extract_shape_assets(shape, spec_dir_path, idx,
                                                   slide_w_emu=sw, slide_h_emu=sh)
                    shape_assets.extend(assets)
                except Exception:
                    pass

            # Update element image counts
            img_count = sum(1 for e in elements if e.element_type.value == "image")

            # Content density
            total_chars = sum(len(e.text) for e in elements if e.text)
            if total_chars < 100:
                density_seq.append("breathing")
            elif total_chars < 500:
                density_seq.append("dense")
            else:
                density_seq.append("anchor")

            # Visual hierarchy
            hierarchy = build_hierarchy_summary(elements)

            # Layout description
            text_count = sum(1 for e in elements if e.element_type.value == "text")
            shape_count = sum(1 for e in elements if e.element_type.value == "shape")
            asset_count = len(shape_assets)

            desc_parts = [
                f"{PageType(ptype).name}: {len(elements)} elements "
                f"({text_count} texts, {img_count} images, {shape_count} shapes).",
                f"Layout: {stype}.",
                f"Hierarchy: {' → '.join(hierarchy) if hierarchy else 'flat'}.",
            ]
            if bg_image:
                desc_parts.append(f"Background image: {bg_image}.")
            elif bg_color and bg_color != "#FFFFFF":
                desc_parts.append(f"Background: {bg_color}.")
            if asset_count > 0:
                desc_parts.append(f"{asset_count} extracted assets.")
            layout_desc = " ".join(desc_parts)

            page = PageSpec(
                page_type=PageType(ptype),
                layout_sub_type=LayoutSubType(stype),
                width_emu=sw,
                height_emu=sh,
                width_inches=sw / 914400,
                height_inches=sh / 914400,
                background_color=bg_color,
                background_image=bg_image,
                background_type=bg_type,
                background_description=bg_desc,
                gradient_stops=gradient_stops,
                visual_hierarchy=hierarchy,
                elements=elements,
                layout_description=layout_desc,
            )
            pages.append(page)
            seq.append(ptype)

        # Logic
        sections = _build_sections(pages)
        logic = PresentationLogic(
            page_sequence=seq,
            sections=sections,
            density_sequence=density_seq,
            avg_content_per_page=sum(
                len(e.text) for p in pages for e in p.elements
            ) // max(total, 1),
        )

        return DesignSpec(
            metadata=metadata,
            palette=palette,
            typography=typography,
            pages=pages,
            logic=logic,
            page_types_found=sorted(page_types),
            layout_sub_types_found=sorted(sub_types),
        )

    def save(self, spec: DesignSpec, base_dir: str = "specs"):
        """Save spec as directory structure with per-page YAML blueprints."""
        spec_dir = Path(base_dir) / spec.metadata["name"]
        spec_dir.mkdir(parents=True, exist_ok=True)

        # spec.yaml
        (spec_dir / "spec.yaml").write_text(
            yaml.dump(spec.to_dict(), default_flow_style=False,
                     sort_keys=False, allow_unicode=True),
            encoding="utf-8",
        )

        # logic.yaml
        (spec_dir / "logic.yaml").write_text(
            yaml.dump(spec.logic.to_dict(), default_flow_style=False,
                     sort_keys=False, allow_unicode=True),
            encoding="utf-8",
        )

        # Per-page YAMLs
        for i, page in enumerate(spec.pages):
            sub = spec_dir / "pages"
            sub = sub / page.page_type.value
            if page.page_type == PageType.CONTENT:
                sub = sub / page.layout_sub_type.value
            sub.mkdir(parents=True, exist_ok=True)

            (sub / f"page_{i}.yaml").write_text(
                yaml.dump(page.to_dict(), default_flow_style=False,
                         sort_keys=False, allow_unicode=True),
                encoding="utf-8",
            )

        print(f"Spec: {spec_dir}/")
        print(f"  Pages: {len(spec.pages)} ({', '.join(spec.page_types_found)})")
        if spec.layout_sub_types_found:
            print(f"  Layouts: {', '.join(spec.layout_sub_types_found)}")


def _build_sections(pages: list[PageSpec]) -> list[dict]:
    sections: list[dict] = []
    cur: dict | None = None
    for i, p in enumerate(pages):
        if p.page_type in (PageType.TRANSITION, PageType.TOC):
            if cur:
                sections.append(cur)
            cur = {"name": f"Section {len(sections) + 1}", "slides": [i], "types": [p.page_type.value]}
        elif p.page_type == PageType.COVER:
            if cur:
                sections.append(cur)
            cur = {"name": "Opening", "slides": [i], "types": [p.page_type.value]}
        elif p.page_type == PageType.END_PAGE:
            if cur:
                cur["slides"].append(i)
                cur["types"].append(p.page_type.value)
                sections.append(cur)
            else:
                sections.append({"name": "Closing", "slides": [i], "types": [p.page_type.value]})
            cur = None
        else:
            if cur is None:
                cur = {"name": f"Section {len(sections) + 1}", "slides": [], "types": []}
            cur["slides"].append(i)
            cur["types"].append(p.page_type.value)
    if cur and cur.get("slides"):
        sections.append(cur)
    return sections
