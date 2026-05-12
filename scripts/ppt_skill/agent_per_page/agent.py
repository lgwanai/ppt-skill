"""PageAgent — one agent per PPT page. Orchestrates tools for single-page generation."""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu
from pathlib import Path
import tempfile, os

from .tools import (
    get_page, extract_spec, plan_assets, plan_layout,
    draw_zone, screenshot_pptx, review_slide,
)

SW, SH = 12192000, 6858000

class PageAgent:
    """Agent responsible for generating one PPT page."""
    
    def __init__(self, outline_path: str, spec_dir: str, page_index: int):
        self.outline_path = outline_path
        self.spec_dir = spec_dir
        self.page_index = page_index
        self.content = None
        self.spec = None
        self.assets = None
        self.layout = None
        self.iterations = 0

    def run(self, prs: Presentation) -> bool:
        """Execute full agent loop for this page. Returns True on success."""
        self.content = get_page(self.outline_path, set())
        if not self.content:
            return False

        self.spec = extract_spec(self.content, self.spec_dir)
        self.assets = plan_assets(self.content, self.spec)
        self.layout = plan_layout(self.content, self.spec, self.assets)

        # Generate slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._draw_layout(slide)

        # Agent loop: screenshot → review → fix → repeat
        for it in range(5):
            self.iterations = it + 1
            
            # Save temp PPTX and screenshot
            tmp_pptx = f'/tmp/agent_page_{self.page_index}_{it}.pptx'
            tmp_prs = Presentation(); tmp_prs.slide_width = SW; tmp_prs.slide_height = SH
            import copy
            # Note: python-pptx can't easily copy slides between presentations
            # For now, save the full prs and screenshot
            prs.save(tmp_pptx)
            
            png = screenshot_pptx(tmp_pptx)
            if not png:
                break
            
            review = review_slide(png, self.spec, self.content)
            if review.get('pass', False):
                break
            
            # Fix issues and redraw
            for issue in review.get('issues', []):
                self._fix_issue(slide, issue)

        return True

    def _draw_layout(self, slide):
        """Render all layout zones with content."""
        zones = self.layout.get('zones', [])
        if not zones:
            return

        # Background
        bg = self.spec.get('background', {})
        if bg.get('color'):
            s = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(SW), Emu(SH))
            s.fill.solid()
            s.fill.fore_color.rgb = RGBColor(*_hx(bg['color']))
            s.line.fill.background()

        # Draw each zone
        body_items = self.content.get('body', [])
        body_idx = 0
        table_idx = 0

        for zone in zones:
            ct = zone.get('content', 'body')
            
            if ct == 'title':
                draw_zone(slide, zone, self.spec, self.content.get('title', ''), 'title')
            elif ct.startswith('body') and body_idx < len(body_items):
                item = body_items[body_idx]
                if item.startswith('__TABLE__'):
                    table_content = '__ROW__'.join(
                        b[9:] for b in body_items if b.startswith('__TABLE__')
                    )
                    draw_zone(slide, zone, self.spec, table_content, 'table')
                    table_idx += 1
                else:
                    draw_zone(slide, zone, self.spec, item, 'text')
                    body_idx += 1
            elif body_idx < len(body_items):
                item = body_items[body_idx]
                if not item.startswith('__TABLE__'):
                    draw_zone(slide, zone, self.spec, item, 'text')
                    body_idx += 1

    def _fix_issue(self, slide, issue: dict):
        """Apply fix from review feedback to the slide."""
        # LLM-driven fix — adjust element properties based on issue type
        issue_type = issue.get('type', '')
        fix_detail = issue.get('fix', '')
        # For now, log issues for manual review
        print(f"  Page {self.page_index} issue: [{issue_type}] {fix_detail}")


def _hx(h):
    h = h.lstrip('#')
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
