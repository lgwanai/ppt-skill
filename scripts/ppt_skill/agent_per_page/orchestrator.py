"""Orchestrator — runs PageAgents sequentially for all pages, assembles final PPTX."""
from pptx import Presentation
from .agent import PageAgent

SW, SH = 12192000, 6858000


def generate_pptx_agent(outline_path: str, spec_dir: str, output_path: str) -> str:
    """Generate PPTX using agent-per-page architecture."""
    from .tools import get_page
    
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    
    completed = set()
    total = len(get_page(outline_path, completed) or [])
    
    # First count pages
    pages = []
    while True:
        page = get_page(outline_path, completed)
        if not page:
            break
        pages.append(page)
        completed.add(page['index'])
    
    completed = set()
    
    for page in pages:
        agent = PageAgent(outline_path, spec_dir, page['index'])
        agent.content = page  # pre-set content to avoid re-parsing
        agent.spec = __import__('ppt_skill.agent_per_page.tools', fromlist=['extract_spec']).extract_spec(page, spec_dir)
        agent.assets = __import__('ppt_skill.agent_per_page.tools', fromlist=['plan_assets']).plan_assets(page, agent.spec)
        agent.layout = __import__('ppt_skill.agent_per_page.tools', fromlist=['plan_layout']).plan_layout(page, agent.spec, agent.assets)
        
        print(f"Page {page['index']} [{page['type']}]: {page['title'][:40]}")
        print(f"  Assets: {len(agent.assets.get('modules',[]))} modules")
        print(f"  Layout: {agent.layout.get('layout','?')} ({len(agent.layout.get('zones',[]))} zones)")
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        agent._draw_layout(slide)
        
        completed.add(page['index'])
    
    prs.save(output_path)
    print(f"\nSaved: {output_path} ({len(prs.slides)} slides)")
    return output_path
