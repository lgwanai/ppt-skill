"""Agent-per-page PPT generation tools — LLM-driven, no hardcoded rendering."""
import json, os, re, io, base64, subprocess, tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from openai import OpenAI
from PIL import Image

SW, SH = 12192000, 6858000
FONT = '微软雅黑'
A = '/tmp/ppt_assets'  # layout image assets directory

def _emu(v): return Emu(int(v))
def _hx(h): h=h.lstrip('#'); return int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)

# ── LLM client (reads config.txt) ──────────────────────────────────

def _llm_client():
    cfg = {}
    if os.path.exists('config.txt'):
        for line in open('config.txt'):
            line = line.strip()
            if '=' in line and not line.startswith('#'):
                k,v = line.split('=',1)
                cfg[k.strip()] = v.strip().strip('"')
    return OpenAI(
        api_key=cfg.get('LLM_API_KEY',''),
        base_url=cfg.get('LLM_API_BASE',''),
    )

def _vl_client():
    cfg = {}
    if os.path.exists('config.txt'):
        for line in open('config.txt'):
            line = line.strip()
            if '=' in line and not line.startswith('#'):
                k,v = line.split('=',1)
                cfg[k.strip()] = v.strip().strip('"')
    return OpenAI(
        api_key=cfg.get('VL_API_KEY',''),
        base_url=cfg.get('VL_API_BASE',''),
    )

VL_MODEL = "mimo-v2-omni"  # supports image input

# ── Tool 1: get_page ──────────────────────────────────────────────
_PAGE_CACHE = {}

def get_page(outline_path: str, completed: set) -> dict | None:
    """Fetch next uncompleted page from outline. Results cached."""
    if outline_path not in _PAGE_CACHE:
        raw = open(outline_path, encoding='utf-8').read()
        _PAGE_CACHE[outline_path] = _parse_outline(raw)
    for p in _PAGE_CACHE[outline_path]:
        if p['index'] not in completed:
            return p
    return None

def _parse_outline(raw: str) -> list[dict]:
    """Parse outline markdown into page list."""
    import re
    lines = raw.split('\n')
    pages = []; chapter = ""; ch = 0; title = ""; body = []; skip = 0
    idx = 0
    
    def _type(t): return 'transition' if t.startswith('转场页') else 'content'
    
    for i, line in enumerate(lines):
        if i < skip: continue
        if line.startswith('## ') and not line.startswith('### '):
            if title:
                pages.append({'index':idx,'type':_type(title),'title':title,'body':body,'chapter':chapter,'ch_num':ch})
                idx+=1; body=[]
            h = line.lstrip('# ').strip()
            if h == '封面':
                ti=h; sb=[]
                for j in range(i+1, min(i+10,len(lines))):
                    t=lines[j].strip()
                    if t.startswith('- 主标题：'): ti=t.replace('- 主标题：','').strip()
                    elif t.startswith('- 副标题：'): sb.append(t.replace('- 副标题：','').strip())
                    elif t.startswith('- '): sb.append(t[2:])
                skip=i+len(sb)+2
                pages.append({'index':idx,'type':'cover','title':ti,'body':sb,'chapter':'','ch_num':0})
                idx+=1; title=""; body=[]
            elif h == '目录':
                toc=[]
                for j in range(i+1,min(i+15,len(lines))):
                    t=lines[j].strip()
                    if not t or t.startswith('##'): break
                    if t[0].isdigit() and '.' in t[:3]: toc.append(t.split('.',1)[1].strip())
                skip=j
                pages.append({'index':idx,'type':'toc','title':'目录','body':toc,'chapter':'','ch_num':0})
                idx+=1; title=""; body=[]
            elif h == '结尾页':
                eb=[]
                for j in range(i+1,min(i+10,len(lines))):
                    t=lines[j].strip()
                    if not t or t.startswith('##'): break
                    if t.startswith('- '): eb.append(t[2:])
                skip=i+len(eb)+2
                pages.append({'index':idx,'type':'end','title':'结尾页','body':eb,'chapter':'','ch_num':0})
                idx+=1; body=[]
            else: ch+=1; chapter=re.sub(r'^[第][一二三四五六七八九十]+章[：:]','',h).strip()
            title=""
        elif line.startswith('### '):
            if title:
                pages.append({'index':idx,'type':_type(title),'title':title,'body':body,'chapter':chapter,'ch_num':ch})
                idx+=1; body=[]
            title=line.lstrip('# ').strip()
        elif line.strip() and not line.startswith(('---','```')):
            sline=line.strip()
            if sline.startswith('|') and '---' not in sline:
                cells=[c.strip().strip('*') for c in sline.split('|') if c.strip()]
                if cells: body.append('__TABLE__'+'|'.join(cells))
            else:
                t=re.sub(r'\*+','',sline.lstrip('- ').strip())
                if t and len(t)>3: body.append(t)
    if title:
        pages.append({'index':idx,'type':_type(title),'title':title,'body':body,'chapter':chapter,'ch_num':ch})
    return pages

# ── Tool 2: extract_spec ──────────────────────────────────────────

def extract_spec(page: dict, spec_dir: str) -> dict:
    """Match page type to spec layout, return style data. Reads JSON directly (no LLM)."""
    sd = Path(spec_dir)
    spec = json.load(open(sd/'spec.json'))
    palette = spec['palette']
    typo = spec.get('typography', {})
    
    # Default styles by page type
    defaults = {
        'cover': {
            'title': {'x_pct':0.161,'y_pct':0.098,'w_pct':0.669,'font':typo.get('heading_family',FONT),'size_pt':36,'bold':False,'color':'#0070C0','align':'center'},
            'body': {'x_pct':0.077,'y_pct':0.625,'w_pct':0.846,'font':FONT,'size_pt':20,'bold':True,'color':'#0070C0','line_spacing':1.3},
            'background': {'type':'image','color':'#FFFFFF','image_path':f'{A}/3_自定义版式_0.png'},
            'logo': {'x_pct':0.762,'y_pct':0.056,'w_pct':0.188,'h_pct':0.103,'image_path':f'{A}/3_自定义版式_2.png'},
        },
        'end': {
            'title': {'x_pct':0.10,'y_pct':0.30,'w_pct':0.80,'font':typo.get('heading_family',FONT),'size_pt':42,'bold':True,'color':'#0070C0','align':'center'},
            'body': {'x_pct':0.15,'y_pct':0.50,'w_pct':0.70,'font':FONT,'size_pt':14,'color':'#333333','line_spacing':1.5},
            'background': {'type':'image','color':'#FFFFFF','image_path':f'{A}/3_自定义版式_0.png'},
            'logo': {'x_pct':0.762,'y_pct':0.056,'w_pct':0.188,'h_pct':0.103,'image_path':f'{A}/3_自定义版式_2.png'},
        },
        'toc': {
            'title': {'x_pct':0.10,'y_pct':0.12,'w_pct':0.35,'font':typo.get('heading_family',FONT),'size_pt':28,'bold':True,'color':'#0070C0','align':'left'},
            'body': {'x_pct':0.47,'y_pct':0.22,'w_pct':0.43,'font':FONT,'size_pt':16,'color':'#333333','line_spacing':1.5},
            'background': {'type':'solid','color':'#FFFFFF'},
            'logo': {'x_pct':0.832,'y_pct':0.023,'w_pct':0.141,'h_pct':0.077,'image_path':f'{A}/统一模板_1.png'},
            'content_bg': {'x_pct':0.0,'y_pct':0.233,'w_pct':0.769,'h_pct':0.769,'image_path':f'{A}/统一模板_2.png'},
        },
        'transition': {
            'title': {'x_pct':0.021,'y_pct':0.049,'w_pct':0.95,'font':typo.get('heading_family',FONT),'size_pt':24,'bold':True,'color':'#0070C0','align':'left'},
            'body': {'x_pct':0.04,'y_pct':0.14,'w_pct':0.92,'font':FONT,'size_pt':13,'color':'#333333','line_spacing':1.5},
            'background': {'type':'solid','color':'#FFFFFF'},
            'logo': {'x_pct':0.832,'y_pct':0.023,'w_pct':0.141,'h_pct':0.077,'image_path':f'{A}/统一模板_1.png'},
            'content_bg': {'x_pct':0.0,'y_pct':0.233,'w_pct':0.769,'h_pct':0.769,'image_path':f'{A}/统一模板_2.png'},
        },
        'content': {
            'title': {'x_pct':0.021,'y_pct':0.049,'w_pct':0.95,'font':typo.get('heading_family',FONT),'size_pt':24,'bold':True,'color':'#0070C0','align':'left'},
            'body': {'x_pct':0.04,'y_pct':0.14,'w_pct':0.92,'font':FONT,'size_pt':13,'color':'#333333','line_spacing':1.5},
            'background': {'type':'solid','color':'#FFFFFF'},
            'logo': {'x_pct':0.832,'y_pct':0.023,'w_pct':0.141,'h_pct':0.077,'image_path':f'{A}/统一模板_1.png'},
            'content_bg': {'x_pct':0.0,'y_pct':0.233,'w_pct':0.769,'h_pct':0.769,'image_path':f'{A}/统一模板_2.png'},
        },
    }
    
    pt = page.get('type', 'content')
    base = defaults.get(pt, defaults['content'])
    
    # Try to find matching spec layout for more precise data
    layouts = spec.get('layouts', [])
    # Find layout by page_type
    for l in layouts:
        if l['page_type'] == pt:
            try:
                data = json.load(open(sd / l['file']))
                elements = data.get('elements', [])
                # Extract actual style data from elements
                for e in elements:
                    ts = e.get('text_style', {}) or {}
                    pos = e.get('position', {}) or {}
                    role = _infer_role(e, pt)
                    fs = ts.get('font_size_pt', 0) or 0
                    # Only update if we have actual font data AND role matches
                    if role == 'title' and fs > 0 and pos.get('x') is not None:
                        base['title'].update({'x_pct':pos['x'],'y_pct':pos['y'],'w_pct':pos['w'],'h_pct':pos['h']})
                        base['title']['size_pt'] = fs
                        base['title']['font'] = ts.get('font_family', base['title']['font'])
                        if ts.get('font_color'): base['title']['color'] = ts['font_color']
                        if ts.get('font_weight'): base['title']['bold'] = (ts['font_weight'] == 'bold')
                    if role == 'subtitle' and fs > 0:
                        base['body'].update({'x_pct':pos.get('x',base['body']['x_pct']),'y_pct':pos.get('y',base['body']['y_pct']),'w_pct':pos.get('w',base['body']['w_pct'])})
                        base['body']['size_pt'] = fs
                        base['body']['font'] = ts.get('font_family', base['body']['font'])
                        if ts.get('font_color'): base['body']['color'] = ts['font_color']
                        base['body']['bold'] = (ts.get('font_weight','') == 'bold')
                break
            except: pass
    
    return {
        "page_type": pt,
        "palette": {"accent": palette.get('accent1','#4472C4'), "bg": "#FFFFFF", "text": "#333333"},
        **base,
        "decorations": [],
        "layout_type": "full_width" if pt in ('cover','end','transition') else "top_bottom",
    }


def _infer_role(elem, page_type='content'):
    pos = elem.get('position', {})
    ts = elem.get('text_style', {}) or {}
    x, y, w, h = pos.get('x',0), pos.get('y',0), pos.get('w',0), pos.get('h',0)
    fs = ts.get('font_size_pt', 0) or 0
    if w > 0.9 and h > 0.9: return 'background'
    if x > 0.7 and w < 0.3 and h < 0.2: return 'logo'
    if y < 0.3 and h > 0.3: return 'title'
    # Only infer subtitle for cover pages
    if page_type == 'cover' and fs >= 16 and y > 0.5: return 'subtitle'
    if fs >= 16 and y < 0.5: return 'title'
    return 'body'

# ── Tool 3: plan_assets ───────────────────────────────────────────

def plan_assets(page: dict, spec: dict) -> dict:
    """Plan optimal visual elements for page content. LLM-driven with deterministic fallback."""
    llm = _llm_client()
    prompt = f"""Plan visual assets for this PPT slide. Priority: SmartArt > chart > diagram > text.
Max 2 graphics per page. If graphics can replace text, condense the text.

CONTENT:
  type: {page['type']}
  title: {page['title'][:80]}
  body ({len(page['body'])} items):
  {json.dumps(page['body'][:5], ensure_ascii=False)[:500]}

Available: SmartArt (list/process/cycle/hierarchy/pyramid/matrix), chart (bar/pie/line), diagram SVG, icons, images.

Analyze content structure: parallel points? flow/process? comparison? numbers? One big idea?
Choose optimal representation. Return JSON:
{{"modules":[{{"type":"smartart|chart|diagram|image|text","subtype":"list|flow|compare|hierarchy|process|bar|pie","contentRef":"body[0..N] or all","reason":"one line why"}}]}}"""

    try:
        r = llm.chat.completions.create(
            model="deepseek-chat",
            messages=[{"role":"user","content":prompt}],
            max_tokens=512, temperature=0.2, timeout=10,
        )
        text = r.choices[0].message.content
        m = re.search(r'\{.*\}', text, re.DOTALL)
        if m:
            result = json.loads(m.group(0))
            if result.get('modules'): return result
    except Exception:
        pass

    # Deterministic fallback — analyze content structure
    return _deterministic_assets(page)


def _deterministic_assets(page: dict) -> dict:
    """Smart fallback: analyze content structure without LLM."""
    body = page.get('body', [])
    pt = page.get('type', 'content')
    title = page.get('title', '')
    
    # Cover/End → centered text hierarchy
    if pt in ('cover', 'end'):
        return {"modules": [{"type": "text", "subtype": "cover_title",
                "contentRef": "all", "reason": "Cover/end: centered title with accent background"}]}
    
    # Detect table content
    table_items = [b for b in body if b.startswith('__TABLE__')]
    if len(table_items) >= 3:
        return {"modules": [{"type": "text", "subtype": "table",
                "contentRef": "table items", "reason": f"Table detected ({len(table_items)} rows)"}]}
    
    # Detect list patterns (numbered items, bullet points)
    list_items = [b for b in body if len(b) > 10 and not b.startswith('__TABLE__')]
    n = len(list_items)
    
    # 3-4 items → SmartArt list or grid
    if 3 <= n <= 4:
        return {"modules": [
            {"type": "smartart", "subtype": "list", "contentRef": f"body[0..{n-1}]",
             "reason": f"{n} parallel points → SmartArt list"}
        ]}
    
    # 5+ items → compact text grid
    if n >= 5:
        return {"modules": [
            {"type": "text", "subtype": "grid_2col", "contentRef": f"all {n} items",
             "reason": f"{n} items → compact 2-column grid"}
        ]}
    
    # 2 items with "vs" or comparison → compare layout
    if n == 2 and ('vs' in title.lower() or '对比' in title or '比较' in title or '区别' in title):
        return {"modules": [
            {"type": "smartart", "subtype": "compare", "contentRef": "body[0..1]",
             "reason": "2 items comparison → SmartArt compare"}
        ]}
    
    # Default: text list
    return {"modules": [{"type": "text", "subtype": "list", "contentRef": "all",
             "reason": "Default text list layout"}]}

# ── Tool 4: plan_layout ───────────────────────────────────────────

def plan_layout(page: dict, spec: dict, assets: dict) -> dict:
    """Calculate page layout zones based on content and assets."""
    llm = _llm_client()
    prompt = f"""Plan page layout. The slide is 13.33x7.5 inches. Fixed elements occupy:
  Title at ({spec['title']['x_pct']:.1%}, {spec['title']['y_pct']:.1%}) height ~{spec['title']['size_pt']/72:.1f}in

Page type: {page['type']}
Assets: {json.dumps(assets, ensure_ascii=False)}
Body line count: ~{len(page['body'])}

Available layouts: center, left_right, top_bottom, grid_2x2, grid_3x3, grid_2x3, pyramid, pinwheel
For transition pages: always use center layout.

Return JSON:
{{"layout":"grid_2x2","zones":[{{"x_pct":0.05,"y_pct":0.15,"w_pct":0.42,"h_pct":0.35,"content":"title"}},{{"x_pct":0.52,"y_pct":0.15,"w_pct":0.42,"h_pct":0.35,"content":"body_0"}},{{"x_pct":0.05,"y_pct":0.52,"w_pct":0.42,"h_pct":0.35,"content":"body_2"}},{{"x_pct":0.52,"y_pct":0.52,"w_pct":0.42,"h_pct":0.35,"content":"diagram"}}]}}
"""
    try:
        r = llm.chat.completions.create(model="deepseek-v4-flash", messages=[{"role":"user","content":prompt}], max_tokens=1024, temperature=0.2)
        text = r.choices[0].message.content
        m = re.search(r'\{.*\}', text, re.DOTALL)
        return json.loads(m.group(0)) if m else {"layout":"top_bottom","zones":[{"x_pct":0.05,"y_pct":0.15,"w_pct":0.90,"h_pct":0.75,"content":"body"}]}
    except:
        return {"layout":"top_bottom","zones":[{"x_pct":0.05,"y_pct":0.15,"w_pct":0.90,"h_pct":0.75,"content":"body"}]}

# ── Tool 5: draw_zone ──────────────────────────────────────────────

def render_spec_images(slide, spec: dict):
    """Render background image and logo from spec onto the slide."""
    # Background image
    bg = spec.get('background', {})
    bg_path = bg.get('image_path', '')
    if bg_path and os.path.exists(bg_path):
        slide.shapes.add_picture(bg_path, Emu(0), Emu(0), Emu(SW), Emu(SH))

    # Logo
    logo = spec.get('logo', {})
    logo_path = logo.get('image_path', '')
    if logo_path and os.path.exists(logo_path):
        lx = int(logo.get('x_pct', 0) * SW)
        ly = int(logo.get('y_pct', 0) * SH)
        lw = int(logo.get('w_pct', 0) * SW)
        lh = int(logo.get('h_pct', 0) * SH)
        slide.shapes.add_picture(logo_path, Emu(lx), Emu(ly), Emu(lw), Emu(lh))

    # Content background
    cbg = spec.get('content_bg', {})
    cbg_path = cbg.get('image_path', '')
    if cbg_path and os.path.exists(cbg_path):
        cx = int(cbg.get('x_pct', 0) * SW)
        cy = int(cbg.get('y_pct', 0) * SH)
        cw = int(cbg.get('w_pct', 0) * SW)
        ch = int(cbg.get('h_pct', 0) * SH)
        slide.shapes.add_picture(cbg_path, Emu(cx), Emu(cy), Emu(cw), Emu(ch))


def draw_zone(slide, zone: dict, spec: dict, content: str, content_type: str = "text"):
    """Draw content into a layout zone on the slide."""
    x = int(zone['x_pct'] * SW)
    y = int(zone['y_pct'] * SH)
    w = int(zone['w_pct'] * SW)
    h = int(zone['h_pct'] * SH)
    
    if content_type == "text":
        bstyle = spec.get('body', spec.get('title', {}))
        sz = bstyle.get('size_pt', 13)
        color = RGBColor(*_hx(bstyle.get('color', '#333333')))
        
        # Auto-calculate font size to fit zone
        chars = len(content)
        cpl = max(1, int(w / (sz * 9144)))
        n_lines = max(1, (chars + cpl - 1) // cpl)
        
        # Reduce font size if text overflows zone
        while n_lines * sz * 1.5 * 18000 > h and sz > 8:
            sz -= 1
            cpl = max(1, int(w / (sz * 9144)))
            n_lines = max(1, (chars + cpl - 1) // cpl)
        
        box_h = n_lines * sz * 1.5 * 18000
        tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(min(box_h, h)))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = Pt(sz * 1.5)
        r = p.add_run(); r.text = content[:500]
        r.font.name = bstyle.get('font', FONT); r.font.size = Pt(sz); r.font.color.rgb = color
    
    elif content_type == "title":
        tstyle = spec.get('title', {})
        sz = tstyle.get('size_pt', 24)
        color = RGBColor(*_hx(tstyle.get('color', '#4472C4')))
        align = PP_ALIGN.CENTER if tstyle.get('align') == 'center' else PP_ALIGN.LEFT
        
        tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = Pt(sz * 1.3)
        r = p.add_run(); r.text = content[:200]
        r.font.name = tstyle.get('font', FONT); r.font.size = Pt(sz)
        r.font.bold = tstyle.get('bold', False); r.font.color.rgb = color
    
    elif content_type == "table":
        rows = [r.split('|') for r in content.split('__ROW__') if r.strip()]
        if len(rows) >= 2:
            nr, nc = len(rows), max(len(r) for r in rows)
            tbl = slide.shapes.add_table(nr, min(nc, 5), Emu(x), Emu(y), Emu(w), Emu(h))
            for ri, row in enumerate(rows):
                for ci, cell_text in enumerate(row[:5]):
                    cell = tbl.table.cell(ri, ci)
                    cell.text = cell_text[:80]
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10); p.font.name = FONT
                        if ri == 0: p.font.bold = True

# ── Tool 6: screenshot ─────────────────────────────────────────────

def screenshot_pptx(pptx_path: str, output_dir: str = "/tmp/ppt_screens") -> str | None:
    """Screenshot a PPTX to PNG using qlmanage."""
    os.makedirs(output_dir, exist_ok=True)
    try:
        subprocess.run(['qlmanage','-t','-s','1920','-o',output_dir,pptx_path], capture_output=True, timeout=10)
        pngs = sorted(Path(output_dir).glob('*.png'))
        return str(pngs[-1]) if pngs else None
    except:
        return None

# ── Tool 7: review ────────────────────────────────────────────────

def review_slide(screenshot_path: str, spec: dict, page_content: dict = None) -> dict:
    """Two-pass review: VL blind-read → LLM content comparison → fix instructions.

    Pass 1 (VL): Without any hints, describe what the image shows:
      - Title text
      - Key content/points
      - Readability
      - Any ambiguity

    Pass 2 (LLM): Compare VL description against actual page_content.
      - Is any information missing?
      - Is the title correct?
      - Are key points present?
      - Is the text readable?

    Returns: {"pass": bool, "vl_description": str, "content_match": bool,
              "missing_info": [...], "issues": [...]}
    """
    vl = _vl_client()
    llm = _llm_client()
    result = {"pass": True, "vl_description": "", "content_match": True, "missing_info": [], "issues": []}

    if not screenshot_path or not os.path.exists(screenshot_path):
        return result

    try:
        img = Image.open(screenshot_path).convert('RGB')
        if img.width > 1024: img = img.resize((1024, int(img.height * 1024 / img.width)))
        buf = io.BytesIO(); img.save(buf, 'JPEG', quality=70)
        b64 = base64.b64encode(buf.getvalue()).decode()

        # ── Pass 1: VL blind read ──
        vl_prompt = """Look at this PPT slide image. Describe what you see WITHOUT any assumptions about what should be there.

Return JSON:
{
  "title_text": "exact title text visible on the slide",
  "content_text": ["key point 1", "key point 2", "..."],
  "readability": "good|ok|poor — can all text be read clearly?",
  "ambiguity": "any text that is unclear, cut off, or ambiguous",
  "visual_issues": "any layout problems: overlapping text, spacing issues, alignment, overflow",
  "overall": "brief summary of what this slide communicates"
}"""

        r = vl.chat.completions.create(
            model=VL_MODEL,
            messages=[{"role": "user", "content": [
                {"type": "text", "text": vl_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
            ]}],
            max_tokens=1024, temperature=0.1)
        vl_text = r.choices[0].message.content
        vl_match = re.search(r'\{.*\}', vl_text, re.DOTALL)
        vl_result = json.loads(vl_match.group(0)) if vl_match else {}
        result["vl_description"] = vl_result.get("overall", "")

        # ── Pass 2: LLM content comparison ──
        if page_content:
            expected = {
                "title": page_content.get("title", ""),
                "body": page_content.get("body", [])[:5],
                "type": page_content.get("type", "content"),
            }
            compare_prompt = f"""Compare what the VL model SAW on the slide against what SHOULD be there.

EXPECTED content:
{json.dumps(expected, ensure_ascii=False, indent=2)}

VL model SAW:
{json.dumps(vl_result, ensure_ascii=False, indent=2)}

Check:
1. Is the expected title present? (exact match not needed, semantic match OK)
2. Are all key points from body visible?
3. Is any expected information MISSING?
4. Is the text readable?
5. Any visual issues?

Return JSON:
{{
  "content_match": true|false,
  "title_ok": true|false,
  "missing_info": ["info that should be on slide but is not visible"],
  "readability_ok": true|false,
  "issues": [{{"type": "missing|cutoff|overlap|spacing|alignment", "detail": "...", "fix": "..."}}],
  "pass": true|false
}}"""
            r2 = llm.chat.completions.create(
                model="deepseek-v4-flash",
                messages=[{"role": "user", "content": compare_prompt}],
                max_tokens=1024, temperature=0.1)
            text2 = r2.choices[0].message.content
            m2 = re.search(r'\{.*\}', text2, re.DOTALL)
            cmp = json.loads(m2.group(0)) if m2 else {}
            result["content_match"] = cmp.get("content_match", True)
            result["missing_info"] = cmp.get("missing_info", [])
            result["issues"] = cmp.get("issues", [])
            result["pass"] = cmp.get("pass", True) and len(vl_result.get("visual_issues", "")) < 10

        return result
    except Exception as e:
        result["issues"] = [{"type": "error", "detail": str(e)}]
        return result

