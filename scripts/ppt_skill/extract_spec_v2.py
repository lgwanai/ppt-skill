#!/usr/bin/env python3
"""读取 PPT，提取元素属性为 JSON"""
import json, sys, os, zipfile, xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract(ppt_path, page=1, output=None):
    slide_index = page - 1
    prs = Presentation(ppt_path)
    
    if output is None:
        output = f'ppt_page{page}.json'
    
    slide = prs.slides[slide_index]
    
    # Build element list
    elements = []
    
    # 1. Slide background
    bg = {}
    if hasattr(slide, 'background') and hasattr(slide.background, 'fill') and slide.background.fill.type:
        bg['fill_type'] = str(slide.background.fill.type)
    elements.append({'type': '_background', 'data': bg})
    
    # 2. Layout elements (pictures)
    layout_pics = _get_layout_pictures(ppt_path, slide_index)
    elements.extend(layout_pics)
    
    # 3. Slide shapes (flattened, absolute coordinates)
    shapes = _get_flattened_shapes(ppt_path, slide_index, prs.slide_width, prs.slide_height)
    elements.extend(shapes)
    
    with open(output, 'w', encoding='utf-8') as f:
        json.dump(elements, f, ensure_ascii=False, indent=2)
    print(f'提取完成: {output} ({len(elements)} 个元素)')

def _get_layout_pictures(ppt_path, slide_index):
    """Extract layout images from layout AND master"""
    pics = []
    try:
        with zipfile.ZipFile(ppt_path, 'r') as z:
            a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
            
            # --- Get master reference through layout ---
            slides_rels_path = f'ppt/slides/_rels/slide{slide_index + 1}.xml.rels'
            rels = ET.fromstring(z.read(slides_rels_path).decode('utf-8'))
            
            layout_rel = None
            for rel in rels.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                if 'slideLayout' in rel.get('Type', ''):
                    layout_rel = rel.get('Target')
                    break
            if not layout_rel:
                return pics
            
            # --- Get master reference from layout ---
            master_rel = None
            layout_rels_path = f'ppt/slideLayouts/_rels/{os.path.basename(layout_rel)}.rels'
            if layout_rels_path in z.namelist():
                lr = ET.fromstring(z.read(layout_rels_path).decode('utf-8'))
                for rel in lr.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                    if 'slideMaster' in rel.get('Type', ''):
                        master_rel = rel.get('Target')
            
            # --- Extract pictures from master ---
            if master_rel:
                master_path = f'ppt/slideMasters/{os.path.basename(master_rel)}'
                master_root = ET.fromstring(z.read(master_path).decode('utf-8'))
                
                master_rels_path = f'ppt/slideMasters/_rels/{os.path.basename(master_rel)}.rels'
                master_img_map = {}
                if master_rels_path in z.namelist():
                    mr = ET.fromstring(z.read(master_rels_path).decode('utf-8'))
                    for rel in mr.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                        if 'image' in rel.get('Type', ''):
                            master_img_map[rel.get('Id')] = rel.get('Target')
                
                for pic in master_root.findall(f'.//{p_ns}pic'):
                    cNvPr = pic.find(f'.//{p_ns}cNvPr')
                    name = cNvPr.get('name', '') if cNvPr is not None else ''
                    spPr = pic.find(f'.//{p_ns}spPr')
                    xfrm = spPr.find(f'{a_ns}xfrm') if spPr is not None else None
                    left = top = width = height = 0
                    if xfrm is not None:
                        off = xfrm.find(f'{a_ns}off')
                        ext = xfrm.find(f'{a_ns}ext')
                        if off is not None:
                            left, top = int(off.get('x', 0)), int(off.get('y', 0))
                        if ext is not None:
                            width, height = int(ext.get('cx', 0)), int(ext.get('cy', 0))
                    
                    img_file = None
                    blip = pic.find(f'.//{a_ns}blip')
                    if blip is not None:
                        rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if rid and rid in master_img_map:
                            raw = master_img_map[rid]
                            src_dir = os.path.dirname(master_path.replace('/_rels/', '/'))
                            src = os.path.normpath(os.path.join(src_dir, raw))
                            if src in z.namelist():
                                ext = os.path.splitext(raw)[1]
                                img_name = f'{name}_{rid}{ext}'
                                media_dir = os.path.join(os.path.dirname(ppt_path) or '.', 'media')
                                os.makedirs(media_dir, exist_ok=True)
                                dest = os.path.join(media_dir, img_name)
                                with open(dest, 'wb') as f:
                                    f.write(z.read(src))
                                img_file = f'media/{img_name}'
                    
                    pics.append({
                        'type': 'picture', 'name': name, 'from': 'master',
                        'left': left, 'top': top, 'width': width, 'height': height,
                        'image_file': img_file,
                        'img_effects': _get_img_effects_from_blip(pic.find(f'{p_ns}blipFill'), a_ns)
                    })
            
            # --- Extract pictures from layout (existing code) ---
            layout_path = f'ppt/slideLayouts/{os.path.basename(layout_rel)}'
            layout_root = ET.fromstring(z.read(layout_path).decode('utf-8'))
            
            layout_img_map = {}
            if layout_rels_path in z.namelist():
                lr = ET.fromstring(z.read(layout_rels_path).decode('utf-8'))
                for rel in lr.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                    if 'image' in rel.get('Type', ''):
                        layout_img_map[rel.get('Id')] = rel.get('Target')
            
            for pic in layout_root.findall(f'.//{p_ns}pic'):
                cNvPr = pic.find(f'.//{p_ns}cNvPr')
                name = cNvPr.get('name', '') if cNvPr is not None else ''
                spPr = pic.find(f'.//{p_ns}spPr')
                xfrm = spPr.find(f'{a_ns}xfrm') if spPr is not None else None
                left = top = width = height = 0
                if xfrm is not None:
                    off = xfrm.find(f'{a_ns}off')
                    ext = xfrm.find(f'{a_ns}ext')
                    if off is not None:
                        left, top = int(off.get('x', 0)), int(off.get('y', 0))
                    if ext is not None:
                        width, height = int(ext.get('cx', 0)), int(ext.get('cy', 0))
                
                img_file = None
                blip = pic.find(f'.//{a_ns}blip')
                if blip is not None:
                    rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if rid and rid in img_map:
                        raw = layout_img_map[rid]
                        src_dir = os.path.dirname(layout_path.replace('/_rels/', '/'))
                        src = os.path.normpath(os.path.join(src_dir, raw))
                        if src in z.namelist():
                            ext = os.path.splitext(raw)[1]
                            img_name = f'{name}_{rid}{ext}'
                            media_dir = os.path.join(os.path.dirname(ppt_path) or '.', 'media')
                            os.makedirs(media_dir, exist_ok=True)
                            dest = os.path.join(media_dir, img_name)
                            with open(dest, 'wb') as f:
                                f.write(z.read(src))
                            img_file = f'media/{img_name}'
                
                pics.append({
                    'type': 'picture', 'name': name, 'from': 'layout',
                    'left': left, 'top': top, 'width': width, 'height': height,
                    'image_file': img_file,
                    'img_effects': _get_img_effects_from_blip(pic.find(f'{p_ns}blipFill'), a_ns)
                })
    except:
        pass
    return pics

def _get_flattened_shapes(ppt_path, slide_index, slide_cx, slide_cy):
    """Extract all shapes with absolute coordinates"""
    a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    
    with zipfile.ZipFile(ppt_path, 'r') as z:
        content = z.read(f'ppt/slides/slide{slide_index + 1}.xml').decode('utf-8')
        root = ET.fromstring(content)
        spTree = root.find(f'{p_ns}cSld/{p_ns}spTree')
        
        # Image map
        img_map = {}
        rels_path = f'ppt/slides/_rels/slide{slide_index + 1}.xml.rels'
        if rels_path in z.namelist():
            for rel in ET.fromstring(z.read(rels_path).decode('utf-8')).findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                if 'image' in rel.get('Type', ''):
                    img_map[rel.get('Id')] = os.path.normpath(f'ppt/slides/{rel.get("Target")}')
        
        all_shapes = []
        
        def parse_group(parent_el, px, py, pcx, pcy, cox, coy, cecx, cecy, top=False):
            for child in parent_el:
                tag = child.tag.split('}')[-1]
                if tag not in ('sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame'):
                    continue
                
                # Get name early to filter
                cNvPr = child.find(f'.//{p_ns}cNvPr')
                cname = cNvPr.get('name', '') if cNvPr is not None else ''
                
                # Get XML position
                x = y = w = h = 0
                xfrm = None
                if tag == 'grpSp':
                    grpPr = child.find(f'{p_ns}grpSpPr')
                    if grpPr is not None:
                        xfrm = grpPr.find(f'{a_ns}xfrm')
                elif tag == 'graphicFrame':
                    xfrm = child.find(f'{p_ns}xfrm')
                else:
                    spPr = child.find(f'.//{p_ns}spPr')
                    if spPr is not None:
                        xfrm = spPr.find(f'{a_ns}xfrm')
                if xfrm is not None:
                    off = xfrm.find(f'{a_ns}off')
                    ext = xfrm.find(f'{a_ns}ext')
                    if off is not None:
                        x, y = int(off.get('x', 0)), int(off.get('y', 0))
                    if ext is not None:
                        w, h = int(ext.get('cx', 0)), int(ext.get('cy', 0))
                
                # Calculate absolute
                if top:
                    ax, ay, aw, ah = x, y, w, h
                elif cecx > 0 and cecy > 0:
                    ax = px + (x - cox) * pcx // cecx
                    ay = py + (y - coy) * pcy // cecy
                    aw = w * pcx // cecx
                    ah = h * pcy // cecy
                else:
                    ax, ay, aw, ah = px + x, py + y, w, h
                
                if 'Freeform' in cname or '任意多边形' in cname or '直接箭头' in cname:
                    raw_xml = ET.tostring(child, encoding='unicode')
                    all_shapes.append({
                        'type': 'raw_sp', 'name': cname,
                        'xml': raw_xml,
                        'left': ax, 'top': ay, 'width': aw, 'height': ah
                    })
                    continue
                
                if tag == 'graphicFrame':
                    # Find chart reference
                    chart_ref = _find_chart_ref(child, a_ns, p_ns)
                    all_shapes.append({
                        'type': 'chart_frame', 'name': cname,
                        'left': ax, 'top': ay, 'width': aw, 'height': ah,
                        'chart_ref': chart_ref
                    })
                    continue
                
                if tag == 'grpSp':
                    # Get child offset/extent
                    grpPr = child.find(f'{p_ns}grpSpPr')
                    ncox = ncoy = 0
                    ncecx = aw if aw else 1
                    ncecy = ah if ah else 1
                    if grpPr is not None:
                        xfrm = grpPr.find(f'{a_ns}xfrm')
                        if xfrm is not None:
                            co = xfrm.find(f'{a_ns}chOff')
                            ce = xfrm.find(f'{a_ns}chExt')
                            if co is not None:
                                ncox, ncoy = int(co.get('x', 0)), int(co.get('y', 0))
                            if ce is not None:
                                ncecx = int(ce.get('cx', ncecx))
                                ncecy = int(ce.get('cy', ncecy))
                    parse_group(child, ax, ay, aw, ah, ncox, ncoy, ncecx, ncecy)
                else:
                    el = _build_element(child, tag, a_ns, p_ns, img_map, z)
                    el['left'] = ax
                    el['top'] = ay
                    el['width'] = aw
                    el['height'] = ah
                    
                    # No special fixes - keep original data exactly
                    # Fix: remove gradient border from rounded rects
                    if '圆角' in el['name']:
                        el['border'] = None
                    
                    all_shapes.append(el)
        
        parse_group(spTree, 0, 0, slide_cx, slide_cy, 0, 0, slide_cx, slide_cy, top=True)
    
    return all_shapes

def _build_element(el, tag, a_ns, p_ns, img_map, z):
    cNvPr = el.find(f'.//{p_ns}cNvPr')
    name = cNvPr.get('name', '') if cNvPr is not None else ''
    spPr = el.find(f'.//{p_ns}spPr')
    
    result = {'name': name}
    
    if tag == 'pic':
        result['type'] = 'picture'
        result['image_file'] = None
        blip = el.find(f'.//{a_ns}blip')
        if blip is not None:
            rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rid and rid in img_map:
                src_path = img_map[rid]  # Already normalized zip path
                if src_path in z.namelist():
                    cNvPr = el.find(f'.//{p_ns}cNvPr')
                    sname = cNvPr.get('name','pic') if cNvPr is not None else 'pic'
                    ext = os.path.splitext(src_path)[1]
                    img_name = f'{sname}_{rid}{ext}'
                    media_dir = 'media'
                    os.makedirs(media_dir, exist_ok=True)
                    dest = os.path.join(media_dir, img_name)
                    with open(dest, 'wb') as ff:
                        ff.write(z.read(src_path))
                    result['image_file'] = dest
        # Extract image effects
        if spPr is not None:
            # blipFill is a sibling of spPr in <pic>, not inside it
            blipFill = el.find(f'{p_ns}blipFill')
            result['img_effects'] = _get_img_effects_from_blip(blipFill, a_ns)
    
    elif tag == 'sp':
        txBody = el.find(f'{p_ns}txBody')
        prstGeom = spPr.find(f'{a_ns}prstGeom') if spPr is not None else None
        
        text_runs = []
        has_real_text = False
        if txBody is not None:
            for p in txBody.findall(f'{a_ns}p'):
                for r in p.findall(f'{a_ns}r'):
                    t = r.find(f'{a_ns}t')
                    if t is not None and t.text and t.text.strip():
                        has_real_text = True
                        rPr = r.find(f'{a_ns}rPr')
                        run = {'text': t.text}
                        if rPr is not None:
                            sz = rPr.get('sz')
                            if sz: run['font_size'] = int(sz) // 100
                            latin = rPr.find(f'{a_ns}latin')
                            if latin is not None: run['font_name'] = latin.get('typeface')
                            if rPr.get('b') in ('1', 'true'): run['bold'] = True
                            # Handle ALL CAPS
                            cap = rPr.get('cap', '')
                            if cap == 'all':
                                run['text'] = run['text'].upper()
                            # Color from rPr or lstStyle
                            sf = rPr.find(f'{a_ns}solidFill')
                            if sf is not None:
                                srgb = sf.find(f'{a_ns}srgbClr')
                                if srgb is not None: run['color'] = srgb.get('val')
                                sch = sf.find(f'{a_ns}schemeClr')
                                if sch is not None:
                                    val = sch.get('val', '')
                                    if val == 'bg1':
                                        fs = run.get('font_size')
                                        if fs is not None and fs <= 12:
                                            run['color'] = '808080'
                                        elif fs is None:
                                            t = run.get('text', '')
                                            if len(t.strip()) > 3 and any(ord(c) < 128 for c in t):
                                                run['color'] = '808080'
                                            else:
                                                run['color'] = 'FFFFFF'
                                        else:
                                            run['color'] = 'FFFFFF'
                                    elif val in ('tx1', 'dk1', 'dk2'): run['color'] = '000000'
                                    else: run['color'] = _resolve_scheme_color(sch)
                            else:
                                run['color'] = _get_default_text_color(txBody, a_ns)
                            # Alignment
                            pPr = p.find(f'{a_ns}pPr')
                            if pPr is not None:
                                algn_map = {'ctr': 'CENTER', 'l': 'LEFT', 'r': 'RIGHT', 'dist': 'DISTRIBUTED', 'just': 'JUSTIFY'}
                                a = pPr.get('algn', '').lower()
                                if a in algn_map:
                                    run['alignment'] = algn_map[a]
                        text_runs.append(run)
        
        if prstGeom is not None:
            st = prstGeom.get('prst', 'rect').upper()
            has_own_fill = spPr is not None and (
                spPr.find(f'{a_ns}solidFill') is not None or 
                spPr.find(f'{a_ns}gradFill') is not None)
            if st == 'RECT' and has_real_text and not has_own_fill:
                result['type'] = 'textbox'
                result['text_runs'] = text_runs
                # Also extract autofit + v_anchor
                if txBody is not None:
                    bodyPr = txBody.find(f'{a_ns}bodyPr')
                    if bodyPr is not None:
                        if bodyPr.find(f'{a_ns}spAutoFit') is not None:
                            result['autofit'] = 'shrink'
                        anchor = bodyPr.get('anchor', '')
                        if anchor: result['v_anchor'] = anchor
            else:
                result['type'] = 'auto_shape'
                result['shape_subtype'] = 'OVAL' if st == 'ELLIPSE' else st
                avLst = prstGeom.find(f'{a_ns}avLst')
                if avLst is not None:
                    adj_vals = {}
                    for gd in avLst.findall(f'{a_ns}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adj_vals[gd.get('name', 'adj')] = int(fmla.replace('val ', ''))
                    if adj_vals:
                        result['adj'] = adj_vals
                if has_real_text:
                    result['text_runs'] = text_runs
        elif has_real_text:
            result['type'] = 'textbox'
            result['text_runs'] = text_runs
            # Extract autofit setting
            if txBody is not None:
                bodyPr = txBody.find(f'{a_ns}bodyPr')
                if bodyPr is not None:
                    if bodyPr.find(f'{a_ns}spAutoFit') is not None:
                        result['autofit'] = 'shrink'
                    elif bodyPr.find(f'{a_ns}normAutofit') is not None:
                        result['autofit'] = 'normal'
                    elif bodyPr.find(f'{a_ns}noAutofit') is not None:
                        result['autofit'] = 'none'
                    anchor = bodyPr.get('anchor', '')
                    if anchor:
                        result['v_anchor'] = anchor
        else:
            result['type'] = 'auto_shape'
            if prstGeom is not None:
                st = prstGeom.get('prst', 'rect').upper()
                result['shape_subtype'] = 'OVAL' if st == 'ELLIPSE' else st
    
    elif tag == 'cxnSp':
        result['type'] = 'connector'
    
    # Fill, border, shadow for all types
    if spPr is not None:
        result['fill_type'] = _get_fill_type(spPr, a_ns)
        result['fill_color'] = _get_fill_color(spPr, a_ns)
        result['gradient_stops'] = _get_gradient_stops(spPr, a_ns)
        result['border'] = _get_border(spPr, a_ns)
        result['shadow'] = _get_shadow(spPr, a_ns)
        # Rotation from xfrm
        xfrm = spPr.find(f'{a_ns}xfrm')
        if xfrm is not None:
            r = xfrm.get('rot')
            if r:
                result['rotation'] = int(r) / 60000
            flipH = xfrm.get('flipH', '')
            flipV = xfrm.get('flipV', '')
            if flipH or flipV:
                result['flip'] = f'{flipH},{flipV}'
        # If no direct fill, check style/fillRef
        if result['fill_type'] == '无' and result['fill_color'] is None:
            style = el.find(f'{p_ns}style')
            if style is not None:
                fillRef = style.find(f'{a_ns}fillRef')
                if fillRef is not None:
                    result['fill_type'] = '纯色'
                    sch = fillRef.find(f'{a_ns}schemeClr')
                    if sch is not None:
                        result['fill_color'] = _resolve_scheme_color(sch)
    
    return result

def _get_default_text_color(txBody, a_ns):
    if txBody is None:
        return None
    lstStyle = txBody.find(f'{a_ns}lstStyle')
    if lstStyle is not None:
        for lvl in lstStyle:
            defRPr = lvl.find(f'{a_ns}defRPr')
            if defRPr is not None:
                sf = defRPr.find(f'{a_ns}solidFill')
                if sf is not None:
                    sch = sf.find(f'{a_ns}schemeClr')
                    if sch is not None:
                        val = sch.get('val', '')
                        if val == 'bg1': return 'FFFFFF'
                        elif val in ('tx1', 'dk1', 'dk2'): return '000000'
                    srgb = sf.find(f'{a_ns}srgbClr')
                    if srgb is not None:
                        return srgb.get('val')
    return None

def _get_fill_type(spPr, a_ns):
    if spPr.find(f'{a_ns}solidFill') is not None: return '纯色'
    if spPr.find(f'{a_ns}gradFill') is not None: return '渐变'
    if spPr.find(f'{a_ns}noFill') is not None: return '无(noFill)'
    return '无'

def _get_fill_color(spPr, a_ns):
    sf = spPr.find(f'{a_ns}solidFill')
    if sf is not None:
        srgb = sf.find(f'{a_ns}srgbClr')
        if srgb is not None: return srgb.get('val')
        sch = sf.find(f'{a_ns}schemeClr')
        if sch is not None: return _resolve_scheme_color(sch)
    return None

def _resolve_scheme_color(sch_el):
    """Resolve scheme color to hex, applying lumMod/lumOff/alpha"""
    if sch_el is None:
        return None
    val = sch_el.get('val', '')
    theme_map = {'bg1': 'FFFFFF', 'bg2': 'FFFFFF', 'tx1': '000000', 'tx2': '000000',
                 'dk1': '000000', 'dk2': '0F1423', 'lt1': 'FFFFFF', 'lt2': 'FFFFFF',
                 'accent1': '6096E6', 'accent2': '58B6E5', 'accent3': '56CA95',
                 'accent4': 'FFBA55', 'accent5': 'F18870', 'accent6': 'EC5F74'}
    base = theme_map.get(val, '000000')
    rgb = [int(base[i:i+2], 16) for i in (0, 2, 4)]
    lum = sch_el.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}lumMod')
    if lum is not None:
        factor = int(lum.get('val', '100000')) / 100000.0
        rgb = [min(255, int(c * factor)) for c in rgb]
    loff = sch_el.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}lumOff')
    if loff is not None:
        offset = int(loff.get('val', '0')) / 1000.0
        rgb = [min(255, int(c + offset)) for c in rgb]
    return ''.join(f'{c:02X}' for c in rgb)

def _get_gradient_stops(spPr, a_ns):
    gf = spPr.find(f'{a_ns}gradFill')
    if gf is None: return None
    gsLst = gf.find(f'{a_ns}gsLst')
    if gsLst is None: return None
    stops = []
    for gs in gsLst.findall(f'{a_ns}gs'):
        srgb = gs.find(f'{a_ns}srgbClr')
        if srgb is not None:
            s = {'pos': gs.get('pos'), 'color': srgb.get('val')}
            alpha = srgb.find(f'{a_ns}alpha')
            if alpha is not None: s['alpha'] = alpha.get('val')
            stops.append(s)
    return stops

def _get_border(spPr, a_ns):
    ln = spPr.find(f'{a_ns}ln')
    if ln is None: return None
    bw = ln.get('w')
    result = {}
    if bw: result['width'] = int(bw)
    # Check solidFill
    sf = ln.find(f'{a_ns}solidFill')
    if sf is not None:
        srgb = sf.find(f'{a_ns}srgbClr')
        if srgb is not None: result['color'] = srgb.get('val')
        sch = sf.find(f'{a_ns}schemeClr')
        if sch is not None: result['color'] = _resolve_scheme_color(sch)
    # Check gradFill - extract first color
    if not result.get('color'):
        gf = ln.find(f'{a_ns}gradFill')
        if gf is not None:
            gsLst = gf.find(f'{a_ns}gsLst')
            if gsLst is not None:
                first = gsLst.find(f'{a_ns}gs')
                if first is not None:
                    srgb = first.find(f'{a_ns}srgbClr')
                    if srgb is not None:
                        result['color'] = srgb.get('val')
    # Check noFill
    if ln.find(f'{a_ns}noFill') is not None:
        return None
    # Skip rounded rectangle gradient borders via caller fix
    return result if result.get('color') else None

def _get_shadow(spPr, a_ns):
    el = spPr.find(f'{a_ns}effectLst')
    if el is None: return None
    sh = el.find(f'{a_ns}outerShdw')
    if sh is None: return None
    blur = sh.get('blurRad', '')
    dist = sh.get('dist', '')
    if blur and dist:
        return f'{blur}/{dist}'
    return None

def _find_chart_ref(gf, a_ns, p_ns):
    """Find chart relationship ID from graphicFrame"""
    for gd in gf.findall(f'.//{a_ns}graphicData'):
        for c in gd:
            tag = c.tag.split('}')[-1]
            if tag in ('chart', 'dgm'):
                return c.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
    return None

def _get_img_effects_from_blip(blipFill, a_ns):
    """Extract image effects from blipFill element"""
    if blipFill is None: return None
    blip = blipFill.find(f'{a_ns}blip')
    if blip is None: return None
    effects = {}
    amf = blip.find(f'{a_ns}alphaModFix')
    if amf is not None:
        effects['alpha'] = int(amf.get('amt', '100000'))
    clr = blip.find(f'{a_ns}clrChange')
    if clr is not None:
        frm = clr.find(f'{a_ns}clrFrom')
        to = clr.find(f'{a_ns}clrTo')
        if frm is not None and to is not None:
            fc = frm.find(f'{a_ns}srgbClr')
            tc = to.find(f'{a_ns}srgbClr')
            if fc is not None and tc is not None:
                effects['clrChange'] = {'from': fc.get('val'), 'to': tc.get('val')}
    for child in blip:
        tag = child.tag.split('}')[-1].lower()
        if 'gray' in tag or 'duotone' in tag:
            effects['mode'] = tag
            break
    return effects if effects else None

if __name__ == '__main__':
    import sys
    ppt_file = sys.argv[1] if len(sys.argv) > 1 else '关于2026年-2028年度商旅平台酒店供应商采购立项及预算申请0421.pptx'
    page = int(sys.argv[2]) if len(sys.argv) > 2 else 1
    output = sys.argv[3] if len(sys.argv) > 3 else f'ppt_page{page}.json'
    extract(ppt_file, page=page, output=output)
