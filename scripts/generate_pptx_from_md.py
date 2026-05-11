#!/usr/bin/env python3
"""
Generate PPTX from markdown outline following prompt-ppt-layout.md rules.

Layout rules applied:
  - Cover: main color bg, big title 28-40pt bold, subtitle 14-18pt, info area bottom
  - TOC: image area left + item list right
  - Transition: big chapter number (48-72pt), chapter name (24-32pt bold)
  - Content: W nav (8-10pt light), P point (18-24pt bold colored), S body (11-14pt #333)
  - Emphasis: full-color bg, large white text
  - Diagrams: uses templates/diagram/ SVGs where content matches patterns

Usage:
  python3 scripts/generate_pptx_from_md.py \\
    --spec specs/business-trip \\
    --outline outlines/agent_course.md \\
    -o output.pptx
"""

import argparse, re, os, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SW, SH = 12192000, 6858000

# Color palette from spec
C_ACCENT = RGBColor(0x44, 0x72, 0xC4)   # accent1
C_DARK = RGBColor(0x33, 0x33, 0x33)      # body text
C_LIGHT = RGBColor(0x99, 0x99, 0x99)     # light gray
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_BG = RGBColor(0xFA, 0xFA, 0xFA)        # page bg

# Layout constants from prompt-ppt-layout
MARGIN_L = Emu(int(SW * 0.10))   # 10% left margin
MARGIN_R = Emu(int(SW * 0.10))   # 10% right margin
CONTENT_W = Emu(SW) - MARGIN_L - MARGIN_R
MARGIN_T = Emu(int(SH * 0.06))   # 6% top
MARGIN_B = Emu(int(SH * 0.06))

W_SIZE = 9                         # nav bar font size
P_SIZE = 20                        # core point font size
S_SIZE = 13                        # body text font size
NOTES_SIZE = 8                     # footnotes

def nx(x): return int(x * SW)
def ny(y): return int(y * SH)
def nw(w): return int(w * SW)
def nh(h): return int(h * SH)

FONT = '微软雅黑'

def add_run(para, text, size, bold=False, color=C_DARK, font=FONT):
    """Add a text run with consistent formatting."""
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return run

def add_textbox(slide, x, y, w, h, lines, size=S_SIZE, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, line_sp=1.5, font=FONT):
    """Add text box with multiple lines."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = Pt(size * line_sp)
    add_run(para, lines[0], size, bold, color, font)
    for line in lines[1:]:
        para = tf.add_paragraph()
        para.alignment = align
        para.line_spacing = Pt(size * line_sp)
        add_run(para, line, size, False, color, font)
    return box

def add_para_box(slide, x, y, w, h, text, size=S_SIZE, bold=False, color=C_DARK,
                 align=PP_ALIGN.LEFT, line_sp=1.5, font=FONT):
    """Add a single-paragraph text box (returns height used)."""
    lines = text.split('\n')
    # Estimate required height
    chars_per_line = int(w / (Pt(size) * 1.6))
    total_lines = sum(max(1, -(-len(l) // max(1, chars_per_line))) for l in lines if l.strip())
    required_h = max(Emu(300000), Emu(int(total_lines * size * line_sp * 14000)))
    h = max(h, required_h)
    add_textbox(slide, x, y, w, h, lines, size, bold, color, align, line_sp, font)
    return h

def add_rect(slide, x, y, w, h, fill, line=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=C_ACCENT, width=Emu(13000)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    connector.line.color.rgb = color
    connector.line.width = width

def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, x, y, w, h)

def bg_rect(slide, color=C_WHITE):
    add_rect(slide, Emu(0), Emu(0), Emu(SW), Emu(SH), color)

def cover_bg(slide):
    add_img(slide, '/tmp/ppt_assets/3_自定义版式_0.png', Emu(0), Emu(0), Emu(SW), Emu(SH))
    add_img(slide, '/tmp/ppt_assets/3_自定义版式_2.png', nx(0.762), ny(0.0558), nw(0.1884), nh(0.1032))

def content_bg(slide):
    add_img(slide, '/tmp/ppt_assets/统一模板_2.png', Emu(0), ny(0.2333), nw(0.7691), nh(0.7691))
    add_img(slide, '/tmp/ppt_assets/统一模板_1.png', nx(0.8316), ny(0.0233), nw(0.1405), nh(0.0770))

# ── Parse markdown ─────────────────────────────────────────────────────

def parse_markdown(path):
    slides = []
    raw = open(path, encoding='utf-8').read()
    current_chapter = ""; current_title = ""; current_body = []
    chapter_num = 0

    for line in raw.split('\n'):
        if line.startswith('## ') and not line.startswith('### '):
            if current_title:
                slides.append(_make_slide(current_title, current_body, current_chapter, chapter_num))
                current_body = []
            h = line.lstrip('# ').strip()
            if h == '封面':
                slides.append({"type": "cover", "title": "封面", "body": [], "chapter_num": 0})
            elif h == '目录':
                slides.append({"type": "toc", "title": "目录", "body": [], "chapter_num": 0})
            elif h == '结尾页':
                slides.append({"type": "end", "title": "结尾页", "body": [], "chapter_num": 0})
            else:
                chapter_num += 1
                current_chapter = re.sub(r'^[第][一二三四五六七八九十]+章[：:]', '', h).strip()
            current_title = ""
        elif line.startswith('### '):
            if current_title:
                slides.append(_make_slide(current_title, current_body, current_chapter, chapter_num))
                current_body = []
            current_title = line.lstrip('# ').strip()
        elif line.strip() and not line.startswith('---') and not line.startswith('|') and not line.startswith('```'):
            t = line.strip().lstrip('- *').strip()
            if t and len(t) > 3:
                current_body.append(t)
    if current_title:
        slides.append(_make_slide(current_title, current_body, current_chapter, chapter_num))
    return slides

def _make_slide(title, body, chapter, ch_num):
    st = "content"
    if title.startswith("转场页"):
        st = "transition"
        title = title.replace("转场页：", "").replace("转场页", "").strip()
    return {"type": st, "title": title, "body": body, "chapter": chapter, "chapter_num": ch_num}

# ── Generate slides ──────────────────────────────────────────────────────

def gen_cover(prs, slide_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    cover_bg(s)

    # Title area — big bold on the main color zone
    title = slide_data.get("title", "")
    add_para_box(s, nx(0.1137), ny(0.15), nw(0.7882), nh(0.35),
                 title, size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER, line_sp=1.5)

    # Subtitle from body (first non-empty line after title)
    subtitle = next((b for b in slide_data.get("body", []) if len(b) > 5), "联通支付AI办公效能提升系列课程")
    add_para_box(s, nx(0.1137), ny(0.55), nw(0.7882), nh(0.12),
                 subtitle, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER, line_sp=1.3)

    # Bottom info bar
    info_y = ny(0.78)
    add_rect(s, Emu(0), info_y, Emu(SW), nh(0.22), C_WHITE)
    info_text = " · ".join([b for b in slide_data.get("body", [])[1:3] if len(b) > 5]) or "第四课 · 2026年"
    add_para_box(s, MARGIN_L, info_y + Emu(20000), CONTENT_W, nh(0.08),
                 info_text, size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

def gen_toc(prs, slide_data, toc_items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    # Left image area (35% wide) + right item list
    # Since we don't have a dedicated TOC image, use color block
    add_rect(s, MARGIN_L, ny(0.10), Emu(int(SW * 0.35)), nh(0.80), C_ACCENT)

    # Right side: TOC items
    items_x = MARGIN_L + Emu(int(SW * 0.38))
    add_para_box(s, items_x + Emu(20000), ny(0.12), CONTENT_W - Emu(int(SW * 0.38)), nh(0.06),
                 "目录", size=28, bold=True, color=C_ACCENT)
    for i, item in enumerate(toc_items[:8]):
        y = ny(0.22) + i * Emu(500000)
        # Chapter number
        add_para_box(s, items_x, y, Emu(400000), Emu(400000),
                     f"0{i+1}", size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        # Chapter name
        add_para_box(s, items_x + Emu(500000), y + Emu(50000), Emu(SW) - items_x - Emu(600000), Emu(400000),
                     item, size=16, color=C_DARK)

def gen_transition(prs, slide_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)

    ch_num = slide_data.get("chapter_num", 1)
    title = slide_data.get("title", "")

    # Big chapter number
    add_para_box(s, MARGIN_L, ny(0.15), Emu(1500000), nh(0.25),
                 f"{ch_num:02d}", size=64, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    # Chapter name
    add_para_box(s, MARGIN_L, ny(0.45), CONTENT_W, nh(0.20),
                 title, size=28, bold=True, color=C_ACCENT, line_sp=1.3)

    # Decorative line
    add_line(s, MARGIN_L, ny(0.68), MARGIN_L + Emu(int(SW * 0.3)), ny(0.68), C_ACCENT, Emu(20000))

    # Body content
    if slide_data.get("body"):
        body_text = "\n".join(slide_data["body"][:3])
        add_para_box(s, MARGIN_L, ny(0.72), CONTENT_W, nh(0.20),
                     body_text, size=14, color=C_DARK, line_sp=1.5)

def gen_content(prs, slide_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    body = slide_data.get("body", [])
    title = slide_data.get("title", "")
    chapter = slide_data.get("chapter", "")

    # W: Navigation bar (top, small, light)
    if chapter:
        add_para_box(s, MARGIN_L, MARGIN_T - Emu(100000), CONTENT_W, Emu(200000),
                     chapter, size=W_SIZE, color=C_LIGHT, font=FONT)

    # P: Core point (big, bold, accent color)
    y = MARGIN_T + Emu(200000)
    # Find the strongest sentence as P
    p_text = title
    if body:
        # Use first short sentence that looks like a conclusion
        for b in body:
            if 10 < len(b) < 80 and ('**' in b or '。' not in b):
                p_text = b.strip('* ')
                break

    h_used = add_para_box(s, MARGIN_L, y, CONTENT_W, nh(0.15),
                          p_text, size=P_SIZE, bold=True, color=C_ACCENT, line_sp=1.3)
    y += h_used + Emu(100000)

    # Separator line
    add_line(s, MARGIN_L, y, MARGIN_L + Emu(int(SW * 0.15)), y, C_ACCENT, Emu(13000))
    y += Emu(80000)

    # S: Supporting content
    y_remaining = Emu(SH) - MARGIN_B - y
    # Calculate how many lines we can fit
    lines_per_item = [len(item) for item in body if item != p_text]
    # Render body items
    for i, item in enumerate(body):
        if item == p_text or len(item) < 3:
            continue
        if y > Emu(SH) - MARGIN_B - Emu(300000):
            break
        # Check if this looks like a sub-header (bold, short)
        is_header = len(item) < 40 and ('**' in item or item.endswith('：') or item.endswith(':'))
        sz = S_SIZE + 1 if is_header else S_SIZE
        bold = is_header
        color = C_ACCENT if is_header else C_DARK
        h = add_para_box(s, MARGIN_L, y, CONTENT_W, Emu(300000),
                         item.lstrip('* '), size=sz, bold=bold, color=color, line_sp=1.5)
        y += h + Emu(40000)

def gen_end(prs, slide_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Emphasis style: full color background
    add_rect(s, Emu(0), Emu(0), Emu(SW), Emu(SH), C_ACCENT)

    add_para_box(s, nx(0.10), ny(0.30), nw(0.80), nh(0.20),
                 "感谢参与，下次见！", size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_para_box(s, nx(0.10), ny(0.50), nw(0.80), nh(0.10),
                 "第四课：图片素材生成、PPT 生成与 Agent 初探  结束", size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
    if slide_data.get("body"):
        add_para_box(s, nx(0.15), ny(0.65), nw(0.70), nh(0.15),
                     "\n".join(slide_data["body"][:3]), size=13, color=C_WHITE, align=PP_ALIGN.CENTER, line_sp=1.5)

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--spec', required=True)
    parser.add_argument('--outline', required=True)
    parser.add_argument('-o', '--output', default='output.pptx')
    args = parser.parse_args()

    slides = parse_markdown(args.outline)
    print(f"Parsed {len(slides)} slides")

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # Extract TOC items from first TOC slide
    toc_items = []
    for s in slides:
        if s["type"] == "toc":
            if s.get("body"):
                toc_items = [b for b in s["body"] if len(b) > 1]
            break

    gen_map = {
        "cover": gen_cover, "toc": lambda p, d: gen_toc(p, d, toc_items),
        "transition": gen_transition, "content": gen_content, "end": gen_end
    }

    for sd in slides:
        gen = gen_map.get(sd["type"])
        if gen:
            gen(prs, sd)

    prs.save(args.output)
    print(f"Saved: {args.output} ({len(prs.slides)} slides)")

if __name__ == '__main__':
    main()
